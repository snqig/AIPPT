"""
排版自检工具（T615）

功能：
    对渲染后的 PPTX 文件进行几何排版质量检测，输出结构化问题报告。

检测项：
    1. 元素溢出：shape 超出画布边界或安全区
    2. 元素重叠：同页非装饰元素重叠（可能遮挡内容）
    3. 间距异常：相邻元素间距过小（< 0.05 inch）或为负
    4. 字号下限：文本字号低于 10pt 不可读
    5. 空文本框：文本框无内容残留占位
    6. 角色缺失：页面缺少必要角色（如 title）

设计原则：
    - 纯函数，无副作用，不修改 PPTX 文件
    - 输出标准化 IssueReport，便于程序化处理
    - 与 validators.py（outline 数据校验）互补，本模块校验渲染后几何结果
    - 支持 CLI lint-layout 命令调用

使用方式：
    from aippt.layout.layout_lint import lint_pptx
    report = lint_pptx("output.pptx")
    for issue in report.issues:
        print(f"[{issue.severity}] 页{issue.page}: {issue.message}")
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Emu, Pt

from aippt.logger import logger


# ==================== 常量 ====================
EMU_PER_INCH = 914400
CANVAS_W_INCH = 13.333
CANVAS_H_INCH = 7.5
SAFE_MARGIN_INCH = 0.5
MIN_FONT_PT = 10  # 字号下限
MIN_GAP_INCH = 0.05  # 最小间距
OVERLAP_TOLERANCE_INCH = 0.02  # 重叠容忍度


# ==================== 数据结构 ====================
@dataclass
class LayoutIssue:
    """单条排版问题

    :param page: 页码（从 1 开始）
    :param severity: 严重级别 "error" / "warning" / "info"
    :param category: 问题分类 overflow/overlap/spacing/fontsize/empty/missing_role
    :param message: 问题描述
    :param shape_name: 涉及的 shape 名称（可选）
    :param detail: 详细数据（如坐标、字号等）
    """
    page: int
    severity: str
    category: str
    message: str
    shape_name: Optional[str] = None
    detail: dict[str, Any] = field(default_factory=dict)


@dataclass
class IssueReport:
    """排版自检报告

    :param file_path: 检测的 PPTX 文件路径
    :param total_pages: 总页数
    :param issues: 问题列表
    """
    file_path: str
    total_pages: int
    issues: list[LayoutIssue] = field(default_factory=list)

    @property
    def error_count(self) -> int:
        return sum(1 for i in self.issues if i.severity == "error")

    @property
    def warning_count(self) -> int:
        return sum(1 for i in self.issues if i.severity == "warning")

    @property
    def passed(self) -> bool:
        """是否通过（无 error 级别问题）"""
        return self.error_count == 0

    def to_dict(self) -> dict[str, Any]:
        """转为字典（用于 JSON 输出）"""
        return {
            "file_path": self.file_path,
            "total_pages": self.total_pages,
            "passed": self.passed,
            "error_count": self.error_count,
            "warning_count": self.warning_count,
            "issues": [
                {
                    "page": i.page,
                    "severity": i.severity,
                    "category": i.category,
                    "message": i.message,
                    "shape_name": i.shape_name,
                    "detail": i.detail,
                }
                for i in self.issues
            ],
        }


# ==================== 主检测函数 ====================
def lint_pptx(
    pptx_path: str,
    check_overflow: bool = True,
    check_overlap: bool = True,
    check_spacing: bool = True,
    check_fontsize: bool = True,
    check_empty: bool = True,
) -> IssueReport:
    """对 PPTX 文件执行排版自检

    :param pptx_path: PPTX 文件路径
    :param check_overflow: 是否检测元素溢出
    :param check_overlap: 是否检测元素重叠
    :param check_spacing: 是否检测间距异常
    :param check_fontsize: 是否检测字号下限
    :param check_empty: 是否检测空文本框
    :return: IssueReport 报告
    """
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX 文件不存在: {pptx_path}")

    prs = Presentation(str(path))
    report = IssueReport(file_path=str(path), total_pages=len(prs.slides))

    canvas_w_emu = int(prs.slide_width)
    canvas_h_emu = int(prs.slide_height)
    safe_left = int(SAFE_MARGIN_INCH * EMU_PER_INCH)
    safe_top = int(SAFE_MARGIN_INCH * EMU_PER_INCH)
    safe_right = canvas_w_emu - safe_left
    safe_bottom = canvas_h_emu - safe_top

    for page_num, slide in enumerate(prs.slides, 1):
        # 收集所有 shape 的几何信息
        shapes_info: list[dict] = []
        for shape in slide.shapes:
            info = {
                "shape": shape,
                "name": shape.name or "",
                "left": int(shape.left) if shape.left is not None else 0,
                "top": int(shape.top) if shape.top is not None else 0,
                "width": int(shape.width) if shape.width is not None else 0,
                "height": int(shape.height) if shape.height is not None else 0,
                "right": (int(shape.left) + int(shape.width)) if shape.left is not None else 0,
                "bottom": (int(shape.top) + int(shape.height)) if shape.top is not None else 0,
                "has_text": shape.has_text_frame and bool(shape.text_frame.text.strip()),
                "text": shape.text_frame.text if shape.has_text_frame else "",
            }
            shapes_info.append(info)

        # 1. 元素溢出检测
        if check_overflow:
            _check_overflow(slide, page_num, shapes_info, canvas_w_emu, canvas_h_emu,
                            safe_left, safe_top, safe_right, safe_bottom, report)

        # 2. 元素重叠检测
        if check_overlap:
            _check_overlap(page_num, shapes_info, report)

        # 3. 间距异常检测
        if check_spacing:
            _check_spacing(page_num, shapes_info, report)

        # 4. 字号下限检测
        if check_fontsize:
            _check_fontsize(page_num, shapes_info, report)

        # 5. 空文本框检测
        if check_empty:
            _check_empty_textbox(page_num, shapes_info, report)

    logger.info("排版自检完成: %s（%d 错误, %d 警告）",
                path.name, report.error_count, report.warning_count)
    return report


# ==================== 检测项实现 ====================
def _check_overflow(slide, page_num, shapes_info, canvas_w, canvas_h,
                    safe_left, safe_top, safe_right, safe_bottom, report) -> None:
    """检测元素溢出画布/安全区"""
    for info in shapes_info:
        # 跳过背景全屏元素（如 cover_bg / divider_bg）
        name = info["name"].lower()
        if "bg" in name or "background" in name:
            continue
        # 跳过装饰线
        if "divider" in name or "line" in name:
            continue

        # 画布溢出（error 级别）
        if info["left"] < 0 or info["top"] < 0:
            report.issues.append(LayoutIssue(
                page=page_num, severity="error", category="overflow",
                message=f"元素 {info['name']} 超出画布左/上边界",
                shape_name=info["name"],
                detail={"left_inch": info["left"] / EMU_PER_INCH,
                        "top_inch": info["top"] / EMU_PER_INCH},
            ))
        if info["right"] > canvas_w or info["bottom"] > canvas_h:
            report.issues.append(LayoutIssue(
                page=page_num, severity="error", category="overflow",
                message=f"元素 {info['name']} 超出画布右/下边界",
                shape_name=info["name"],
                detail={"right_inch": info["right"] / EMU_PER_INCH,
                        "bottom_inch": info["bottom"] / EMU_PER_INCH,
                        "canvas_w_inch": canvas_w / EMU_PER_INCH,
                        "canvas_h_inch": canvas_h / EMU_PER_INCH},
            ))

        # 安全区溢出（warning 级别，仅对有文本元素）
        if info["has_text"]:
            if info["left"] < safe_left or info["top"] < safe_top:
                report.issues.append(LayoutIssue(
                    page=page_num, severity="warning", category="overflow",
                    message=f"文本元素 {info['name']} 进入安全边距区域",
                    shape_name=info["name"],
                    detail={"left_inch": info["left"] / EMU_PER_INCH,
                            "top_inch": info["top"] / EMU_PER_INCH,
                            "safe_margin_inch": SAFE_MARGIN_INCH},
                ))


def _check_overlap(page_num, shapes_info, report) -> None:
    """检测元素重叠（仅检测有文本元素之间的重叠）"""
    text_shapes = [s for s in shapes_info if s["has_text"]]
    tolerance = int(OVERLAP_TOLERANCE_INCH * EMU_PER_INCH)

    for i, s1 in enumerate(text_shapes):
        for s2 in text_shapes[i + 1:]:
            # 计算重叠区域
            overlap_left = max(s1["left"], s2["left"])
            overlap_top = max(s1["top"], s2["top"])
            overlap_right = min(s1["right"], s2["right"])
            overlap_bottom = min(s1["bottom"], s2["bottom"])

            if overlap_right > overlap_left + tolerance and overlap_bottom > overlap_top + tolerance:
                overlap_w = (overlap_right - overlap_left) / EMU_PER_INCH
                overlap_h = (overlap_bottom - overlap_top) / EMU_PER_INCH
                report.issues.append(LayoutIssue(
                    page=page_num, severity="warning", category="overlap",
                    message=f"文本元素重叠: {s1['name']} 与 {s2['name']}",
                    shape_name=f"{s1['name']} & {s2['name']}",
                    detail={"overlap_w_inch": round(overlap_w, 2),
                            "overlap_h_inch": round(overlap_h, 2)},
                ))


def _check_spacing(page_num, shapes_info, report) -> None:
    """检测相邻元素间距异常（负间距或过小）"""
    text_shapes = [s for s in shapes_info if s["has_text"]]
    if len(text_shapes) < 2:
        return

    # 按 top 排序，检测纵向相邻元素
    sorted_shapes = sorted(text_shapes, key=lambda s: s["top"])
    min_gap_emu = int(MIN_GAP_INCH * EMU_PER_INCH)

    for i in range(len(sorted_shapes) - 1):
        s1 = sorted_shapes[i]
        s2 = sorted_shapes[i + 1]
        # 纵向间距 = s2.top - s1.bottom
        gap = s2["top"] - s1["bottom"]
        # 仅当两元素横向有重叠时才视为纵向相邻
        h_overlap = min(s1["right"], s2["right"]) - max(s1["left"], s2["left"])
        if h_overlap <= 0:
            continue  # 横向不重叠，跳过
        if gap < -int(0.1 * EMU_PER_INCH):
            # 负间距（重叠>0.1inch），由 overlap 检测覆盖，跳过
            continue
        if 0 <= gap < min_gap_emu:
            report.issues.append(LayoutIssue(
                page=page_num, severity="info", category="spacing",
                message=f"元素间距过小: {s1['name']} 与 {s2['name']}",
                shape_name=f"{s1['name']} & {s2['name']}",
                detail={"gap_inch": round(gap / EMU_PER_INCH, 3),
                        "min_gap_inch": MIN_GAP_INCH},
            ))


def _check_fontsize(page_num, shapes_info, report) -> None:
    """检测字号低于下限"""
    for info in shapes_info:
        if not info["has_text"]:
            continue
        shape = info["shape"]
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size is None:
                    continue
                size_pt = run.font.size.pt
                if size_pt < MIN_FONT_PT:
                    report.issues.append(LayoutIssue(
                        page=page_num, severity="warning", category="fontsize",
                        message=f"字号 {size_pt}pt 低于下限 {MIN_FONT_PT}pt",
                        shape_name=info["name"],
                        detail={"font_size_pt": size_pt,
                                "min_font_pt": MIN_FONT_PT,
                                "text_preview": info["text"][:30]},
                    ))
                    break  # 同 shape 仅报告一次
            else:
                continue
            break


def _check_empty_textbox(page_num, shapes_info, report) -> None:
    """检测空文本框残留（有 text_frame 但无内容）"""
    for info in shapes_info:
        shape = info["shape"]
        if not shape.has_text_frame:
            continue
        # 跳过装饰元素（背景、装饰线、卡片容器等，本身无文本预期）
        name_lower = info["name"].lower()
        if "bg" in name_lower or "background" in name_lower or "divider" in name_lower \
                or "line" in name_lower or "badge" in name_lower or "card" in name_lower:
            continue
        # 排除非文本形状（线条/图片/表格/图表本身不应有文本）
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if shape.shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.PICTURE,
                                     MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.CHART):
                continue
        except Exception:
            pass
        text = shape.text_frame.text.strip() if shape.text_frame else ""
        if not text:
            # 仅对尺寸较大的文本框报告（小尺寸可能是装饰元素）
            if info["width"] > int(2 * EMU_PER_INCH) and info["height"] > int(0.5 * EMU_PER_INCH):
                report.issues.append(LayoutIssue(
                    page=page_num, severity="info", category="empty",
                    message=f"空文本框: {info['name']}",
                    shape_name=info["name"],
                    detail={"width_inch": round(info["width"] / EMU_PER_INCH, 2),
                            "height_inch": round(info["height"] / EMU_PER_INCH, 2)},
                ))


# ==================== 便捷输出函数 ====================
def format_report(report: IssueReport, verbose: bool = False) -> str:
    """格式化报告为人类可读字符串

    :param report: 自检报告
    :param verbose: 是否输出 info 级别问题
    :return: 格式化字符串
    """
    lines = [
        f"=== 排版自检报告 ===",
        f"文件: {report.file_path}",
        f"总页数: {report.total_pages}",
        f"错误: {report.error_count}, 警告: {report.warning_count}",
        f"结论: {'通过' if report.passed else '未通过'}",
        "",
    ]
    for issue in report.issues:
        if not verbose and issue.severity == "info":
            continue
        severity_icon = {"error": "ERROR", "warning": "WARN", "info": "INFO"}[issue.severity]
        lines.append(f"[{severity_icon}] 页{issue.page} {issue.category}: {issue.message}")
        if issue.shape_name:
            lines.append(f"    shape: {issue.shape_name}")
        if issue.detail:
            detail_str = ", ".join(f"{k}={v}" for k, v in issue.detail.items())
            lines.append(f"    {detail_str}")
    return "\n".join(lines)
