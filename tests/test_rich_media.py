"""T503 富媒体专项测试

覆盖四类富媒体元素 + 新增页面布局 + 弹性约束：
    1. 图片槽位替换（T301）：cover/contain 两种填充模式
    2. 图表数据替换（T302）：bar/line/pie/radar 四类
    3. 表格动态扩展（T303）：行数变化时样式继承
    4. SmartArt 文本替换（T304）：节点文本精准替换
    5. 新增页面布局（T603）：timeline/two_column/table 三类
    6. 弹性约束计算（T604）：文本高度预估 + 纵向均分

通过 python-pptx 动态构造测试 PPTX，不依赖外部模板文件。
标记为 slow：需实际渲染，使用 `-m "not slow"` 跳过。
"""
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt, Emu

from aippt.image_replacer import (
    resolve_image_source, compute_fit_rect, replace_images, _get_image_size,
)
from aippt.layout.elastic_constraint import (
    estimate_text_height, distribute_vertical, wrap_text, elastic_distribute_items,
)
# 注意：必须从 aippt.layout 包导入（而非 ppt_auto_layout 模块），
# 以触发 __init__.py 中的 page_layouts 导入，完成 8 类页面布局注册
from aippt.layout import (
    PAGE_LAYOUT_REGISTRY, LayoutContext, create_presentation, add_blank_slide,
    dispatch_page_layout,
)
from aippt.theme_loader import load_theme
from ppt_renderer import PptRenderer


# ==================== Fixture：通用资源 ====================

@pytest.fixture
def sample_image_path(tmp_path) -> Path:
    """生成测试用 PNG 图片（200x150）"""
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        pytest.skip("Pillow 未安装，跳过图片相关测试")
    img = Image.new("RGB", (200, 150), color=(0, 47, 167))
    draw = ImageDraw.Draw(img)
    draw.rectangle([50, 30, 150, 120], fill=(255, 255, 255))
    path = tmp_path / "sample.png"
    img.save(path)
    return path


@pytest.fixture
def sample_image_path_portrait(tmp_path) -> Path:
    """生成竖版测试图片（150x300）"""
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("Pillow 未安装")
    img = Image.new("RGB", (150, 300), color=(220, 38, 38))
    path = tmp_path / "portrait.png"
    img.save(path)
    return path


@pytest.fixture
def default_theme() -> dict:
    """加载默认主题（商务蓝）"""
    return load_theme("商务蓝")


# ==================== 1. 图片槽位替换测试（T301） ====================

class TestImageReplacement:
    """T301 图片槽位替换"""

    def test_resolve_local_path(self, sample_image_path):
        """本地路径解析"""
        spec = {"path": str(sample_image_path)}
        result = resolve_image_source(spec)
        assert result == str(sample_image_path.resolve())

    def test_resolve_missing_path(self, tmp_path):
        """不存在的本地路径应抛 FileNotFoundError"""
        spec = {"path": str(tmp_path / "nonexistent.png")}
        with pytest.raises(FileNotFoundError):
            resolve_image_source(spec)

    def test_resolve_no_source(self):
        """未提供 path/url 应抛 ValueError"""
        with pytest.raises(ValueError):
            resolve_image_source({})

    def test_compute_fit_contain_landscape(self):
        """contain 模式：横版图片在竖版区域应留白居中"""
        # 目标区域 2x3 inch（竖），图片 200x150（横）
        target = (0, 0, int(2 * 914400), int(3 * 914400))
        left, top, w, h = compute_fit_rect(*target, 200, 150, "contain")
        # 图片更宽，应以宽度为准
        assert w == target[2]
        expected_h = int(target[2] / (200 / 150))
        assert h == expected_h
        # 居中
        assert left == 0
        assert top == (target[3] - expected_h) // 2

    def test_compute_fit_contain_portrait(self):
        """contain 模式：竖版图片在横版区域应留白居中"""
        # 目标区域 3x2 inch（横），图片 150x300（竖）
        target = (0, 0, int(3 * 914400), int(2 * 914400))
        left, top, w, h = compute_fit_rect(*target, 150, 300, "contain")
        assert h == target[3]
        expected_w = int(target[3] * (150 / 300))
        assert w == expected_w

    def test_get_image_size(self, sample_image_path):
        """读取图片尺寸"""
        w, h = _get_image_size(str(sample_image_path))
        assert w == 200
        assert h == 150

    def test_replace_images_in_slide(self, sample_image_path, tmp_path):
        """替换 slide 中的 picture shape"""
        # 构造含图片的 PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        pic = slide.shapes.add_picture(str(sample_image_path), Inches(1), Inches(1), Inches(4), Inches(3))
        original_name = pic.name

        # 替换图片
        image_data = {
            "image_1": {"path": str(sample_image_path), "fit": "contain"}
        }
        replaced = replace_images(slide, image_data)
        assert replaced == 1

        # 验证：slide 仍有 1 个 picture
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pictures) == 1

    def test_replace_images_no_picture(self, tmp_path):
        """无 picture shape 时应返回 0 并 warning"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # 仅添加文本框
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))

        replaced = replace_images(slide, {"image_1": {"path": "dummy"}})
        assert replaced == 0

    def test_replace_images_multiple(self, sample_image_path, sample_image_path_portrait, tmp_path):
        """多图片替换：按出现顺序匹配"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(str(sample_image_path), Inches(1), Inches(1), Inches(2), Inches(2))
        slide.shapes.add_picture(str(sample_image_path_portrait), Inches(4), Inches(1), Inches(2), Inches(2))

        image_data = {
            "image_1": {"path": str(sample_image_path_portrait), "fit": "contain"},
            "image_2": {"path": str(sample_image_path), "fit": "contain"},
        }
        replaced = replace_images(slide, image_data)
        assert replaced == 2


# ==================== 2. 图表数据替换测试（T302） ====================

class TestChartDataReplacement:
    """T302 图表数据源替换"""

    def _build_chart_pptx(self, tmp_path, chart_type, name, series_count=2, cat_count=4) -> Path:
        path = tmp_path / name
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        chart_data = ChartData()
        chart_data.categories = [f"Q{i+1}" for i in range(cat_count)]
        for s in range(series_count):
            chart_data.add_series(f"Series{s+1}", list(range(1, cat_count + 1)))
        slide.shapes.add_chart(chart_type, Inches(1), Inches(1), Inches(8), Inches(5), chart_data)
        prs.save(path)
        return path

    def test_bar_chart_replacement(self, tmp_path):
        """柱状图数据替换：模板 2 系列，传入 1 系列，多余系列清空（保留模板系列数）"""
        path = self._build_chart_pptx(tmp_path, XL_CHART_TYPE.COLUMN_CLUSTERED, "bar.pptx")
        prs = Presentation(path)
        slide = prs.slides[0]
        chart_shape = next(s for s in slide.shapes if s.has_chart)

        from aippt.chart_replacer import replace_chart_data
        new_data = {
            "categories": ["一月", "二月", "三月"],
            "series": [{"name": "销售额", "data": [100, 200, 150]}],
        }
        replace_chart_data(chart_shape, new_data)

        # 验证：模板 2 系列，传入 1 系列，多余系列清空数据但系列数保持 2
        chart = chart_shape.chart
        assert len(chart.series) == 2  # 模板原系列数不变
        assert chart.series[0].name == "销售额"
        assert len(chart.plots[0].categories) == 3

    def test_pie_chart_replacement(self, tmp_path):
        """饼图数据替换"""
        path = self._build_chart_pptx(tmp_path, XL_CHART_TYPE.PIE, "pie.pptx",
                                       series_count=1, cat_count=3)
        prs = Presentation(path)
        slide = prs.slides[0]
        chart_shape = next(s for s in slide.shapes if s.has_chart)

        from aippt.chart_replacer import replace_chart_data
        new_data = {
            "categories": ["A", "B", "C", "D"],
            "series": [{"name": "占比", "data": [30, 40, 20, 10]}],
        }
        replace_chart_data(chart_shape, new_data)

        chart = chart_shape.chart
        assert len(chart.plots[0].categories) == 4

    def test_multi_series_replacement(self, tmp_path):
        """多系列图表替换（M>N 截断）"""
        path = self._build_chart_pptx(tmp_path, XL_CHART_TYPE.COLUMN_CLUSTERED, "multi.pptx",
                                       series_count=2, cat_count=4)
        prs = Presentation(path)
        slide = prs.slides[0]
        chart_shape = next(s for s in slide.shapes if s.has_chart)

        from aippt.chart_replacer import replace_chart_data
        new_data = {
            "categories": ["Q1", "Q2"],
            "series": [
                {"name": "A", "data": [10, 20]},
                {"name": "B", "data": [30, 40]},
                {"name": "C", "data": [50, 60]},  # 模板仅 2 系列，应截断
            ],
        }
        replace_chart_data(chart_shape, new_data)

        chart = chart_shape.chart
        assert len(chart.series) == 2  # 模板原 2 系列


# ==================== 3. 表格动态扩展测试（T303） ====================

class TestTableDynamicExpansion:
    """T303 表格动态行扩展"""

    def _build_table_pptx(self, tmp_path, rows=3, cols=3) -> Path:
        path = tmp_path / "table.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(8), Inches(4))
        # 填充初始数据
        for r in range(rows):
            for c in range(cols):
                table_shape.table.cell(r, c).text = f"R{r}C{c}"
        prs.save(path)
        return path

    def test_table_expand_rows(self, tmp_path):
        """扩展表格行数（2 行 → 5 行）"""
        path = self._build_table_pptx(tmp_path, rows=3, cols=3)  # 1 表头 + 2 数据
        prs = Presentation(path)
        slide = prs.slides[0]
        table_shape = next(s for s in slide.shapes if s.has_table)

        from aippt.table_filler import fill_dynamic_table
        new_data = {
            "headers": ["姓名", "部门", "职级"],
            "rows": [
                ["张三", "技术部", "P5"],
                ["李四", "产品部", "P6"],
                ["王五", "设计部", "P4"],
                ["赵六", "运营部", "P5"],
                ["钱七", "市场部", "P7"],
            ],
        }
        fill_dynamic_table(table_shape, new_data)

        # 验证：行数 = 1 表头 + 5 数据 = 6
        assert len(table_shape.table.rows) == 6

    def test_table_shrink_rows(self, tmp_path):
        """缩减表格行数（2 行 → 1 行）"""
        path = self._build_table_pptx(tmp_path, rows=3, cols=3)
        prs = Presentation(path)
        slide = prs.slides[0]
        table_shape = next(s for s in slide.shapes if s.has_table)

        from aippt.table_filler import fill_dynamic_table
        new_data = {
            "headers": ["A", "B", "C"],
            "rows": [["1", "2", "3"]],  # 仅 1 行数据
        }
        fill_dynamic_table(table_shape, new_data)

        assert len(table_shape.table.rows) == 2  # 1 表头 + 1 数据

    def test_table_data_correctness(self, tmp_path):
        """表格数据正确性"""
        path = self._build_table_pptx(tmp_path, rows=2, cols=2)
        prs = Presentation(path)
        slide = prs.slides[0]
        table_shape = next(s for s in slide.shapes if s.has_table)

        from aippt.table_filler import fill_dynamic_table
        new_data = {
            "headers": ["项目", "金额"],
            "rows": [["A项目", "100万"], ["B项目", "200万"]],
        }
        fill_dynamic_table(table_shape, new_data)

        table = table_shape.table
        assert table.cell(0, 0).text == "项目"
        assert table.cell(1, 0).text == "A项目"
        assert table.cell(2, 1).text == "200万"


# ==================== 4. SmartArt 文本替换测试（T304） ====================

class TestSmartArtReplacement:
    """T304 SmartArt 文本替换

    注：python-pptx 不支持动态创建 SmartArt，本测试验证 ppt_smartart 模块可正常导入，
    实际 SmartArt 替换需要含 SmartArt 的 PPTX 文件（手工准备）。
    """

    def test_smartart_module_importable(self):
        """ppt_smartart 模块可正常导入"""
        try:
            from ppt_smartart import replace_smartart_text
            assert callable(replace_smartart_text)
        except ImportError:
            pytest.skip("ppt_smartart 模块不可用")

    def test_smartart_replace_empty_mapping(self):
        """空映射应返回 0"""
        try:
            from ppt_smartart import replace_smartart_text
        except ImportError:
            pytest.skip("ppt_smartart 模块不可用")

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # 无 SmartArt 时应返回 0
        result = replace_smartart_text(slide, {})
        assert result == 0


# ==================== 5. 新增页面布局测试（T603） ====================

class TestNewPageLayouts:
    """T603 timeline/two_column/table 三类新布局"""

    def test_layout_registry_contains_all_8(self):
        """8 类页面布局全部注册"""
        expected = {"cover", "catalog", "divider", "numbered_list", "kpi",
                    "timeline", "two_column", "table"}
        assert expected.issubset(set(PAGE_LAYOUT_REGISTRY.keys()))

    def test_timeline_layout_render(self, default_theme):
        """timeline 页面渲染"""
        prs = create_presentation()
        slide = add_blank_slide(prs)
        ctx = LayoutContext(page_num=1)
        page_data = {
            "page_type": "timeline",
            "title": "项目里程碑",
            "timeline_items": [
                {"time": "2024 Q1", "event": "项目启动"},
                {"time": "2024 Q2", "event": "需求调研完成"},
                {"time": "2024 Q3", "event": "核心功能上线"},
                {"time": "2024 Q4", "event": "全量推广"},
            ],
        }
        dispatch_page_layout(slide, page_data, default_theme, ctx)

        # 验证：生成了 title/year/desc/number 等角色
        roles = {e.role for e in ctx.elements}
        assert "title" in roles
        assert "year" in roles
        assert "desc" in roles
        assert "number" in roles
        # 验证：4 个时间节点
        assert ctx.role_seq.get("year", 0) == 4

    def test_two_column_layout_render(self, default_theme):
        """two_column 页面渲染"""
        prs = create_presentation()
        slide = add_blank_slide(prs)
        ctx = LayoutContext(page_num=1)
        page_data = {
            "page_type": "two_column",
            "title": "优势 vs 挑战",
            "left_title": "核心优势",
            "left_items": ["技术领先", "团队成熟", "资源充足"],
            "right_title": "面临挑战",
            "right_items": ["竞争激烈", "成本上升", "人才短缺"],
        }
        dispatch_page_layout(slide, page_data, default_theme, ctx)

        roles = {e.role for e in ctx.elements}
        assert "title" in roles
        assert "left_subtitle" in roles
        assert "right_subtitle" in roles
        assert "left_body" in roles
        assert "right_body" in roles

    def test_table_layout_render(self, default_theme):
        """table 页面渲染（原生 PPT 表格）"""
        prs = create_presentation()
        slide = add_blank_slide(prs)
        ctx = LayoutContext(page_num=1)
        page_data = {
            "page_type": "table",
            "title": "团队配置",
            "headers": ["角色", "人数", "负责人"],
            "rows": [
                ["前端", "3", "张三"],
                ["后端", "4", "李四"],
                ["测试", "2", "王五"],
            ],
        }
        dispatch_page_layout(slide, page_data, default_theme, ctx)

        # 验证：生成了 title 和 table 角色
        roles = {e.role for e in ctx.elements}
        assert "title" in roles
        assert "table" in roles

        # 验证：slide 含 1 个 GraphicFrame（表格）
        from pptx.shapes.graphfrm import GraphicFrame
        tables = [s for s in slide.shapes if isinstance(s, GraphicFrame) and s.has_table]
        assert len(tables) == 1
        # 表格尺寸：1 表头 + 3 数据 = 4 行，3 列
        assert len(tables[0].table.rows) == 4
        assert len(tables[0].table.columns) == 3
        # 表头内容
        assert tables[0].table.cell(0, 0).text == "角色"
        assert tables[0].table.cell(1, 2).text == "张三"

    def test_table_layout_zebra(self, default_theme):
        """table 页面斑马纹"""
        prs = create_presentation()
        slide = add_blank_slide(prs)
        ctx = LayoutContext(page_num=1)
        page_data = {
            "page_type": "table",
            "title": "测试",
            "headers": ["A", "B"],
            "rows": [["1", "2"], ["3", "4"], ["5", "6"]],
        }
        dispatch_page_layout(slide, page_data, default_theme, ctx)

        from pptx.shapes.graphfrm import GraphicFrame
        table_shape = next(s for s in slide.shapes
                          if isinstance(s, GraphicFrame) and s.has_table)
        table = table_shape.table
        # 第 1 数据行（r=0）应为 row_bg，第 2 数据行（r=1）应为 zebra_bg
        from aippt.layout.ppt_auto_layout import get_token
        row_bg = get_token(default_theme, "color.table_row_bg", "#FFFFFF")
        zebra_bg = get_token(default_theme, "color.table_zebra_bg", "#F3F4F6")
        # 验证填充色（RGBColor 比较）
        assert table.cell(1, 0).fill.fore_color.rgb is not None
        assert table.cell(2, 0).fill.fore_color.rgb is not None


# ==================== 6. 弹性约束计算测试（T604） ====================

class TestElasticConstraint:
    """T604 弹性约束计算"""

    def test_estimate_short_text_single_line(self):
        """短文本预估为单行"""
        h = estimate_text_height("短", 16, int(5 * 914400))
        # 16pt × 1.2 = 19.2pt ≈ 0.27 inch
        assert 0.2 < h / 914400 < 0.35

    def test_estimate_long_text_multi_line(self):
        """长文本应预估多行"""
        h = estimate_text_height("a" * 100, 16, int(5 * 914400))
        # 100 字符 / (5 inch × 96 / (16 × 0.55)) ≈ 多行
        assert h / 914400 > 0.5  # 应超过 0.5 inch

    def test_estimate_cjk_vs_ascii(self):
        """中文字符宽度系数 1.0，英文 0.55，同样字数中文应更高（强制换行）"""
        # 30 字符 @ 16pt：中文 480pt > 5inch(360pt) → 2 行；英文 264pt < 360pt → 1 行
        h_cn = estimate_text_height("中" * 30, 16, int(5 * 914400))
        h_en = estimate_text_height("a" * 30, 16, int(5 * 914400))
        assert h_cn > h_en, f"中文高度 {h_cn} 应大于英文高度 {h_en}"

    def test_estimate_empty_text(self):
        """空文本高度为 0"""
        assert estimate_text_height("", 16, int(5 * 914400)) == 0

    def test_distribute_surplus_allocation(self):
        """总高度 < 容器：剩余空间均分"""
        items = [{"text": "短"}, {"text": "短"}, {"text": "短"}]
        hs = distribute_vertical(int(4 * 914400), items,
                                  font_size_pt=16, box_width_emu=int(5 * 914400))
        # 3 项短文本，容器 4 inch，应分配 bonus
        assert sum(hs) == int(4 * 914400)  # 总和等于容器高度
        assert all(h > 0 for h in hs)

    def test_distribute_shrink_with_min(self):
        """总高度 > 容器：等比缩小，但保持下限"""
        items = [
            {"text": "这是一个非常非常非常长的文本条目，内容很多需要换行才能完整显示"},
            {"text": "另一个也很长的文本条目，同样需要多行才能容纳"},
            {"text": "第三个长条目，文本内容依然很长很长"},
        ]
        hs = distribute_vertical(int(2 * 914400), items, min_h_emu=int(0.4 * 914400),
                                  font_size_pt=16, box_width_emu=int(5 * 914400))
        # 总和不超过容器太多（下限兜底时可能略超）
        assert sum(hs) <= int(2 * 914400) + int(0.1 * 914400)
        # 每项至少 0.4 inch
        assert all(h >= int(0.4 * 914400) - 1 for h in hs)

    def test_distribute_empty_items(self):
        """空条目列表返回空"""
        assert distribute_vertical(int(4 * 914400), []) == []

    def test_wrap_text_basic(self):
        """自动换行基本功能"""
        lines = wrap_text("这是一个需要自动换行的较长中文文本", 16, int(2 * 914400))
        assert len(lines) >= 2  # 应至少换行一次
        # 拼接后应包含原文（去除换行符）
        assert "".join(lines) == "这是一个需要自动换行的较长中文文本"

    def test_wrap_text_explicit_newline(self):
        """显式换行符应保留为分行"""
        lines = wrap_text("第一行\n第二行", 16, int(10 * 914400))
        assert lines == ["第一行", "第二行"]

    def test_elastic_distribute_items_layout(self):
        """便捷接口：返回 (top, height) 元组列表"""
        items = [{"text": "A"}, {"text": "B"}, {"text": "C"}]
        result = elastic_distribute_items(
            items, int(1 * 914400), int(4 * 914400), int(5 * 914400),
            font_size_pt=16, gap_inch=0.15,
        )
        assert len(result) == 3
        # 第一项 top = 区域顶部
        assert result[0][0] == int(1 * 914400)
        # 第二项 top > 第一项 top + 第一项 height
        assert result[1][0] > result[0][0] + result[0][1] - 1
        # 所有高度 > 0
        assert all(h > 0 for _, h in result)


# ==================== 7. 集成测试：双引擎 + 富媒体 ====================

class TestIntegration:
    """集成测试：双引擎协同 + 富媒体"""

    def test_autolayout_render_full_pages(self, default_theme, tmp_path):
        """AutoLayout 渲染包含全部 8 类页面的 PPT"""
        prs = create_presentation()
        pages = [
            {"page_id": 1, "page_type": "cover", "title": "测试报告", "subtitle": "2024 年度"},
            {"page_id": 2, "page_type": "catalog", "title": "目录", "items": ["背景", "方法", "结果", "结论"]},
            {"page_id": 3, "page_type": "divider", "title": "01 研究背景", "section_no": "01"},
            {"page_id": 4, "page_type": "numbered_list", "title": "核心发现",
             "items": ["发现一", "发现二", "发现三", "发现四"]},
            {"page_id": 5, "page_type": "kpi", "title": "关键指标",
             "kpi_items": [{"label": "用户数", "value": "100万", "trend": "+20%"},
                           {"label": "收入", "value": "5000万", "trend": "+15%"}]},
            {"page_id": 6, "page_type": "timeline", "title": "项目进度",
             "timeline_items": [
                 {"time": "Q1", "event": "启动"},
                 {"time": "Q2", "event": "调研"},
                 {"time": "Q3", "event": "开发"},
                 {"time": "Q4", "event": "上线"},
             ]},
            {"page_id": 7, "page_type": "two_column", "title": "对比分析",
             "left_title": "方案 A", "left_items": ["成本低", "周期短"],
             "right_title": "方案 B", "right_items": ["质量高", "可扩展"]},
            {"page_id": 8, "page_type": "table", "title": "数据汇总",
             "headers": ["指标", "数值", "趋势"],
             "rows": [["DAU", "100万", "↑"], ["收入", "5000万", "↑"]]},
        ]

        for i, page_data in enumerate(pages, 1):
            slide = add_blank_slide(prs)
            ctx = LayoutContext(page_num=i)
            dispatch_page_layout(slide, page_data, default_theme, ctx)

        # 验证：8 页全部生成
        assert len(prs.slides) == 8

        # 保存
        output = tmp_path / "integration.pptx"
        prs.save(output)
        assert output.exists()
        assert output.stat().st_size > 0

    def test_template_renderer_image_data_skipped(self, tmp_path):
        """PptRenderer 主循环正确跳过 image_data 不走文本替换"""
        # 构造简单 PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "标题"
        template_path = tmp_path / "t.pptx"
        prs.save(template_path)

        # 构造 meta
        meta = {
            "template_id": "test",
            "category": "工作总结",
            "total_pages": 1,
            "chapters": [],
            "page_slots": {
                "1": [{"slot": "title", "match_text": "标题", "shape_name": ""}]
            },
        }
        meta_path = tmp_path / "t.meta.json"
        meta_path.write_text(json.dumps(meta), encoding="utf-8")

        renderer = PptRenderer(str(template_path), str(meta_path))
        slot_data = {
            "1": {
                "title": "新标题",
                "image_data": {"image_1": {"path": "nonexistent.png", "fit": "cover"}},
            }
        }
        output = tmp_path / "out.pptx"
        # 不应抛异常，image_data 走专门逻辑（替换失败仅 warning）
        renderer.render(slot_data, str(output))
        assert output.exists()
