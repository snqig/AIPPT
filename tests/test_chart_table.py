"""3.2.3 图表表格专项测试

覆盖 4 类图表（bar/line/pie/radar）数据源替换 + 表格动态行扩展 + 不同数据量场景。
通过 python-pptx 动态构造含图表/表格的测试 PPTX，不依赖外部模板文件。

标记为 slow：需实际渲染 PPTX，耗时较长，使用 `-m "not slow"` 跳过。
"""
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from insert_tables import (
    list_charts_and_tables,
    run_chart_replacement,
    run_table_replacement,
)


# ==================== Fixture：动态构造含图表/表格的 PPTX ====================

@pytest.fixture
def chart_pptx_bar(tmp_path) -> Path:
    """构造含柱状图的 PPTX（2 系列 × 4 分类）"""
    return _build_chart_pptx(tmp_path, XL_CHART_TYPE.COLUMN_CLUSTERED, "bar_chart.pptx",
                              series_count=2, cat_count=4)


@pytest.fixture
def chart_pptx_line(tmp_path) -> Path:
    """构造含折线图的 PPTX（1 系列 × 5 分类）"""
    return _build_chart_pptx(tmp_path, XL_CHART_TYPE.LINE, "line_chart.pptx",
                              series_count=1, cat_count=5)


@pytest.fixture
def chart_pptx_pie(tmp_path) -> Path:
    """构造含饼图的 PPTX（1 系列 × 3 分类）"""
    return _build_chart_pptx(tmp_path, XL_CHART_TYPE.PIE, "pie_chart.pptx",
                              series_count=1, cat_count=3)


@pytest.fixture
def chart_pptx_radar(tmp_path) -> Path:
    """构造含雷达图的 PPTX（2 系列 × 4 分类）"""
    return _build_chart_pptx(tmp_path, XL_CHART_TYPE.RADAR, "radar_chart.pptx",
                              series_count=2, cat_count=4)


@pytest.fixture
def table_pptx(tmp_path) -> Path:
    """构造含表格的 PPTX（表头 + 2 数据行，3 列）"""
    path = tmp_path / "table.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    rows, cols = 3, 3  # 1 header + 2 data rows
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1),
                                          Inches(8), Inches(3))
    table = table_shape.table
    # 表头
    for c in range(cols):
        table.cell(0, c).text = f"列{c+1}"
    # 数据行
    for r in range(1, rows):
        for c in range(cols):
            table.cell(r, c).text = f"r{r}c{c}"

    prs.save(str(path))
    return path


def _build_chart_pptx(tmp_path: Path, chart_type: XL_CHART_TYPE,
                       filename: str, series_count: int, cat_count: int) -> Path:
    """构造含指定类型图表的 PPTX"""
    path = tmp_path / filename
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    chart_data = ChartData()
    chart_data.categories = [f"Q{i+1}" for i in range(cat_count)]
    for s in range(series_count):
        chart_data.add_series(f"系列{s+1}", [10 * (s + 1) + i for i in range(cat_count)])

    slide.shapes.add_chart(chart_type, Inches(1), Inches(1), Inches(8), Inches(4.5),
                            chart_data)
    prs.save(str(path))
    return path


# ==================== 图表数据替换测试（4 类） ====================

@pytest.mark.slow
def test_chart_bar_replacement(chart_pptx_bar, tmp_path):
    """柱状图数据源替换：2系列→2系列，4分类→4分类"""
    out = tmp_path / "bar_out.pptx"
    chart_data = {
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [
            {"name": "营收", "data": [1200, 1500, 1800, 2100]},
            {"name": "利润", "data": [300, 450, 600, 800]},
        ],
    }
    result = run_chart_replacement(chart_pptx_bar, chart_data, out)
    assert result["test_pass"], f"柱状图替换失败: {result.get('error')}"
    assert result["before"]["chart_type"] == "bar"
    assert result["after"]["chart_type"] == "bar"
    assert result["after"]["series_count"] == 2
    assert result["after"]["category_count"] == 4
    assert out.exists()


@pytest.mark.slow
def test_chart_line_replacement(chart_pptx_line, tmp_path):
    """折线图数据源替换：1系列→1系列，5分类→5分类"""
    out = tmp_path / "line_out.pptx"
    chart_data = {
        "categories": ["1月", "2月", "3月", "4月", "5月"],
        "series": [{"name": "用户增长", "data": [100, 250, 480, 760, 1100]}],
    }
    result = run_chart_replacement(chart_pptx_line, chart_data, out)
    assert result["test_pass"], f"折线图替换失败: {result.get('error')}"
    assert result["before"]["chart_type"] == "line"
    assert result["after"]["chart_type"] == "line"
    assert result["after"]["series_count"] == 1


@pytest.mark.slow
def test_chart_pie_replacement(chart_pptx_pie, tmp_path):
    """饼图数据源替换：1系列，3分类→3分类"""
    out = tmp_path / "pie_out.pptx"
    chart_data = {
        "categories": ["产品A", "产品B", "产品C"],
        "series": [{"name": "销量占比", "data": [45, 30, 25]}],
    }
    result = run_chart_replacement(chart_pptx_pie, chart_data, out)
    assert result["test_pass"], f"饼图替换失败: {result.get('error')}"
    assert result["before"]["chart_type"] == "pie"
    assert result["after"]["chart_type"] == "pie"


@pytest.mark.slow
def test_chart_radar_replacement(chart_pptx_radar, tmp_path):
    """雷达图数据源替换：2系列×4分类"""
    out = tmp_path / "radar_out.pptx"
    chart_data = {
        "categories": ["技术", "沟通", "管理", "创新"],
        "series": [
            {"name": "当前能力", "data": [80, 65, 70, 75]},
            {"name": "目标能力", "data": [95, 85, 90, 92]},
        ],
    }
    result = run_chart_replacement(chart_pptx_radar, chart_data, out)
    assert result["test_pass"], f"雷达图替换失败: {result.get('error')}"
    assert result["before"]["chart_type"] == "radar"


# ==================== 多系列适配测试 ====================

@pytest.mark.slow
def test_chart_series_expansion(chart_pptx_bar, tmp_path):
    """图表系列数变化：模板2系列，传入1系列（多余系列清空）"""
    out = tmp_path / "series_shrink.pptx"
    chart_data = {
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [{"name": "营收", "data": [1200, 1500, 1800, 2100]}],
    }
    result = run_chart_replacement(chart_pptx_bar, chart_data, out)
    assert result["test_pass"], f"系列缩减测试失败: {result.get('error')}"
    # 模板 2 系列，传入 1 系列，应保持 2 系列（多余清零）
    assert result["after"]["series_count"] == 2


@pytest.mark.slow
def test_chart_category_change(chart_pptx_bar, tmp_path):
    """图表分类数变化：模板4分类，传入6分类"""
    out = tmp_path / "cat_expand.pptx"
    chart_data = {
        "categories": ["Q1", "Q2", "Q3", "Q4", "Q5", "Q6"],
        "series": [
            {"name": "营收", "data": [1200, 1500, 1800, 2100, 2400, 2700]},
            {"name": "利润", "data": [300, 450, 600, 800, 950, 1100]},
        ],
    }
    result = run_chart_replacement(chart_pptx_bar, chart_data, out)
    assert result["test_pass"], f"分类扩展测试失败: {result.get('error')}"
    assert result["after"]["category_count"] == 6


# ==================== 表格动态行扩展测试 ====================

@pytest.mark.slow
def test_table_row_expansion(table_pptx, tmp_path):
    """表格行扩展：模板3行（1表头+2数据），传入5数据行 → 扩展至6行"""
    out = tmp_path / "table_expand.pptx"
    table_data = {
        "headers": ["产品", "价格", "销量"],
        "rows": [
            ["产品A", "99元", "1.2万"],
            ["产品B", "199元", "0.8万"],
            ["产品C", "299元", "0.6万"],
            ["产品D", "399元", "0.4万"],
            ["产品E", "499元", "0.3万"],
        ],
    }
    result = run_table_replacement(table_pptx, table_data, out)
    assert result["test_pass"], f"表格扩展失败: {result.get('error')}"
    assert result["before"]["rows"] == 3
    assert result["after"]["rows"] == 6  # 1 表头 + 5 数据
    assert result["rows_delta"] == 3
    assert out.exists()


@pytest.mark.slow
def test_table_row_shrink(table_pptx, tmp_path):
    """表格行缩减：模板3行，传入1数据行 → 缩减至2行"""
    out = tmp_path / "table_shrink.pptx"
    table_data = {
        "headers": ["产品", "价格", "销量"],
        "rows": [["产品A", "99元", "1.2万"]],
    }
    result = run_table_replacement(table_pptx, table_data, out)
    assert result["test_pass"], f"表格缩减失败: {result.get('error')}"
    assert result["after"]["rows"] == 2  # 1 表头 + 1 数据
    assert result["rows_delta"] == -1


@pytest.mark.slow
def test_table_empty_rows(table_pptx, tmp_path):
    """表格空数据行：传入0数据行（仅保留表头）"""
    out = tmp_path / "table_empty.pptx"
    table_data = {
        "headers": ["产品", "价格", "销量"],
        "rows": [],
    }
    result = run_table_replacement(table_pptx, table_data, out)
    assert result["test_pass"], f"空数据测试失败: {result.get('error')}"
    assert result["after"]["rows"] == 1  # 仅表头


@pytest.mark.slow
def test_table_large_data(table_pptx, tmp_path):
    """表格大数据量：传入10数据行"""
    out = tmp_path / "table_large.pptx"
    rows = [[f"产品{i+1}", f"{(i+1)*100}元", f"{(10-i)*0.1:.1f}万"] for i in range(10)]
    table_data = {
        "headers": ["产品", "价格", "销量"],
        "rows": rows,
    }
    result = run_table_replacement(table_pptx, table_data, out)
    assert result["test_pass"], f"大数据量测试失败: {result.get('error')}"
    assert result["after"]["rows"] == 11  # 1 表头 + 10 数据


# ==================== list 命令测试 ====================

def test_list_charts_and_tables(chart_pptx_bar, table_pptx):
    """list 命令正确识别图表与表格"""
    chart_info = list_charts_and_tables(chart_pptx_bar)
    assert chart_info["total"] >= 1
    assert chart_info["charts"][0]["chart_type"] == "bar"
    assert chart_info["charts"][0]["series_count"] == 2

    table_info = list_charts_and_tables(table_pptx)
    assert table_info["total"] >= 1
    assert table_info["tables"][0]["rows"] == 3
    assert table_info["tables"][0]["cols"] == 3


def test_list_no_charts(tmp_path):
    """list 命令对无图表/表格的 PPTX 返回空"""
    path = tmp_path / "empty.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.save(str(path))

    info = list_charts_and_tables(path)
    assert info["total"] == 0
    assert info["charts"] == []
    assert info["tables"] == []


# ==================== 错误场景测试 ====================

def test_chart_replacement_no_chart(tmp_path):
    """无图表时返回明确错误"""
    path = tmp_path / "no_chart.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.save(str(path))

    result = run_chart_replacement(path, {"categories": [], "series": []}, tmp_path / "out.pptx")
    assert not result["test_pass"]
    assert "未找到图表" in result["error"]


def test_table_replacement_no_table(tmp_path):
    """无表格时返回明确错误"""
    path = tmp_path / "no_table.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.save(str(path))

    result = run_table_replacement(path, {"headers": [], "rows": []}, tmp_path / "out.pptx")
    assert not result["test_pass"]
    assert "未找到表格" in result["error"]
