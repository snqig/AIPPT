"""
PPT 元素分类器模块

基于字号 / 位置 / 格式特征对幻灯片内的形状（shape）进行角色分类，
输出 role / confidence / reason 三元组，并对低置信度元素产生告警。

角色取值：
    title | subtitle | body | kpi_value | kpi_label | note | decorative

判定规则（详见 classify_element）：
    - 字号 ≥ 36 且 top < 30% 高度 → title（0.9）
    - 字号 24-32 且 top < 40% → subtitle（0.85）
    - 字号 < 18 且 top > 70% → note（0.8）
    - 纯数字或含 %/$ → kpi_value（0.9）
    - 紧邻 kpi_value 的短文本（< 10 字）→ kpi_label（0.85，由 classify_page 二次判定）
    - 字号 18-24 且 run.text 长度 > 15 → body（0.85）
    - 置信度 < 0.6 → decorative（0.5）
"""
import re
from typing import Any, Iterator

from aippt.logger import logger


# 16:9 默认幻灯片尺寸（EMU），用于反推失败时的兜底
_DEFAULT_SLIDE_WIDTH_EMU = 12192000
_DEFAULT_SLIDE_HEIGHT_EMU = 6858000

# 紧邻判定的最大垂直间距（EMU），0.5 英寸
_ADJACENT_GAP_EMU = 457200

# 纯数字（允许前导负号、小数点、千分位逗号）
_NUMERIC_RE = re.compile(r'^-?[\d,]+(\.\d+)?$')


# ==================== 内部工具函数 ====================
def _safe_get(obj: Any, attr: str, default: Any = None) -> Any:
    """安全读取对象属性，异常或不存在时返回默认值"""
    try:
        return getattr(obj, attr, default)
    except Exception:
        return default


def _get_shape_text(shape: Any) -> str:
    """读取形状文本（去首尾空白），无文本框返回空串"""
    try:
        if shape.has_text_frame:
            return shape.text_frame.text.strip()
    except Exception:
        return ""
    return ""


def _get_font_size_pt(shape: Any) -> Any:
    """
    提取形状字号（pt）
    优先取首个 run 的显式字号，回退到段落级 defRPr，均无则返回 None
    """
    try:
        if not shape.has_text_frame:
            return None
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    return run.font.size.pt
            # 段落级默认字号
            if para.font.size is not None:
                return para.font.size.pt
    except Exception:
        return None
    return None


def _get_presentation(pkg: Any) -> Any:
    """从 Package 中查找 PresentationPart 并返回 Presentation 对象"""
    try:
        from pptx.parts.presentation import PresentationPart
        parts = list(pkg.iter_parts()) if hasattr(pkg, 'iter_parts') else []
        for part in parts:
            if isinstance(part, PresentationPart):
                return part.presentation
    except Exception as e:
        logger.debug("查找 PresentationPart 失败: %s", e)
    return None


def _get_slide_dimensions(shape: Any) -> tuple[int, int]:
    """从形状反推所属幻灯片尺寸（EMU），失败回退 16:9 默认值"""
    try:
        pkg = shape.part.package
        prs = _get_presentation(pkg)
        if prs is not None and prs.slide_width and prs.slide_height:
            return int(prs.slide_width), int(prs.slide_height)
    except Exception:
        pass
    return _DEFAULT_SLIDE_WIDTH_EMU, _DEFAULT_SLIDE_HEIGHT_EMU


def _get_slide_index_and_total(slide: Any) -> tuple[int, int]:
    """
    反推 slide 在所属 Presentation 中的 0 基页码与总页数
    失败时返回 (0, 1)，保证调用方可用
    """
    try:
        part = slide.part
        pkg = part.package
        prs = _get_presentation(pkg)
        if prs is None:
            return 0, 1
        total = len(prs.slides)
        for i, s in enumerate(prs.slides):
            if s.part is part:
                return i, total
        return 0, total
    except Exception:
        return 0, 1


def _is_kpi_value_text(text: str) -> bool:
    """判断文本是否为 KPI 数值：纯数字或含 %/$"""
    if not text:
        return False
    if '%' in text or '$' in text:
        return True
    t = text.strip()
    if _NUMERIC_RE.match(t):
        return True
    return False


def _build_position(shape: Any) -> dict[str, Any]:
    """构建形状位置/尺寸信息（EMU）"""
    return {
        "left": _safe_get(shape, 'left'),
        "top": _safe_get(shape, 'top'),
        "width": _safe_get(shape, 'width'),
        "height": _safe_get(shape, 'height'),
    }


def _iter_leaf_shapes(shapes: Any) -> Iterator[Any]:
    """递归展开分组形状，产出所有叶子形状"""
    for shape in shapes:
        try:
            # GROUP = 6，递归进入子形状
            if shape.shape_type == 6:
                yield from _iter_leaf_shapes(shape.shapes)
            else:
                yield shape
        except Exception:
            yield shape


def _is_adjacent(a: Any, b: Any) -> bool:
    """
    判断形状 a 是否紧邻形状 b：
    水平区间有重叠，且垂直间距 ≤ 0.5 英寸（典型 KPI 数值与标签的上下关系）
    """
    a_l = _safe_get(a, 'left'); a_t = _safe_get(a, 'top')
    a_w = _safe_get(a, 'width'); a_h = _safe_get(a, 'height')
    b_l = _safe_get(b, 'left'); b_t = _safe_get(b, 'top')
    b_w = _safe_get(b, 'width'); b_h = _safe_get(b, 'height')
    if None in (a_l, a_t, a_w, a_h, b_l, b_t, b_w, b_h):
        return False
    a_r = a_l + a_w
    a_b = a_t + a_h
    b_r = b_l + b_w
    b_b = b_t + b_h
    h_overlap = min(a_r, b_r) - max(a_l, b_l)
    v_gap = max(0, max(a_t, b_t) - min(a_b, b_b))
    return h_overlap > 0 and v_gap <= _ADJACENT_GAP_EMU


# ==================== 对外核心函数 ====================
def classify_element(shape: Any, page_idx: int = 0, total_pages: int = 1) -> dict[str, Any]:
    """
    对单个形状进行角色分类

    :param shape: python-pptx Shape 对象
    :param page_idx: 该页 0 基页码（保留供扩展规则使用）
    :param total_pages: 总页数（保留供扩展规则使用）
    :return: {"role": str, "confidence": float, "reason": str}
    """
    text = _get_shape_text(shape)

    # 无文本 → 装饰
    if not text:
        return {"role": "decorative", "confidence": 0.5, "reason": "无文本内容"}

    font_size = _get_font_size_pt(shape)
    _, slide_h = _get_slide_dimensions(shape)
    top = _safe_get(shape, 'top')
    top_ratio = (top / slide_h) if (top is not None and slide_h) else None

    # 1. KPI 数值：纯数字或含 %/$（强格式信号，优先判定）
    if _is_kpi_value_text(text):
        return {
            "role": "kpi_value",
            "confidence": 0.9,
            "reason": f"文本为数值/百分比/金额: {text[:20]}",
        }

    # 2. 标题：字号 ≥ 36 且 top < 30% 高度
    if font_size is not None and font_size >= 36 and top_ratio is not None and top_ratio < 0.30:
        return {
            "role": "title",
            "confidence": 0.9,
            "reason": f"字号 {font_size}pt 且位于顶部 {top_ratio:.0%}",
        }

    # 3. 副标题：字号 24-32 且 top < 40%
    if font_size is not None and 24 <= font_size <= 32 and top_ratio is not None and top_ratio < 0.40:
        return {
            "role": "subtitle",
            "confidence": 0.85,
            "reason": f"字号 {font_size}pt 且位于上部 {top_ratio:.0%}",
        }

    # 4. 备注：字号 < 18 且 top > 70%
    if font_size is not None and font_size < 18 and top_ratio is not None and top_ratio > 0.70:
        return {
            "role": "note",
            "confidence": 0.8,
            "reason": f"字号 {font_size}pt 且位于底部 {top_ratio:.0%}",
        }

    # 5. 正文：字号 18-24 且 run.text 长度 > 15
    if font_size is not None and 18 <= font_size <= 24 and len(text) > 15:
        return {
            "role": "body",
            "confidence": 0.85,
            "reason": f"字号 {font_size}pt 且文本较长({len(text)}字)",
        }

    # 6. 兜底：置信度 < 0.6 → 装饰
    return {
        "role": "decorative",
        "confidence": 0.5,
        "reason": f"未命中强特征规则(字号={font_size}, top_ratio={top_ratio})",
    }


def classify_page(slide: Any) -> dict[str, Any]:
    """
    对整页所有形状进行分类，并产出低置信度告警

    :param slide: python-pptx Slide 对象
    :return: {
        "elements": [{"role","confidence","reason","shape_name","text","position"}],
        "low_confidence_warnings": [{"shape_name","position","reason"}]
    }
    """
    elements: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []

    try:
        page_idx, total_pages = _get_slide_index_and_total(slide)
    except Exception:
        page_idx, total_pages = 0, 1

    leaf_shapes = list(_iter_leaf_shapes(slide.shapes))

    # 第一遍：逐形状分类
    classified: list[tuple[Any, dict[str, Any]]] = []
    for shape in leaf_shapes:
        info = classify_element(shape, page_idx, total_pages)
        entry = {
            "role": info["role"],
            "confidence": info["confidence"],
            "reason": info["reason"],
            "shape_name": _safe_get(shape, 'name', '') or '',
            "text": _get_shape_text(shape)[:50],
            "position": _build_position(shape),
        }
        elements.append(entry)
        classified.append((shape, info))

    # 第二遍：kpi_label 邻接判定
    # 短文本（< 10 字）且紧邻 kpi_value 的元素，重判为 kpi_label
    kpi_value_shapes = [sh for sh, inf in classified if inf["role"] == "kpi_value"]
    if kpi_value_shapes:
        for i, (shape, info) in enumerate(classified):
            if info["role"] == "kpi_value":
                continue
            # 已高置信度判为 title 的不重判，避免把真标题变 label
            if info["role"] == "title" and info["confidence"] >= 0.9:
                continue
            text = _get_shape_text(shape)
            if not text or len(text) >= 10:
                continue
            for kv_shape in kpi_value_shapes:
                if _is_adjacent(shape, kv_shape):
                    new_info = {
                        "role": "kpi_label",
                        "confidence": 0.85,
                        "reason": "短文本紧邻 KPI 数值，判定为标签",
                    }
                    elements[i].update({
                        "role": new_info["role"],
                        "confidence": new_info["confidence"],
                        "reason": new_info["reason"],
                    })
                    classified[i] = (shape, new_info)
                    break

    # 低置信度告警（confidence < 0.6）
    for shape, info in classified:
        if info["confidence"] < 0.6:
            warnings.append({
                "shape_name": _safe_get(shape, 'name', '') or '',
                "position": _build_position(shape),
                "reason": info["reason"],
            })

    return {"elements": elements, "low_confidence_warnings": warnings}
