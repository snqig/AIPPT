"""
文本替换子模块 - 渲染引擎拆分
功能：递归遍历 shape、按槽位信息定位 shape、文本替换（保留格式）、长文本自适应
依赖：python-pptx
"""
from typing import Any, Optional

from pptx.shapes.base import BaseShape
from pptx.util import Pt

from aippt.logger import logger


def iter_all_shapes(shapes: Any) -> Any:
    """递归遍历所有 shape，含 GROUP 内的子 shape

    :param shapes: slide.shapes 或 group.shapes
    :yield: shape 对象
    """
    for s in shapes:
        yield s
        if s.shape_type == 6:  # GROUP
            yield from iter_all_shapes(s.shapes)


def find_shape(shapes: Any, slot_info: dict[str, Any],
               used: set[int]) -> Optional[BaseShape]:
    """根据槽位信息定位目标 shape

    匹配策略（按优先级）：
    1. shape_name + match_text 双匹配（最精确）
    2. match_text 文本匹配
    3. shape_name 匹配（兜底）

    :param shapes: 待搜索的 shapes 集合
    :param slot_info: 槽位信息，含 match_text / shape_name 字段
    :param used: 已使用的 shape element id 集合，避免重复匹配
    :return: 匹配到的 shape，未找到返回 None
    """
    match_text = slot_info.get('match_text', '').strip()
    shape_name = slot_info.get('shape_name', '')
    all_shapes = list(iter_all_shapes(shapes))

    # 策略1：shape_name + match_text 双匹配（最精确）
    if shape_name and match_text:
        for s in all_shapes:
            if id(s._element) in used:
                continue
            if not s.has_text_frame:
                continue
            if s.name == shape_name and match_text in s.text_frame.text:
                return s

    # 策略2：match_text 文本匹配
    if match_text:
        for s in all_shapes:
            if id(s._element) in used:
                continue
            if not s.has_text_frame:
                continue
            if match_text in s.text_frame.text:
                return s

    # 策略3：shape_name 匹配（兜底）
    if shape_name:
        for s in all_shapes:
            if id(s._element) in used:
                continue
            if not s.has_text_frame:
                continue
            if s.name == shape_name:
                return s

    return None


def replace_text(shape: BaseShape, new_text: str) -> None:
    """替换 shape 的文本，保留首个 run 的格式

    策略：保留首段首 run 的格式，仅替换其 text，删除同段其余 run 和其余段落

    :param shape: 目标 shape
    :param new_text: 新文本内容
    """
    tf = shape.text_frame
    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        return

    first_para = paragraphs[0]

    # 保留第一个 run 的格式，仅改其 text，删除同段其余 run
    if first_para.runs:
        first_para.runs[0].text = new_text
        for run in first_para.runs[1:]:
            run._r.getparent().remove(run._r)
    else:
        run = first_para.add_run()
        run.text = new_text

    # 删除其余段落（保留首段段落属性）
    for para in paragraphs[1:]:
        para._p.getparent().remove(para._p)


def auto_fit(shape: BaseShape, new_text: str, original_text: str,
             slot_info: Optional[dict[str, Any]] = None,
             page_str: Optional[str] = None,
             slot_name: Optional[str] = None) -> None:
    """长文本字号自适应：超容量时按比例缩小字号（下限 8pt）

    策略：
    1. 优先使用 meta 中的 capacity.total_chars（事前预警）
    2. 回退到运行时几何估算（基于 shape 宽高 + 字号）

    :param shape: 目标 shape
    :param new_text: 新文本
    :param original_text: 原文本（保留参数，便于扩展）
    :param slot_info: 槽位信息，可含 capacity 字段
    :param page_str: 页码字符串（用于日志）
    :param slot_name: 槽位名（用于日志）
    """
    tf = shape.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        return
    run = tf.paragraphs[0].runs[0]
    if run.font.size is None:
        return
    original_size_pt = run.font.size.pt

    # 策略1：优先使用 meta 中的 capacity（事前预警）
    capacity = None
    if slot_info and isinstance(slot_info.get('capacity'), dict):
        capacity = slot_info['capacity'].get('total_chars')

    # 策略2：回退到运行时几何估算
    if capacity is None:
        try:
            box_w = shape.width        # EMU
            box_h = shape.height       # EMU
        except Exception:
            return
        if not box_w or not box_h:
            return
        EMU_PER_PT = 12700
        char_w = original_size_pt * 0.55 * EMU_PER_PT
        line_h = original_size_pt * 1.2 * EMU_PER_PT
        chars_per_line = max(1, int(box_w / char_w))
        lines_available = max(1, int(box_h / line_h))
        capacity = chars_per_line * lines_available

    text_len = len(new_text)
    if text_len <= capacity:
        return

    if page_str and slot_name:
        logger.warning("页%s 槽位 %s: 文本长度 %d 超过容量 %d，自动缩字号",
                       page_str, slot_name, text_len, capacity)

    # 按比例缩小字号，下限 8pt
    ratio = (capacity / text_len) ** 0.5
    new_size = max(8, int(original_size_pt * ratio))
    if new_size < original_size_pt:
        run.font.size = Pt(new_size)
