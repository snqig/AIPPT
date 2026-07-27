"""
自动布局引擎骨架（T002）

模块分层：
    1. 主题加载器（theme_loader，T003 实现）
    2. 栅格坐标计算工具（12 列栅格、安全区、分区计算）
    3. 原子绘制组件：add_text_box / add_shape / add_line / add_card / render_number_badge
    4. 页面分发入口：根据 page_type 调用对应布局函数
    5. 元素角色自动标记逻辑（role / shape_id）

设计约束（严格遵守前置约束）：
    - 单位统一英寸（Inches），对齐 python-pptx 原生坐标体系
    - 所有样式从主题 Design Token 读取，禁止硬编码颜色/字号/间距
    - 自动生成元素必须附加 role、唯一 shape_id，供动画模块匹配
    - 渲染输出全部为原生可编辑 PPTX 元素
    - 图片嵌入仅在 --enable-assets 模式下启用，默认不启用
    - 16:9 画布默认 13.333 x 7.5 英寸

页面布局函数注册表 PAGE_LAYOUT_REGISTRY：
    page_type → layout_func(slide, page_data, theme, ctx)
    T004 起逐步注册 cover/catalog/divider/numbered_list/kpi
    T101 起补齐 timeline/two_column/table
"""
from __future__ import annotations

from dataclasses import dataclass, field
from typing import Any, Callable, Optional

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from aippt.logger import logger


# ==================== 画布与栅格常量 ====================
#: 16:9 画布宽度（英寸）
CANVAS_WIDTH_INCH: float = 13.333
#: 16:9 画布高度（英寸）
CANVAS_HEIGHT_INCH: float = 7.5

#: 12 列栅格列数
GRID_COLUMNS: int = 12
#: 栅格列间距（英寸）
GRID_GUTTER_INCH: float = 0.25

#: 安全边距（英寸）— 上下左右统一
SAFE_MARGIN_INCH: float = 0.5


# ==================== 数据结构 ====================
@dataclass
class GridRect:
    """栅格坐标矩形（单位：英寸）

    表示一个由 left/top/width/height 描述的矩形区域，
    所有值均为英寸，对齐 python-pptx 原生坐标体系。

    :param left: 左上角 x 坐标（英寸）
    :param top: 左上角 y 坐标（英寸）
    :param width: 宽度（英寸）
    :param height: 高度（英寸）
    """
    left: float
    top: float
    width: float
    height: float

    def to_emu(self) -> tuple[int, int, int, int]:
        """转换为 python-pptx 所需的 EMU 整数元组

        :return: (left_emu, top_emu, width_emu, height_emu)
        """
        return (Inches(self.left), Inches(self.top),
                Inches(self.width), Inches(self.height))


@dataclass
class ElementMeta:
    """元素元数据（role + shape_id 标记）

    自动布局生成的每个元素必须附加 ElementMeta，用于：
    1. 动画模块按 role 匹配（title/subtitle/desc/number/year 等）
    2. shape_id 全局唯一，便于后续追溯与质量校验

    :param shape_id: 全局唯一 shape 标识（格式：p{page_num}_{role}_{seq}）
    :param role: 元素角色（title/subtitle/desc/number/year/body/card 等）
    :param page_num: 所属页码（1-based）
    :param seq: 同页同角色序号（0-based）
    """
    shape_id: str
    role: str
    page_num: int
    seq: int


@dataclass
class LayoutContext:
    """布局上下文（跨原子组件传递）

    封装当前页码、已生成元素列表、shape_id 计数器等，
    避免原子组件函数签名过长。

    :param page_num: 当前页码（1-based）
    :param elements: 当前页已生成元素的 ElementMeta 列表
    :param role_seq: 同页同角色序号计数器 {role: seq}
    """
    page_num: int
    elements: list[ElementMeta] = field(default_factory=list)
    role_seq: dict[str, int] = field(default_factory=dict)

    def next_seq(self, role: str) -> int:
        """获取指定角色的下一个序号

        :param role: 元素角色
        :return: 序号（0-based，每次调用自增）
        """
        seq = self.role_seq.get(role, 0)
        self.role_seq[role] = seq + 1
        return seq

    def make_shape_id(self, role: str) -> str:
        """生成全局唯一 shape_id

        格式：p{page_num}_{role}_{seq}，如 p1_title_0、p3_desc_2

        :param role: 元素角色
        :return: 唯一 shape_id 字符串
        """
        seq = self.next_seq(role)
        return f"p{self.page_num}_{role}_{seq}"

    def register(self, role: str, shape: BaseShape) -> ElementMeta:
        """注册已生成的 shape，附加 role + shape_id 标记

        通过 shape.name 字段写入 shape_id，便于动画模块按 name 匹配。
        返回 ElementMeta 供上层收集。

        :param role: 元素角色
        :param shape: python-pptx shape 对象
        :return: ElementMeta 元数据
        """
        shape_id = self.make_shape_id(role)
        try:
            shape.name = shape_id  # python-pptx 支持 shape.name 读写
        except Exception:
            pass  # 部分 shape 不支持改名，忽略
        meta = ElementMeta(shape_id=shape_id, role=role,
                           page_num=self.page_num, seq=self.role_seq.get(role, 0) - 1)
        self.elements.append(meta)
        return meta


# ==================== 栅格坐标计算工具 ====================
def safe_area(margin: float = SAFE_MARGIN_INCH) -> GridRect:
    """计算安全区矩形（去除四周边距后的可用区域）

    :param margin: 安全边距（英寸），默认 0.5
    :return: GridRect 安全区矩形
    """
    return GridRect(
        left=margin,
        top=margin,
        width=CANVAS_WIDTH_INCH - 2 * margin,
        height=CANVAS_HEIGHT_INCH - 2 * margin,
    )


def column_x(col: int, col_span: int = 1, margin: float = SAFE_MARGIN_INCH,
             gutter: float = GRID_GUTTER_INCH) -> tuple[float, float]:
    """计算 12 列栅格中指定列的 x 坐标与宽度

    :param col: 起始列号（0-based，0~11）
    :param col_span: 跨列数（默认 1）
    :param margin: 左右边距（英寸）
    :param gutter: 列间距（英寸）
    :return: (x_left, width) 单位英寸
    """
    if col < 0 or col >= GRID_COLUMNS:
        raise ValueError(f"col 越界: {col}（应为 0~{GRID_COLUMNS - 1}）")
    if col_span < 1 or col + col_span > GRID_COLUMNS:
        raise ValueError(f"col_span 越界: col={col}, span={col_span}")
    available_width = CANVAS_WIDTH_INCH - 2 * margin
    col_width = (available_width - (GRID_COLUMNS - 1) * gutter) / GRID_COLUMNS
    x_left = margin + col * (col_width + gutter)
    width = col_span * col_width + (col_span - 1) * gutter
    return x_left, width


def row_y(row: int, row_span: float = 1.0, margin: float = SAFE_MARGIN_INCH,
          row_height: float = 0.6) -> tuple[float, float]:
    """计算行 y 坐标与高度（基于固定行高）

    :param row: 行号（0-based）
    :param row_span: 跨行数（默认 1.0，可为小数）
    :param margin: 上下边距（英寸）
    :param row_height: 单行高度（英寸），默认 0.6
    :return: (y_top, height) 单位英寸
    """
    y_top = margin + row * row_height
    height = row_span * row_height
    return y_top, height


def split_horizontal(area: GridRect, ratio: float = 0.5,
                     gutter: float = GRID_GUTTER_INCH) -> tuple[GridRect, GridRect]:
    """将矩形水平切分为左右两块

    :param area: 原矩形
    :param ratio: 左块占比（0~1），默认 0.5
    :param gutter: 中间间距（英寸）
    :return: (left_rect, right_rect)
    """
    if not 0 < ratio < 1:
        raise ValueError(f"ratio 应在 (0,1)，当前: {ratio}")
    left_width = (area.width - gutter) * ratio
    right_width = (area.width - gutter) * (1 - ratio)
    left = GridRect(left=area.left, top=area.top,
                    width=left_width, height=area.height)
    right = GridRect(left=area.left + left_width + gutter, top=area.top,
                     width=right_width, height=area.height)
    return left, right


def split_vertical(area: GridRect, ratio: float = 0.5,
                   gutter: float = GRID_GUTTER_INCH) -> tuple[GridRect, GridRect]:
    """将矩形垂直切分为上下两块

    :param area: 原矩形
    :param ratio: 上块占比（0~1），默认 0.5
    :param gutter: 中间间距（英寸）
    :return: (top_rect, bottom_rect)
    """
    if not 0 < ratio < 1:
        raise ValueError(f"ratio 应在 (0,1)，当前: {ratio}")
    top_height = (area.height - gutter) * ratio
    bottom_height = (area.height - gutter) * (1 - ratio)
    top = GridRect(left=area.left, top=area.top,
                   width=area.width, height=top_height)
    bottom = GridRect(left=area.left, top=area.top + top_height + gutter,
                      width=area.width, height=bottom_height)
    return top, bottom


def grid_matrix(area: GridRect, rows: int, cols: int,
                gutter: float = GRID_GUTTER_INCH) -> list[list[GridRect]]:
    """将矩形切分为 rows x cols 网格矩阵

    :param area: 原矩形
    :param rows: 行数
    :param cols: 列数
    :param gutter: 单元格间距（英寸）
    :return: 二维 GridRect 矩阵 [row][col]
    """
    if rows < 1 or cols < 1:
        raise ValueError(f"rows/cols 应 >=1，当前: {rows}x{cols}")
    cell_w = (area.width - (cols - 1) * gutter) / cols
    cell_h = (area.height - (rows - 1) * gutter) / rows
    matrix: list[list[GridRect]] = []
    for r in range(rows):
        row_cells: list[GridRect] = []
        for c in range(cols):
            cell = GridRect(
                left=area.left + c * (cell_w + gutter),
                top=area.top + r * (cell_h + gutter),
                width=cell_w,
                height=cell_h,
            )
            row_cells.append(cell)
        matrix.append(row_cells)
    return matrix


# ==================== 主题令牌读取工具（T003 前的轻量兜底）====================
def get_token(theme: dict[str, Any], path: str, default: Any = None) -> Any:
    """从主题令牌字典按点分路径读取值

    示例：
        get_token(theme, "color.primary", "#1A56DB")
        get_token(theme, "font.title.size_pt", 36)

    T003 完成后由 theme_loader 提供完整接口，本函数为骨架阶段的轻量实现。

    :param theme: 主题字典
    :param path: 点分路径（如 "color.primary"）
    :param default: 路径不存在时的默认值
    :return: 令牌值或 default
    """
    if not theme:
        return default
    keys = path.split(".")
    cur: Any = theme
    for k in keys:
        if not isinstance(cur, dict):
            return default
        cur = cur.get(k)
        if cur is None:
            return default
    return cur


def hex_to_rgb(hex_str: str) -> RGBColor:
    """十六进制颜色字符串转 RGBColor

    :param hex_str: 十六进制颜色（如 "#1A56DB" 或 "1A56DB"）
    :return: RGBColor 实例
    """
    s = hex_str.lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) != 6:
        raise ValueError(f"非法颜色值: {hex_str}")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


# ==================== 原子绘制组件 ====================
def add_text_box(
    slide: Slide,
    rect: GridRect,
    text: str,
    role: str,
    ctx: LayoutContext,
    theme: dict[str, Any],
    font_size_pt: Optional[float] = None,
    bold: bool = False,
    color: Optional[str] = None,
    align: str = "left",
    anchor: str = "top",
    font_name: Optional[str] = None,
) -> BaseShape:
    """添加文本框原子组件

    样式全部从主题令牌读取：
    - font_size_pt: None 时按 role 从主题读取（如 font.title.size_pt）
    - color: None 时按 role 从主题读取（如 color.text_primary）
    - font_name: None 时从主题读取 font.family

    :param slide: 目标 slide
    :param rect: 文本框位置与尺寸
    :param text: 文本内容
    :param role: 元素角色（用于动画匹配与主题令牌查找）
    :param ctx: 布局上下文
    :param theme: 主题字典
    :param font_size_pt: 字号（pt），None 时从主题读取
    :param bold: 是否粗体
    :param color: 字体颜色（hex），None 时从主题读取
    :param align: 对齐方式（left/center/right）
    :param anchor: 垂直锚点（top/middle/bottom）
    :param font_name: 字体名，None 时从主题读取
    :return: python-pptx shape 对象
    """
    left, top, width, height = rect.to_emu()
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    # 垂直锚点
    anchor_map = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }
    tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)

    # 写入文本（保留首段格式）
    p = tf.paragraphs[0]
    p.text = text
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)

    # 应用字体样式
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    if font_size_pt is None:
        font_size_pt = get_token(theme, f"font.{role}.size_pt",
                                 get_token(theme, "font.body.size_pt", 18))
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold

    if color is None:
        color = get_token(theme, f"color.{role}",
                          get_token(theme, "color.text_primary", "#333333"))
    run.font.color.rgb = hex_to_rgb(color)

    if font_name is None:
        font_name = get_token(theme, "font.family", "Microsoft YaHei")
    run.font.name = font_name

    ctx.register(role, tb)
    return tb


def add_shape(
    slide: Slide,
    rect: GridRect,
    shape_type: int,
    role: str,
    ctx: LayoutContext,
    theme: dict[str, Any],
    fill_color: Optional[str] = None,
    line_color: Optional[str] = None,
    line_width_pt: float = 0.75,
) -> BaseShape:
    """添加形状原子组件（矩形/圆角矩形/线条等）

    :param slide: 目标 slide
    :param rect: 形状位置与尺寸
    :param shape_type: MSO_SHAPE 枚举值
    :param role: 元素角色
    :param ctx: 布局上下文
    :param theme: 主题字典
    :param fill_color: 填充色（hex），None 时从主题读取 color.{role}_bg
    :param line_color: 边框色（hex），None 时从主题读取 color.{role}_border
    :param line_width_pt: 边框宽度（pt）
    :return: python-pptx shape 对象
    """
    left, top, width, height = rect.to_emu()
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)

    if fill_color is None:
        fill_color = get_token(theme, f"color.{role}_bg",
                               get_token(theme, "color.card_bg", "#FFFFFF"))
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(fill_color)

    if line_color is None:
        line_color = get_token(theme, f"color.{role}_border",
                               get_token(theme, "color.card_border", "#E5E7EB"))
    shp.line.color.rgb = hex_to_rgb(line_color)
    shp.line.width = Pt(line_width_pt)

    # 关闭阴影（主题令牌控制）
    shadow_enabled = get_token(theme, "effect.card_shadow", False)
    if not shadow_enabled:
        shp.shadow.inherit = False

    ctx.register(role, shp)
    return shp


def add_line(
    slide: Slide,
    x1: float, y1: float, x2: float, y2: float,
    role: str,
    ctx: LayoutContext,
    theme: dict[str, Any],
    color: Optional[str] = None,
    width_pt: float = 1.0,
) -> BaseShape:
    """添加连接线原子组件

    :param slide: 目标 slide
    :param x1: 起点 x（英寸）
    :param y1: 起点 y（英寸）
    :param x2: 终点 x（英寸）
    :param y2: 终点 y（英寸）
    :param role: 元素角色
    :param ctx: 布局上下文
    :param theme: 主题字典
    :param color: 线条颜色（hex），None 时从主题读取 color.divider
    :param width_pt: 线宽（pt）
    :return: python-pptx shape 对象
    """
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    if color is None:
        color = get_token(theme, "color.divider", "#E5E7EB")
    connector.line.color.rgb = hex_to_rgb(color)
    connector.line.width = Pt(width_pt)
    ctx.register(role, connector)
    return connector


def add_image_box(
    slide: Slide,
    rect: GridRect,
    image_path: str,
    role: str,
    ctx: LayoutContext,
    theme: dict[str, Any],
    mask: str = "none",
    overlay: bool = False,
) -> Optional[BaseShape]:
    left, top, width, height = rect.to_emu()
    try:
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        ctx.register(role, pic)
        if overlay:
            _add_image_overlay(slide, rect, theme, ctx)
        return pic
    except Exception as e:
        logger.warning("add_image_box failed: %s", e)
        return None


def _add_image_overlay(
    slide: Slide,
    rect: GridRect,
    theme: dict[str, Any],
    ctx: LayoutContext,
) -> BaseShape:
    left, top, width, height = rect.to_emu()
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    overlay_color = get_token(theme, "color.image_overlay", "#000000")
    overlay_opacity = get_token(theme, "effect.image_overlay_opacity", 35000)
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(overlay_color)
    from pptx.oxml.ns import qn
    alpha_elem = shp.fill.fore_color._color
    alpha_elem.set(qn("a:alpha"), str(int(overlay_opacity)))
    shp.line.fill.background()
    ctx.register("image_overlay", shp)
    return shp


def add_card(
    slide: Slide,
    rect: GridRect,
    role: str,
    ctx: LayoutContext,
    theme: dict[str, Any],
    radius: bool = True,
) -> BaseShape:
    """添加卡片容器（圆角矩形 + 阴影 + 主题配色）

    :param slide: 目标 slide
    :param rect: 卡片位置与尺寸
    :param role: 元素角色（如 card / kpi_card）
    :param ctx: 布局上下文
    :param theme: 主题字典
    :param radius: 是否圆角（True=RoundedRectangle，False=Rectangle）
    :return: python-pptx shape 对象
    """
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    return add_shape(slide, rect, shape_type, role, ctx, theme)


def render_number_badge(
    slide: Slide,
    rect: GridRect,
    number: str,
    role: str,
    ctx: LayoutContext,
    theme: dict[str, Any],
) -> BaseShape:
    """渲染数字徽章（圆形/方形 + 大号数字）

    用于 numbered_list 页的序号、kpi 页的数值等。

    :param slide: 目标 slide
    :param rect: 徽章位置与尺寸
    :param number: 数字文本（如 "01" / "67%"）
    :param role: 元素角色（如 number / badge）
    :param ctx: 布局上下文
    :param theme: 主题字典
    :return: python-pptx shape 对象
    """
    # 圆形背景
    badge = add_shape(slide, rect, MSO_SHAPE.OVAL, role, ctx, theme,
                      fill_color=get_token(theme, f"color.{role}_bg",
                                           get_token(theme, "color.primary", "#1A56DB")))
    # 数字文本
    tf = badge.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    run.text = number
    run.font.size = Pt(get_token(theme, f"font.{role}.size_pt", 28))
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(get_token(theme, f"color.{role}_fg",
                                              get_token(theme, "color.text_on_primary", "#FFFFFF")))
    run.font.name = get_token(theme, "font.family", "Microsoft YaHei")
    return badge


# ==================== 页面布局函数注册表 ====================
#: 二级注册表：page_type → {variant_name → layout_func}
#: T613 升级：支持 layout_variant 版式变体机制
#: 每种 page_type 至少含 "default" 变体，可扩展 "grid_2x2" / "single_column" 等
PAGE_LAYOUT_REGISTRY: dict[str, dict[str, Callable[[Slide, dict, dict, LayoutContext], None]]] = {}

#: 默认变体名
DEFAULT_VARIANT = "default"


def register_layout(page_type: str, variant: str = DEFAULT_VARIANT) -> Callable:
    """页面布局函数注册装饰器（T613 支持 variant）

    用法：
        # 默认变体
        @register_layout("cover")
        def layout_cover(slide, page_data, theme, ctx):
            ...

        # 指定变体（如 KPI 页 2×2 网格变体）
        @register_layout("kpi", variant="grid_2x2")
        def layout_kpi_grid(slide, page_data, theme, ctx):
            ...

    :param page_type: 页面类型（cover/catalog/divider/numbered_list/kpi/...）
    :param variant: 版式变体名，默认 "default"
    :return: 装饰器函数
    """
    def decorator(func: Callable) -> Callable:
        variants = PAGE_LAYOUT_REGISTRY.setdefault(page_type, {})
        if variant in variants:
            logger.warning("page_type=%s variant=%s 已注册，将被覆盖", page_type, variant)
        variants[variant] = func
        return func
    return decorator


def dispatch_page_layout(
    slide: Slide,
    page_data: dict[str, Any],
    theme: dict[str, Any],
    ctx: LayoutContext,
    variant_override: Optional[dict[str, str]] = None,
) -> list[ElementMeta]:
    """页面分发入口：根据 page_type + variant 调用对应布局函数（T613 支持 variant）

    variant 选择优先级（高 → 低）：
        1. variant_override[page_type]：CLI 全局 --layout-variant-override 参数
        2. page_data["layout_variant"]：单页 outline 内显式指定
        3. "default"：默认变体

    :param slide: 目标 slide
    :param page_data: outline.json 中的单页数据
    :param theme: 主题字典
    :param ctx: 布局上下文
    :param variant_override: 全局变体覆盖映射 {page_type: variant_name}，可选
    :return: 当前页生成的元素 ElementMeta 列表
    :raises ValueError: 未注册的 page_type 且无 numbered_list 兜底
    """
    page_type = page_data.get("page_type", "numbered_list")
    variants = PAGE_LAYOUT_REGISTRY.get(page_type)

    if variants is None:
        # 兜底：使用 numbered_list 布局
        logger.warning("未注册的 page_type: %s，降级为 numbered_list", page_type)
        variants = PAGE_LAYOUT_REGISTRY.get("numbered_list")
        if variants is None:
            raise ValueError(f"page_type {page_type} 未注册且无 numbered_list 兜底")
        page_type = "numbered_list"

    # 选择 variant：全局 override > 单页显式 > default
    variant_name = DEFAULT_VARIANT
    if variant_override and page_type in variant_override:
        variant_name = variant_override[page_type]
    elif "layout_variant" in page_data:
        variant_name = page_data["layout_variant"]

    layout_func = variants.get(variant_name)
    if layout_func is None:
        # 变体未注册，回退 default
        if variant_name != DEFAULT_VARIANT:
            logger.warning("page_type=%s variant=%s 未注册，回退 default",
                          page_type, variant_name)
        layout_func = variants.get(DEFAULT_VARIANT)
        if layout_func is None:
            # 极端情况：该 page_type 无 default 变体，取第一个注册的
            layout_func = next(iter(variants.values()))
    layout_func(slide, page_data, theme, ctx)
    return ctx.elements


def list_layout_variants(page_type: Optional[str] = None) -> dict[str, list[str]]:
    """查询已注册的页面布局变体（T613 新增，供 CLI list-variants 调用）

    :param page_type: 指定页面类型，None 查询全部
    :return: {page_type: [variant_name, ...]} 字典
    """
    if page_type:
        variants = PAGE_LAYOUT_REGISTRY.get(page_type, {})
        return {page_type: list(variants.keys())}
    return {pt: list(variants.keys()) for pt, variants in PAGE_LAYOUT_REGISTRY.items()}


def create_presentation(slide_width_inch: float = CANVAS_WIDTH_INCH,
                        slide_height_inch: float = CANVAS_HEIGHT_INCH) -> Presentation:
    """创建 16:9 空白 Presentation

    :param slide_width_inch: 画布宽度（英寸），默认 13.333
    :param slide_height_inch: 画布高度（英寸），默认 7.5
    :return: Presentation 对象
    """
    prs = Presentation()
    prs.slide_width = Inches(slide_width_inch)
    prs.slide_height = Inches(slide_height_inch)
    return prs


def add_blank_slide(prs: Presentation) -> Slide:
    """添加空白 slide（使用 Blank 版式）

    :param prs: Presentation 对象
    :return: 新 slide
    """
    blank_layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(blank_layout)
