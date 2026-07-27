"""
8 类核心页面标准布局（T004 + T603）

页面类型：
    1. cover        封面页：上下分割版式
    2. catalog      目录页：左标题 + 右编号列表
    3. divider      章节分隔页：左右分栏大号序号
    4. numbered_list 编号列表页：条目卡片化布局
    5. kpi          KPI 页：等宽横向卡片矩阵
    6. timeline     时间轴页：水平时间轴 + 节点事件（T603 新增）
    7. two_column   双栏页：左右分栏对照（T603 新增）
    8. table        表格页：原生 PPT 表格 + 表头样式继承（T603 新增）

设计规范：
    - 严格执行安全边距、栅格布局、字体层级、卡片圆角/阴影、间距令牌、文本自适应
    - 所有样式从主题令牌读取，禁止硬编码
    - 元素自动标记 role + shape_id（供动画模块匹配）
    - 渲染输出为原生可编辑 PPTX 元素

字段约定（outline.json page 对象）：
    - cover: title, subtitle
    - catalog: title, items[]
    - divider: title, section_no
    - numbered_list: title, items[]
    - kpi: title, kpi_items[{label, value, trend}]
    - timeline: title, timeline_items[{time, event}]
    - two_column: title, left_title, left_items[], right_title, right_items[]
    - table: title, headers[], rows[][]
"""
from __future__ import annotations

from typing import Any

from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from aippt.logger import logger
from aippt.layout.ppt_auto_layout import (
    GridRect, LayoutContext, ElementMeta,
    safe_area, column_x, row_y,
    split_horizontal, split_vertical, grid_matrix,
    get_token, hex_to_rgb,
    add_text_box, add_shape, add_line, add_card, render_number_badge,
    register_layout,
)
from aippt.theme_loader import resolve_color
from aippt.layout.elastic_constraint import (
    elastic_distribute_items,
    estimate_text_height,
    EMU_PER_INCH,
)


# ==================== cover 封面页：上下分割版式 ====================
@register_layout("cover")
def layout_cover(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """封面页布局：上下分割

    版式：
        - 上半部分（45%）：主题色背景，大号标题居中
        - 下半部分（55%）：白色背景，副标题居中
    字段：
        - title: 主标题（必须）
        - subtitle: 副标题（可选）
    元素角色：
        - cover_bg: 背景矩形
        - title: 主标题
        - subtitle: 副标题
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))
    split_ratio = get_token(theme, "layout.cover_split_ratio", 0.45)
    top_rect, bottom_rect = split_vertical(area, ratio=split_ratio, gutter=0)

    # 上半部分主题色背景（铺满整宽，超出安全区到画布顶部）
    cover_bg_rect = GridRect(
        left=0, top=0,
        width=13.333,
        height=split_ratio * area.height + get_token(theme, "spacing.safe_margin_inch", 0.5),
    )
    cover_bg = add_shape(
        slide, cover_bg_rect, 1, "cover_bg", ctx, theme,  # 1 = MSO_SHAPE.RECTANGLE
        fill_color=resolve_color(theme, get_token(theme, "color.cover_bg", "#1A56DB")),
        line_color=None,
        line_width_pt=0,
    )
    # 去掉边框
    cover_bg.line.fill.background()

    # 主标题（上半部分居中）
    title_text = page_data.get("title", "")
    title_rect = GridRect(
        left=area.left, top=cover_bg_rect.top + cover_bg_rect.height * 0.3,
        width=area.width, height=cover_bg_rect.height * 0.4,
    )
    add_text_box(
        slide, title_rect, title_text, "cover_title", ctx, theme,
        font_size_pt=get_token(theme, "font.cover_title.size_pt", 44),
        bold=get_token(theme, "font.cover_title.bold", True),
        color=resolve_color(theme, get_token(theme, "color.cover_title", "#FFFFFF")),
        align="center", anchor="middle",
    )

    # 副标题（下半部分居中）
    subtitle_text = page_data.get("subtitle", "")
    if subtitle_text:
        sub_rect = GridRect(
            left=area.left, top=cover_bg_rect.height + 0.5,
            width=area.width, height=1.0,
        )
        add_text_box(
            slide, sub_rect, subtitle_text, "cover_subtitle", ctx, theme,
            font_size_pt=get_token(theme, "font.cover_subtitle.size_pt", 22),
            bold=get_token(theme, "font.cover_subtitle.bold", False),
            color=resolve_color(theme, get_token(theme, "color.cover_subtitle", "#DBEAFE")),
            align="center", anchor="middle",
        )


# ==================== catalog 目录页：左标题 + 右编号列表 ====================
@register_layout("catalog")
def layout_catalog(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """目录页布局：左标题 + 右编号列表

    版式：
        - 左侧（35%）：大号"目录"标题 + 副标题
        - 右侧（65%）：编号列表，每项含序号 + 标题
    字段：
        - title: 目录标题（如"目录"或"CONTENTS"）
        - items[]: 目录条目列表
    元素角色：
        - title: 左侧大标题
        - body: 右侧每个目录条目
        - number: 每个条目的序号
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))
    left_rect, right_rect = split_horizontal(area, ratio=0.35)

    # 左侧大标题
    title_text = page_data.get("title", "目录")
    add_text_box(
        slide, GridRect(left_rect.left, left_rect.top + 1.5, left_rect.width, 1.5),
        title_text, "title", ctx, theme,
        font_size_pt=get_token(theme, "font.title.size_pt", 36),
        bold=True,
        color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
        align="left", anchor="middle",
    )

    # 左侧装饰线
    add_line(
        slide, left_rect.left, left_rect.top + 3.2,
        left_rect.left + 1.5, left_rect.top + 3.2,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    # 右侧编号列表
    items = page_data.get("items", [])
    if not items:
        return

    # T612：弹性约束接入 - 长短条目自适应高度
    # 通过 theme.layout.elastic_enabled 开关控制（默认 True），可回退静态等高
    elastic_enabled = get_token(theme, "layout.elastic_enabled", True)
    item_height = 0.7
    gap = get_token(theme, "spacing.element_gap_inch", 0.15)
    start_y = right_rect.top + 0.3
    list_area_h = area.top + area.height - start_y

    if elastic_enabled and items:
        # 弹性模式：按文本长度分配高度
        items_data = [{"text": str(it)} for it in items]
        positions = elastic_distribute_items(
            items_data,
            area_top_emu=int(start_y * EMU_PER_INCH),
            area_height_emu=int(list_area_h * EMU_PER_INCH),
            area_width_emu=int(right_rect.width * EMU_PER_INCH),
            font_size_pt=get_token(theme, "font.body.size_pt", 16),
            min_item_h_inch=0.5,
            gap_inch=gap,
        )
    else:
        # 静态模式：固定高度等分（原逻辑）
        positions = None

    for idx, item in enumerate(items):
        if positions is not None:
            top_emu, h_emu = positions[idx]
            item_y = top_emu / EMU_PER_INCH
            item_height = h_emu / EMU_PER_INCH
            if item_y + item_height > area.top + area.height:
                break  # 防止溢出
        else:
            item_y = start_y + idx * (item_height + gap)
            if item_y + item_height > area.top + area.height:
                break  # 防止溢出

        # 序号徽章
        badge_size = 0.5
        badge_rect = GridRect(
            left=right_rect.left, top=item_y,
            width=badge_size, height=badge_size,
        )
        render_number_badge(
            slide, badge_rect, f"{idx + 1:02d}", "number", ctx, theme,
        )

        # 条目文本
        text_rect = GridRect(
            left=right_rect.left + badge_size + 0.2, top=item_y,
            width=right_rect.width - badge_size - 0.2, height=item_height,
        )
        add_text_box(
            slide, text_rect, str(item), "body", ctx, theme,
            font_size_pt=get_token(theme, "font.body.size_pt", 16),
            bold=False,
            color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
            align="left", anchor="middle",
        )


# ==================== divider 章节分隔页：左右分栏大号序号 ====================
@register_layout("divider")
def layout_divider(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """章节分隔页布局：左右分栏大号序号

    版式：
        - 左侧（35%）：超大号章节序号（如 01）
        - 右侧（65%）：章节标题
    字段：
        - section_no: 章节序号（如 "01" / "PART ONE"）
        - title: 章节标题
    元素角色：
        - number: 左侧大号序号
        - title: 右侧章节标题
        - divider_bg: 全页背景
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

    # 全页背景
    bg_rect = GridRect(left=0, top=0, width=13.333, height=7.5)
    bg = add_shape(
        slide, bg_rect, 1, "divider_bg", ctx, theme,
        fill_color=resolve_color(theme, get_token(theme, "color.divider_bg", "#1E40AF")),
        line_color=None,
    )
    bg.line.fill.background()

    # 左右分栏
    left_rect, right_rect = split_horizontal(area, ratio=0.35)

    # 左侧大号序号
    section_no = page_data.get("section_no", "")
    if not section_no:
        # 从 title 前缀提取（如 "01. 项目背景" → "01"）
        title = page_data.get("title", "")
        if title and len(title) >= 2 and title[:2].isdigit():
            section_no = title[:2]
        else:
            section_no = ""

    if section_no:
        add_text_box(
            slide, GridRect(left_rect.left, left_rect.top + 1.5, left_rect.width, 3.5),
            section_no, "divider_number", ctx, theme,
            font_size_pt=get_token(theme, "font.divider_number.size_pt", 72),
            bold=get_token(theme, "font.divider_number.bold", True),
            color=resolve_color(theme, get_token(theme, "color.divider_number", "#FFFFFF")),
            align="center", anchor="middle",
        )

    # 右侧章节标题
    title_text = page_data.get("title", "")
    # 去掉前缀序号
    if title_text and section_no and title_text.startswith(section_no):
        title_text = title_text[len(section_no):].lstrip(".、 ")

    add_text_box(
        slide, GridRect(right_rect.left, right_rect.top + 2.5, right_rect.width, 2.0),
        title_text, "divider_title", ctx, theme,
        font_size_pt=get_token(theme, "font.divider_title.size_pt", 40),
        bold=get_token(theme, "font.divider_title.bold", True),
        color=resolve_color(theme, get_token(theme, "color.divider_title", "#FFFFFF")),
        align="left", anchor="middle",
    )


# ==================== numbered_list 编号列表页：条目卡片化布局 ====================
@register_layout("numbered_list")
def layout_numbered_list(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """编号列表页布局：条目卡片化

    版式：
        - 顶部标题区（占 20%）
        - 下方卡片网格（2 列布局，每卡片含序号 + 标题）
    字段：
        - title: 页面标题
        - items[]: 条目列表
    元素角色：
        - title: 页面标题
        - card: 每个条目卡片
        - number: 卡片序号
        - body: 卡片文本
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

    # 顶部标题区
    title_rect = GridRect(area.left, area.top, area.width, 1.0)
    add_text_box(
        slide, title_rect, page_data.get("title", ""), "title", ctx, theme,
        font_size_pt=get_token(theme, "font.title.size_pt", 36),
        bold=True,
        color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
        align="left", anchor="middle",
    )

    # 装饰下划线
    add_line(
        slide, area.left, area.top + 1.1,
        area.left + 1.5, area.top + 1.1,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    # 卡片网格区域
    items = page_data.get("items", [])
    if not items:
        return

    cards_area = GridRect(
        left=area.left, top=area.top + 1.5,
        width=area.width, height=area.height - 1.5,
    )
    per_row = get_token(theme, "layout.numbered_list_card_per_row", 2)
    rows = (len(items) + per_row - 1) // per_row
    matrix = grid_matrix(cards_area, rows=rows, cols=per_row,
                         gutter=get_token(theme, "spacing.grid_gutter_inch", 0.25))

    for idx, item in enumerate(items):
        r = idx // per_row
        c = idx % per_row
        if r >= len(matrix):
            break
        cell = matrix[r][c]

        # 卡片背景
        card = add_card(slide, cell, "card", ctx, theme,
                        radius=get_token(theme, "effect.card_radius", True))

        # 序号徽章（左上角）
        badge_size = 0.5
        badge_rect = GridRect(
            left=cell.left + 0.15, top=cell.top + 0.15,
            width=badge_size, height=badge_size,
        )
        render_number_badge(
            slide, badge_rect, f"{idx + 1:02d}", "number", ctx, theme,
        )

        # 卡片文本
        text_rect = GridRect(
            left=cell.left + 0.15 + badge_size + 0.15, top=cell.top + 0.15,
            width=cell.width - badge_size - 0.45, height=badge_size,
        )
        add_text_box(
            slide, text_rect, str(item), "body", ctx, theme,
            font_size_pt=get_token(theme, "font.body.size_pt", 16),
            bold=False,
            color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
            align="left", anchor="middle",
        )


# ==================== kpi KPI 页：等宽横向卡片矩阵 ====================
@register_layout("kpi")
def layout_kpi(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """KPI 页布局：等宽横向卡片矩阵

    版式：
        - 顶部标题区（占 20%）
        - 下方 KPI 卡片矩阵（1 行 N 列，N 由 kpi_items 数量决定，最多 4 列）
    字段：
        - title: 页面标题
        - kpi_items[{label, value, trend}]: KPI 数据
    元素角色：
        - title: 页面标题
        - card: 每个 KPI 卡片
        - number: KPI 数值
        - desc: KPI 标签
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

    # 顶部标题区
    title_rect = GridRect(area.left, area.top, area.width, 1.0)
    add_text_box(
        slide, title_rect, page_data.get("title", ""), "title", ctx, theme,
        font_size_pt=get_token(theme, "font.title.size_pt", 36),
        bold=True,
        color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
        align="left", anchor="middle",
    )

    # 装饰下划线
    add_line(
        slide, area.left, area.top + 1.1,
        area.left + 1.5, area.top + 1.1,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    # KPI 卡片矩阵
    kpi_items = page_data.get("kpi_items", [])
    if not kpi_items:
        return

    cards_area = GridRect(
        left=area.left, top=area.top + 1.5,
        width=area.width, height=area.height - 1.5,
    )
    # KPI 卡片固定 1 行，列数 = KPI 数量（最多 4）
    cols = min(len(kpi_items), 4)
    matrix = grid_matrix(cards_area, rows=1, cols=cols,
                         gutter=get_token(theme, "spacing.grid_gutter_inch", 0.25))

    for idx, item in enumerate(kpi_items):
        if idx >= cols:
            break
        cell = matrix[0][idx]

        # KPI 卡片背景
        add_card(slide, cell, "card", ctx, theme,
                 radius=get_token(theme, "effect.card_radius", True))

        # KPI 数值（大号居中，上半部分）
        value_text = str(item.get("value", ""))
        value_rect = GridRect(
            left=cell.left, top=cell.top + cell.height * 0.2,
            width=cell.width, height=cell.height * 0.5,
        )
        add_text_box(
            slide, value_rect, value_text, "kpi_number", ctx, theme,
            font_size_pt=get_token(theme, "font.kpi_number.size_pt", 48),
            bold=get_token(theme, "font.kpi_number.bold", True),
            color=resolve_color(theme, get_token(theme, "color.kpi_number", "#1A56DB")),
            align="center", anchor="middle",
        )

        # KPI 标签（小号居中，下半部分）
        label_text = str(item.get("label", ""))
        label_rect = GridRect(
            left=cell.left, top=cell.top + cell.height * 0.7,
            width=cell.width, height=cell.height * 0.25,
        )
        add_text_box(
            slide, label_rect, label_text, "desc", ctx, theme,
            font_size_pt=get_token(theme, "font.kpi_label.size_pt", 14),
            bold=False,
            color=resolve_color(theme, get_token(theme, "color.text_secondary", "#6B7280")),
            align="center", anchor="middle",
        )


# ==================== timeline 时间轴页：水平时间轴 + 节点事件（T603 新增） ====================
@register_layout("timeline")
def layout_timeline(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """时间轴页布局：水平时间轴 + 节点事件

    版式：
        - 顶部标题区（占 18%）
        - 中部水平轴线 + 时间节点（圆点徽章 + 年份/时间标签）
        - 下方事件描述卡片（交替上下排列，避免拥挤）
    字段：
        - title: 页面标题
        - timeline_items[{time, event}]: 时间节点列表（3~6 项）
    元素角色：
        - title: 页面标题
        - year: 每个节点的时间标签
        - desc: 每个节点的事件描述
        - number: 节点圆点徽章
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

    # 顶部标题区
    title_rect = GridRect(area.left, area.top, area.width, 1.0)
    add_text_box(
        slide, title_rect, page_data.get("title", ""), "title", ctx, theme,
        font_size_pt=get_token(theme, "font.title.size_pt", 36),
        bold=True,
        color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
        align="left", anchor="middle",
    )
    add_line(
        slide, area.left, area.top + 1.1,
        area.left + 1.5, area.top + 1.1,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    items = page_data.get("timeline_items", [])
    if not items:
        return

    # 时间轴区域
    axis_area = GridRect(
        left=area.left, top=area.top + 1.8,
        width=area.width, height=area.height - 2.0,
    )
    n = len(items)
    # 水平等分 N 个节点
    cell_width = axis_area.width / n
    axis_y = axis_area.top + axis_area.height * 0.5

    # 主轴线
    add_line(
        slide, axis_area.left, axis_y,
        axis_area.left + axis_area.width, axis_y,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    badge_size = 0.35
    for idx, item in enumerate(items):
        cx = axis_area.left + cell_width * (idx + 0.5)
        # 交替上下排列：偶数在下，奇数在上
        below = (idx % 2 == 0)

        # 节点圆点
        badge_rect = GridRect(
            left=cx - badge_size / 2, top=axis_y - badge_size / 2,
            width=badge_size, height=badge_size,
        )
        render_number_badge(
            slide, badge_rect, f"{idx + 1:02d}", "number", ctx, theme,
        )

        # 时间标签（紧贴轴线）
        time_text = str(item.get("time", ""))
        time_rect_h = 0.4
        time_rect = GridRect(
            left=cx - cell_width * 0.45, top=axis_y + badge_size * 0.6,
            width=cell_width * 0.9, height=time_rect_h,
        ) if below else GridRect(
            left=cx - cell_width * 0.45, top=axis_y - badge_size * 0.6 - time_rect_h,
            width=cell_width * 0.9, height=time_rect_h,
        )
        add_text_box(
            slide, time_rect, time_text, "year", ctx, theme,
            font_size_pt=get_token(theme, "font.year.size_pt", 18),
            bold=get_token(theme, "font.year.bold", True),
            color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
            align="center", anchor="middle",
        )

        # 事件描述（远离轴线）
        event_text = str(item.get("event", ""))
        # T612：弹性约束接入 - 描述区高度按文本长度估算
        # 默认 1.3 inch，长文本自动增高（上限 2.0 inch），短文本自动缩减（下限 0.8 inch）
        elastic_enabled = get_token(theme, "layout.elastic_enabled", True)
        if elastic_enabled and event_text:
            body_font_pt = get_token(theme, "font.body.size_pt", 14)
            box_w_emu = int(cell_width * 0.9 * EMU_PER_INCH)
            est_h_emu = estimate_text_height(event_text, body_font_pt, box_w_emu)
            # 加上下内边距 0.2 inch
            est_h_inch = est_h_emu / EMU_PER_INCH + 0.2
            desc_rect_h = max(0.8, min(2.0, est_h_inch))
        else:
            desc_rect_h = 1.3
        desc_rect = GridRect(
            left=cx - cell_width * 0.45, top=time_rect.top + time_rect_h + 0.1,
            width=cell_width * 0.9, height=desc_rect_h,
        ) if below else GridRect(
            left=cx - cell_width * 0.45, top=time_rect.top - desc_rect_h - 0.1,
            width=cell_width * 0.9, height=desc_rect_h,
        )
        # 超出底部时截断
        if desc_rect.top + desc_rect.height > area.top + area.height:
            desc_rect = GridRect(
                desc_rect.left, area.top + area.height - desc_rect.height,
                desc_rect.width, desc_rect.height,
            )
        add_text_box(
            slide, desc_rect, event_text, "desc", ctx, theme,
            font_size_pt=get_token(theme, "font.body.size_pt", 14),
            bold=False,
            color=resolve_color(theme, get_token(theme, "color.text_secondary", "#6B7280")),
            align="center", anchor="top",
        )


# ==================== two_column 双栏页：左右分栏对照（T603 新增） ====================
@register_layout("two_column")
def layout_two_column(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """双栏页布局：左右分栏对照

    版式：
        - 顶部标题区（占 18%）
        - 左右等分两栏，每栏含小标题 + 列表条目
        - 中间垂直分割线
    字段：
        - title: 页面标题
        - left_title: 左栏标题
        - left_items[]: 左栏条目
        - right_title: 右栏标题
        - right_items[]: 右栏条目
    元素角色：
        - title: 页面标题
        - subtitle: 左/右栏小标题
        - body: 每个条目
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

    # 顶部标题区
    title_rect = GridRect(area.left, area.top, area.width, 1.0)
    add_text_box(
        slide, title_rect, page_data.get("title", ""), "title", ctx, theme,
        font_size_pt=get_token(theme, "font.title.size_pt", 36),
        bold=True,
        color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
        align="left", anchor="middle",
    )
    add_line(
        slide, area.left, area.top + 1.1,
        area.left + 1.5, area.top + 1.1,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    # 左右等分
    content_area = GridRect(
        left=area.left, top=area.top + 1.6,
        width=area.width, height=area.height - 1.6,
    )
    gutter = get_token(theme, "spacing.grid_gutter_inch", 0.3)
    half_w = (content_area.width - gutter) / 2
    left_rect = GridRect(content_area.left, content_area.top, half_w, content_area.height)
    right_rect = GridRect(
        content_area.left + half_w + gutter, content_area.top,
        half_w, content_area.height,
    )

    # 中间垂直分割线
    mid_x = content_area.left + half_w + gutter / 2
    add_line(
        slide, mid_x, content_area.top,
        mid_x, content_area.top + content_area.height,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.divider", "#E5E7EB")),
        width_pt=1.0,
    )

    # 渲染单栏（复用逻辑）
    def _render_column(rect: GridRect, col_title: str, items: list, role_prefix: str) -> None:
        # 小标题
        sub_rect = GridRect(rect.left, rect.top, rect.width, 0.7)
        add_text_box(
            slide, sub_rect, col_title, f"{role_prefix}_subtitle", ctx, theme,
            font_size_pt=get_token(theme, "font.subtitle.size_pt", 22),
            bold=True,
            color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
            align="left", anchor="middle",
        )
        # 小标题下划线
        add_line(
            slide, rect.left, rect.top + 0.8,
            rect.left + 0.8, rect.top + 0.8,
            "divider", ctx, theme,
            color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
            width_pt=1.5,
        )
        # 条目列表
        if not items:
            return
        # T612：弹性约束接入 - 长短条目自适应高度
        elastic_enabled = get_token(theme, "layout.elastic_enabled", True)
        item_h = 0.6
        gap = get_token(theme, "spacing.element_gap_inch", 0.15)
        start_y = rect.top + 1.1
        list_area_h = rect.top + rect.height - start_y

        if elastic_enabled:
            items_data = [{"text": str(it)} for it in items]
            positions = elastic_distribute_items(
                items_data,
                area_top_emu=int(start_y * EMU_PER_INCH),
                area_height_emu=int(list_area_h * EMU_PER_INCH),
                area_width_emu=int(rect.width * EMU_PER_INCH),
                font_size_pt=get_token(theme, "font.body.size_pt", 16),
                min_item_h_inch=0.4,
                gap_inch=gap,
            )
        else:
            positions = None

        for idx, item in enumerate(items):
            if positions is not None:
                top_emu, h_emu = positions[idx]
                cur_y = top_emu / EMU_PER_INCH
                cur_h = h_emu / EMU_PER_INCH
                if cur_y + cur_h > rect.top + rect.height:
                    break
            else:
                cur_y = start_y + idx * (item_h + gap)
                cur_h = item_h
                if cur_y + cur_h > rect.top + rect.height:
                    break
            item_rect = GridRect(rect.left, cur_y, rect.width, cur_h)
            add_text_box(
                slide, item_rect, f"· {item}", f"{role_prefix}_body", ctx, theme,
                font_size_pt=get_token(theme, "font.body.size_pt", 16),
                bold=False,
                color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
                align="left", anchor="middle",
            )

    _render_column(left_rect, page_data.get("left_title", ""), page_data.get("left_items", []), "left")
    _render_column(right_rect, page_data.get("right_title", ""), page_data.get("right_items", []), "right")


# ==================== table 表格页：原生 PPT 表格 + 表头样式继承（T603 新增） ====================
@register_layout("table")
def layout_table(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """表格页布局：原生 PPT 表格

    版式：
        - 顶部标题区（占 18%）
        - 下方原生 GraphicFrame 表格（headers + rows）
        - 表头主题色背景 + 白字，正文行斑马纹
    字段：
        - title: 页面标题
        - headers[]: 表头列表
        - rows[][]: 数据行二维数组
    元素角色：
        - title: 页面标题
        - table: 表格本体（GraphicFrame）
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

    # 顶部标题区
    title_rect = GridRect(area.left, area.top, area.width, 1.0)
    add_text_box(
        slide, title_rect, page_data.get("title", ""), "title", ctx, theme,
        font_size_pt=get_token(theme, "font.title.size_pt", 36),
        bold=True,
        color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
        align="left", anchor="middle",
    )
    add_line(
        slide, area.left, area.top + 1.1,
        area.left + 1.5, area.top + 1.1,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    headers = page_data.get("headers", [])
    rows = page_data.get("rows", [])
    if not headers or not rows:
        return

    # 表格区域
    table_area = GridRect(
        left=area.left, top=area.top + 1.6,
        width=area.width, height=area.height - 1.6,
    )
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    left, top, width, height = table_area.to_emu()
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table

    # 标记 GraphicFrame 的 role + shape_id（用于动画匹配）
    # 使用 ctx.register 正式注册到 elements 列表，供动画模块匹配
    ctx.register("table", graphic_frame)

    # 关闭表格默认样式（避免主题样式干扰）
    tbl = table._tbl
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')

    # 主题色
    header_bg = get_token(theme, "color.table_header_bg",
                          get_token(theme, "color.primary", "#1A56DB"))
    header_text = get_token(theme, "color.table_header_text", "#FFFFFF")
    row_bg = get_token(theme, "color.table_row_bg", "#FFFFFF")
    zebra_bg = get_token(theme, "color.table_zebra_bg", "#F3F4F6")
    row_text = get_token(theme, "color.text_primary", "#1F2937")
    border_color = get_token(theme, "color.table_border", "#E5E7EB")
    header_font_pt = get_token(theme, "font.table_header.size_pt", 14)
    body_font_pt = get_token(theme, "font.table_body.size_pt", 12)

    # 表头行
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(header_bg)
        _style_cell_text(cell, hex_to_rgb(header_text), header_font_pt, bold=True,
                         align="center", anchor="middle")

    # 数据行
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            if c >= n_cols:
                break
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            # 斑马纹
            cell.fill.fore_color.rgb = hex_to_rgb(row_bg if r % 2 == 0 else zebra_bg)
            _style_cell_text(cell, hex_to_rgb(row_text), body_font_pt, bold=False,
                             align="left", anchor="middle")

    # 应用边框（通过 XML 注入，确保所有单元格统一边框样式）
    _apply_table_borders(table, hex_to_rgb(border_color), width_pt=0.5)


def _style_cell_text(cell, color_rgb, size_pt: float, bold: bool,
                     align: str = "left", anchor: str = "middle") -> None:
    """样式化单元格文本

    :param cell: python-pptx _Cell 对象
    :param color_rgb: RGBColor 文本色
    :param size_pt: 字号 pt
    :param bold: 是否加粗
    :param align: 对齐方式 left/center/right
    :param anchor: 垂直对齐 top/middle/bottom
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf = cell.text_frame
    tf.word_wrap = True
    cell.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.MIDDLE)
    for para in tf.paragraphs:
        para.alignment = align_map.get(align, PP_ALIGN.LEFT)
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = color_rgb
            # 字体名继承主题
            run.font.name = None


def _apply_table_borders(table, color_rgb, width_pt: float = 0.5) -> None:
    """为表格所有单元格应用统一边框（通过 XML 注入）

    python-pptx 未提供高级边框 API，需通过 lxml 直接操作 a:tcPr 子元素。

    :param table: python-pptx Table 对象
    :param color_rgb: RGBColor 边框色
    :param width_pt: 边框宽度 pt
    """
    from pptx.util import Pt
    width_emu = int(Pt(width_pt))
    hex_str = f"{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"
    for row in table.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            # 清除已有 lnL/lnR/lnT/lnB
            for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
                existing = tcPr.find(qn(tag))
                if existing is not None:
                    tcPr.remove(existing)
            # 注入四边边框
            for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
                ln = tcPr.makeelement(qn(tag), {'w': str(width_emu), 'cap': 'flat'})
                fill = ln.makeelement(qn('a:solidFill'), {})
                clr = fill.makeelement(qn('a:srgbClr'), {'val': hex_str})
                fill.append(clr)
                ln.append(fill)
                tcPr.append(ln)


# ==================== T613 layout_variant 版式变体 ====================
# 以下为各 page_type 的非 default 变体实现，通过 @register_layout(name, variant=...)
# 注册到二级注册表，dispatch 时按 variant_override 或 page_data.layout_variant 选择


@register_layout("kpi", variant="grid_2x2")
def layout_kpi_grid_2x2(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """KPI 页变体：2×2 网格布局（最多 4 个 KPI，每个占用更大空间）

    适用于 KPI 数量 ≤4 且数值重要性高的场景，
    相比 default 的 1×4 横向排布，2×2 网格让每个卡片更大更醒目。
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

    # 顶部标题区
    title_rect = GridRect(area.left, area.top, area.width, 1.0)
    add_text_box(
        slide, title_rect, page_data.get("title", ""), "title", ctx, theme,
        font_size_pt=get_token(theme, "font.title.size_pt", 36),
        bold=True,
        color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
        align="left", anchor="middle",
    )
    add_line(
        slide, area.left, area.top + 1.1,
        area.left + 1.5, area.top + 1.1,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    kpi_items = page_data.get("kpi_items", [])
    if not kpi_items:
        return

    # 2×2 网格区域
    cards_area = GridRect(
        left=area.left, top=area.top + 1.5,
        width=area.width, height=area.height - 1.5,
    )
    # 固定 2 行 2 列，最多展示 4 个 KPI
    matrix = grid_matrix(cards_area, rows=2, cols=2,
                         gutter=get_token(theme, "spacing.grid_gutter_inch", 0.25))

    for idx, item in enumerate(kpi_items):
        if idx >= 4:
            break
        r = idx // 2
        c = idx % 2
        cell = matrix[r][c]

        # KPI 卡片背景
        add_card(slide, cell, "card", ctx, theme,
                 radius=get_token(theme, "effect.card_radius", True))

        # KPI 数值（大号居中，上半部分）
        value_text = str(item.get("value", ""))
        value_rect = GridRect(
            left=cell.left, top=cell.top + cell.height * 0.2,
            width=cell.width, height=cell.height * 0.5,
        )
        add_text_box(
            slide, value_rect, value_text, "kpi_number", ctx, theme,
            font_size_pt=get_token(theme, "font.kpi_number.size_pt", 48),
            bold=get_token(theme, "font.kpi_number.bold", True),
            color=resolve_color(theme, get_token(theme, "color.kpi_number", "#1A56DB")),
            align="center", anchor="middle",
        )

        # KPI 标签（小号居中，下半部分）
        label_text = str(item.get("label", ""))
        label_rect = GridRect(
            left=cell.left, top=cell.top + cell.height * 0.7,
            width=cell.width, height=cell.height * 0.25,
        )
        add_text_box(
            slide, label_rect, label_text, "desc", ctx, theme,
            font_size_pt=get_token(theme, "font.kpi_label.size_pt", 14),
            bold=False,
            color=resolve_color(theme, get_token(theme, "color.text_secondary", "#6B7280")),
            align="center", anchor="middle",
        )


@register_layout("numbered_list", variant="single_column")
def layout_numbered_list_single_column(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """编号列表页变体：单列纵向列表（条目数较多时使用）

    相比 default 的 2 列卡片网格，单列布局适合 5~8 个条目的场景，
    每条占用整行宽度，文本可读性更好。
    """
    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

    # 顶部标题区
    title_rect = GridRect(area.left, area.top, area.width, 1.0)
    add_text_box(
        slide, title_rect, page_data.get("title", ""), "title", ctx, theme,
        font_size_pt=get_token(theme, "font.title.size_pt", 36),
        bold=True,
        color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
        align="left", anchor="middle",
    )
    add_line(
        slide, area.left, area.top + 1.1,
        area.left + 1.5, area.top + 1.1,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    items = page_data.get("items", [])
    if not items:
        return

    # 单列卡片区域
    cards_area = GridRect(
        left=area.left, top=area.top + 1.5,
        width=area.width, height=area.height - 1.5,
    )
    n = len(items)
    matrix = grid_matrix(cards_area, rows=n, cols=1,
                         gutter=get_token(theme, "spacing.grid_gutter_inch", 0.2))

    for idx, item in enumerate(items):
        if idx >= n:
            break
        cell = matrix[idx][0]

        # 卡片背景
        add_card(slide, cell, "card", ctx, theme,
                 radius=get_token(theme, "effect.card_radius", True))

        # 序号徽章（左上角）
        badge_size = 0.5
        badge_rect = GridRect(
            left=cell.left + 0.15, top=cell.top + (cell.height - badge_size) / 2,
            width=badge_size, height=badge_size,
        )
        render_number_badge(
            slide, badge_rect, f"{idx + 1:02d}", "number", ctx, theme,
        )

        # 卡片文本（垂直居中）
        text_rect = GridRect(
            left=cell.left + 0.15 + badge_size + 0.2, top=cell.top,
            width=cell.width - badge_size - 0.5, height=cell.height,
        )
        add_text_box(
            slide, text_rect, str(item), "body", ctx, theme,
            font_size_pt=get_token(theme, "font.body.size_pt", 16),
            bold=False,
            color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
            align="left", anchor="middle",
        )


# ==================== T614 image 图片页：单图 / 多图网格 ====================
@register_layout("image")
def layout_image(slide: Slide, page_data: dict, theme: dict, ctx: LayoutContext) -> None:
    """图片页布局：单图大图 / 双图 / 四图网格

    支持两种输入：
        1. image_path：单图模式，大图居中展示
        2. image_items[]：多图模式，按 image_layout 自动选择 grid_2 / grid_2x2

    版式选择规则：
        - image_layout="single" 或仅 image_path：单图大图
        - image_layout="grid_2" 或 image_items 长度 ≤2：双图并排
        - image_layout="grid_2x2" 或 image_items 长度 ≥3：四图 2×2 网格

    字段：
        - title: 页面标题
        - image_path: 单图路径（与 image_items 二选一）
        - image_items[{path, caption}]: 多图列表（与 image_path 二选一）
        - image_layout: 显式版式 "single"/"grid_2"/"grid_2x2"（可选）

    元素角色：
        - title: 页面标题
        - picture: 图片元素
        - desc: 图片说明文字
    """
    from PIL import Image as PILImage
    from aippt.image_replacer import resolve_image_source, compute_fit_rect

    area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

    # 顶部标题区
    title_rect = GridRect(area.left, area.top, area.width, 1.0)
    add_text_box(
        slide, title_rect, page_data.get("title", ""), "title", ctx, theme,
        font_size_pt=get_token(theme, "font.title.size_pt", 36),
        bold=True,
        color=resolve_color(theme, get_token(theme, "color.text_primary", "#1F2937")),
        align="left", anchor="middle",
    )
    add_line(
        slide, area.left, area.top + 1.1,
        area.left + 1.5, area.top + 1.1,
        "divider", ctx, theme,
        color=resolve_color(theme, get_token(theme, "color.primary", "#1A56DB")),
        width_pt=2.0,
    )

    # 图片区域
    image_area = GridRect(
        left=area.left, top=area.top + 1.5,
        width=area.width, height=area.height - 1.5,
    )

    # 收集图片项，统一为 [{path, caption}] 结构
    image_items_field = page_data.get("image_items")
    single_path = page_data.get("image_path")
    explicit_layout = page_data.get("image_layout")

    if image_items_field:
        items = [{"path": it.get("path", ""), "caption": it.get("caption", "")}
                 for it in image_items_field if it.get("path")]
    elif single_path:
        items = [{"path": single_path, "caption": ""}]
    else:
        logger.warning("image 页无 image_path 或 image_items，跳过渲染")
        return

    if not items:
        return

    # 决定版式
    n = len(items)
    if explicit_layout:
        layout_mode = explicit_layout
    elif n == 1:
        layout_mode = "single"
    elif n <= 2:
        layout_mode = "grid_2"
    else:
        layout_mode = "grid_2x2"

    # 计算每个图片单元的 GridRect
    gutter = get_token(theme, "spacing.grid_gutter_inch", 0.25)
    if layout_mode == "single":
        cells = [image_area]
    elif layout_mode == "grid_2":
        left_cell, right_cell = split_horizontal(image_area, ratio=0.5, gutter=gutter)
        cells = [left_cell, right_cell]
    else:  # grid_2x2
        matrix = grid_matrix(image_area, rows=2, cols=2, gutter=gutter)
        cells = [matrix[r][c] for r in range(2) for c in range(2)]

    # 渲染每个图片
    for idx, item in enumerate(items):
        if idx >= len(cells):
            break
        cell = cells[idx]
        # 图片框预留底部说明文字空间（如有 caption）
        caption_text = item.get("caption", "")
        if caption_text:
            pic_rect = GridRect(cell.left, cell.top, cell.width, cell.height - 0.4)
            cap_rect = GridRect(cell.left, cell.top + cell.height - 0.35,
                                cell.width, 0.35)
        else:
            pic_rect = cell
            cap_rect = None

        # 解析图片源（支持本地路径 / URL）
        try:
            img_path = resolve_image_source(item, cache_dir=None)
        except Exception as e:
            logger.warning("图片 %s 解析失败: %s，跳过", item.get("path"), e)
            # 渲染占位框
            add_text_box(
                slide, pic_rect, f"[图片加载失败: {item.get('path', '')}]",
                "picture", ctx, theme,
                font_size_pt=12, bold=False,
                color=resolve_color(theme, get_token(theme, "color.text_secondary", "#6B7280")),
                align="center", anchor="middle",
            )
            continue

        # 等比缩放：保持图片原始宽高比，contain 模式适配 pic_rect
        try:
            with PILImage.open(img_path) as pil_img:
                orig_w, orig_h = pil_img.size
        except Exception:
            # 无法读取尺寸，直接按 rect 填充
            orig_w, orig_h = int(pic_rect.width * 914400), int(pic_rect.height * 914400)

        from pptx.util import Emu
        box_w_emu = int(pic_rect.width * 914400)
        box_h_emu = int(pic_rect.height * 914400)
        final_left, final_top, final_w, final_h = compute_fit_rect(
            target_left=int(pic_rect.left * 914400),
            target_top=int(pic_rect.top * 914400),
            target_w=box_w_emu,
            target_h=box_h_emu,
            img_w=orig_w,
            img_h=orig_h,
            fit="contain",
        )

        # 插入图片
        pic = slide.shapes.add_picture(
            img_path,
            Emu(final_left), Emu(final_top),
            Emu(final_w), Emu(final_h),
        )
        # 标记角色（用于动画匹配）
        ctx.register("picture", pic)

        # 图片说明文字
        if cap_rect and caption_text:
            add_text_box(
                slide, cap_rect, caption_text, "desc", ctx, theme,
                font_size_pt=get_token(theme, "font.caption.size_pt", 12),
                bold=False,
                color=resolve_color(theme, get_token(theme, "color.text_secondary", "#6B7280")),
                align="center", anchor="middle",
            )
