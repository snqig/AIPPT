import os
import re
import json
import argparse
import tempfile
import colorsys
from pptx import Presentation
from pathlib import Path
from collections import defaultdict
from typing import Any, Optional
from lxml import etree

from aippt.config import PAGE_TYPE_KEYWORDS, MODELS_ROOT
from aippt.constants import SLOT_MATCH_KEYWORDS
from aippt.logger import logger
from aippt.ppt_element_classifier import classify_page

# meta必填字段校验
META_REQUIRED_FIELDS = ['template_id', 'category', 'total_pages', 'chapters', 'page_slots']

# DrawingML 命名空间
_NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# ==================== 核心解析函数 ====================
def extract_page_texts(slide) -> list[dict[str, Any]]:
    texts: list[dict[str, Any]] = []

    def _extract_from_shape(shape) -> None:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                width_emu: Optional[int] = None
                height_emu: Optional[int] = None
                try:
                    width_emu = int(shape.width)
                    height_emu = int(shape.height)
                except Exception:
                    pass
                font_size_pt: Optional[float] = None
                try:
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            fs = para.runs[0].font.size
                            if fs is not None:
                                font_size_pt = fs.pt
                                break
                except Exception:
                    pass

                texts.append({
                    'text': text,
                    'shape_name': shape.name,
                    'is_placeholder': shape.is_placeholder,
                    'width_emu': width_emu,
                    'height_emu': height_emu,
                    'font_size_pt': font_size_pt,
                })
        if shape.shape_type == 6:
            for sub in shape.shapes:
                _extract_from_shape(sub)

    for shape in slide.shapes:
        _extract_from_shape(shape)
    return texts


def detect_page_shapes(slide) -> dict[str, Any]:
    flags: dict[str, Any] = {"has_chart": False, "has_table": False, "has_picture": False,
                              "chart_type": None, "picture_count": 0,
                              "chart_shape_name": None, "table_shape_name": None}

    def _scan(shape) -> None:
        # GRAPH_CHART = 3
        if shape.shape_type == 3 or (hasattr(shape, "has_chart") and shape.has_chart):
            flags["has_chart"] = True
            # 记录第一个图表形状名（用于 chart_slot 标识）
            if not flags["chart_shape_name"]:
                flags["chart_shape_name"] = shape.name
            try:
                if hasattr(shape, "chart") and shape.chart is not None:
                    flags["chart_type"] = str(shape.chart.chart_type)
            except Exception:
                pass
        # TABLE = 19 / has_table
        elif (hasattr(shape, "has_table") and shape.has_table) or shape.shape_type == 19:
            flags["has_table"] = True
            # 记录第一个表格形状名（用于 table_slot 标识）
            if not flags["table_shape_name"]:
                flags["table_shape_name"] = shape.name
        # PICTURE = 13
        elif shape.shape_type == 13:
            flags["has_picture"] = True
            flags["picture_count"] += 1
        # GROUP 递归
        if shape.shape_type == 6:
            for sub in shape.shapes:
                _scan(sub)

    for shape in slide.shapes:
        try:
            _scan(shape)
        except Exception:
            continue
    return flags

def detect_page_type(page_texts: list[dict[str, Any]], page_idx: int, total_pages: int) -> str:
    all_text = ' '.join([t['text'] for t in page_texts])

    # 首页强制识别为封面
    if page_idx == 0:
        return 'cover'

    # 目录页优先判断（CONTENTS/目录是强信号）：放在 end 之前，避免目录页含 THANK YOU 装饰文本被误判为 end
    # 放宽文本框数量限制（许多模板目录页有15+文本框）
    catalog_strong = ['CONTENTS', '目 录', '目录页', '目录', '目\n录']
    for kw in catalog_strong:
        if kw in all_text and len(page_texts) <= 25:
            # 但若同时含 PART ONE/第N章 等强章节信号，且文本框少，则可能是章节页带目录样式
            if len(page_texts) <= 6 and any(k in all_text for k in ['PART ONE', 'PART TWO', 'PART.0', 'Part.0']):
                pass  # 让章节判断处理
            else:
                return 'catalog'

    # 末尾3页判断版权页（含版权关键词）
    if page_idx >= total_pages - 3:
        for kw in PAGE_TYPE_KEYWORDS['copyright']:
            if kw.lower() in all_text.lower():
                return 'copyright'
        # 末页若与首页样式重复（封面/装饰页重复），识别为 copyright
        if page_idx == total_pages - 1:
            cover_kw = ['WORK REPORT', 'WORK SUMMARY', 'BUSINESS', 'YOUR LOGO', 'COMPANY']
            # 末页若非结束页且无明显内容，归为版权/装饰页
            if not any(k in all_text for k in ['感谢聆听', '谢谢聆听', 'THANK YOU', '感谢各位', '谢谢观看', '感谢观看']):
                if any(k in all_text for k in cover_kw) or 'CONTENTS' in all_text:
                    return 'copyright'

    # 结束页：要求在末5页内，避免目录页 THANK YOU 装饰文本误命中
    if page_idx >= total_pages - 5:
        end_strong = ['感谢聆听', '谢谢聆听', '谢谢观看', '感谢观看', 'THANK YOU', '感谢各位']
        for kw in end_strong:
            if kw in all_text:
                return 'end'
    if page_idx >= total_pages - 3:
        for kw in ['感谢', '致谢', '谢谢']:
            if kw in all_text:
                return 'end'

    # 章节页优先判断（含 PART ONE/第N章/Part.01 等强信号）
    chapter_strong = ['PART ONE', 'PART TWO', 'PART THREE', 'PART FOUR', 'PART FIVE', 'PART SIX', 'PART SEVEN',
                      'Part.0', 'Part.1', 'Part.01', 'Part.02', 'Part.03', 'Part.04', 'Part.05', 'Part.06', 'Part.07',
                      'PART.0', 'PART.1', 'PART.01', 'PART.02', 'PART.03', 'PART.04', 'PART.05', 'PART.06', 'PART.07',
                      '第1部分', '第2部分', '第3部分', '第4部分', '第1章', '第2章', '第3章', '第4章',
                      'COMPETENCY', 'JOB AWARENESS', 'PART 01', 'PART 02', 'PART 03', 'PART 04']
    for kw in chapter_strong:
        if kw in all_text:
            return 'chapter'
    # 通用章节关键词（需文本框少，避免误判内容页）
    for kw in ['PART ', '第', '章']:
        if kw in all_text and len(page_texts) <= 6:
            return 'chapter'

    # 章节分隔页（数字+标题模式）：文本框少且含纯数字 01-09 + 短中文标题，无目录关键词
    if len(page_texts) <= 6:
        has_number = any(re.match(r'^0[1-9]$', t['text'].strip()) for t in page_texts)
        has_short_cn_title = any(2 <= len(t['text'].strip()) <= 10 and re.search(r'[\u4e00-\u9fa5]', t['text'])
                                 and t['text'].strip() not in ['目录', '目 录'] for t in page_texts)
        no_catalog_kw = not any(k in all_text for k in catalog_strong)
        no_end_kw = not any(k in all_text for k in ['感谢', 'THANK', '致谢'])
        if has_number and has_short_cn_title and no_catalog_kw and no_end_kw:
            return 'chapter'

    return 'content'

def extract_slots(page_texts: list[dict[str, Any]], page_idx: int) -> list[dict[str, Any]]:
    slots = []
    title_counter = 0
    desc_counter = 0
    item_counter = 0
    year_counter = 0
    percent_counter = 0
    number_counter = 0

    for item in page_texts:
        text = item['text']
        # 判断是否为可替换槽位
        is_slot = any(kw in text for kw in SLOT_MATCH_KEYWORDS)
        is_short_title = len(text) < 15 and '\n' not in text

        if not (is_slot or is_short_title):
            continue

        stripped = text.strip()
        # 语义化命名（优先匹配强语义模式，避免数字/年份/百分比被误归为 title）
        if '汇报人' in text or '姓名' in text:
            slot_name = 'reporter'
        elif '年度' in text or '2O2X' in text or '202X' in text:
            slot_name = 'period'
        elif re.match(r'^\d{4}$', stripped):  # 纯年份，如 2019、2027
            year_counter += 1
            slot_name = f'year_{year_counter}' if year_counter > 1 else 'year'
        elif re.match(r'^\d+%$', stripped):  # 百分比，如 67%
            percent_counter += 1
            slot_name = f'percent_{percent_counter}' if percent_counter > 1 else 'percent'
        elif re.match(r'^\d+$', stripped):  # 纯数字，如 01、02
            number_counter += 1
            slot_name = f'number_{number_counter}' if number_counter > 1 else 'number'
        elif '标题' in text or 'TITLE' in text.upper() or '某某' in text or is_short_title:
            title_counter += 1
            slot_name = f'title_{title_counter}' if title_counter > 1 else 'title'
        elif '内容' in text or '录入' in text or '输入' in text or 'text content' in text.lower() or 'your text' in text.lower() or len(text) > 20:
            desc_counter += 1
            slot_name = f'desc_{desc_counter}' if desc_counter > 1 else 'desc'
        elif '项目' in text and ('名称' in text or 'PROJECT' in text.upper()):
            item_counter += 1
            slot_name = f'project_name_{item_counter}'
        elif '%' in text:
            item_counter += 1
            slot_name = f'progress_{item_counter}'
        else:
            item_counter += 1
            slot_name = f'item_{item_counter}'

        # 取前20字符作为匹配特征
        match_text = text[:20].replace('\n', ' ').strip()
        if match_text:
            slot_entry = {
                'slot': slot_name,
                'match_text': match_text,
                'shape_name': item.get('shape_name', '')
            }
            # 计算 capacity（事前预警用）：基于 shape 几何和字号
            capacity = _compute_capacity(item)
            if capacity:
                slot_entry['capacity'] = capacity
            slots.append(slot_entry)

    return slots


def _compute_capacity(item: dict[str, Any]) -> Optional[dict[str, int]]:
    width_emu = item.get('width_emu')
    height_emu = item.get('height_emu')
    font_size_pt = item.get('font_size_pt')
    if not width_emu or not height_emu or not font_size_pt:
        return None

    EMU_PER_PT = 12700
    char_w = font_size_pt * 0.55 * EMU_PER_PT
    line_h = font_size_pt * 1.2 * EMU_PER_PT
    if char_w <= 0 or line_h <= 0:
        return None
    max_chars_per_line = max(1, int(width_emu / char_w))
    max_lines = max(1, int(height_emu / line_h))
    total_chars = max_chars_per_line * max_lines
    # 容量过小（<5）通常意味着 shape 是装饰性小框/角标，capacity 不可靠
    # 此时不写入 capacity，让渲染器回退到运行时几何估算，避免误报预警
    if total_chars < 5:
        return None
    return {
        "max_chars_per_line": max_chars_per_line,
        "max_lines": max_lines,
        "total_chars": total_chars,
    }

def detect_chapters(slides, page_types: list[str]) -> list[dict[str, Any]]:
    chapters = []
    chapter_idx = 0
    chapter_start = None
    chapter_name = ''
    
    for i, ptype in enumerate(page_types):
        if ptype == 'chapter':
            if chapter_start is not None:
                chapters.append({
                    'key': f'chapter_{chapter_idx}',
                    'start_page': chapter_start + 1,
                    'end_page': i,
                    'name': chapter_name
                })
            chapter_idx += 1
            chapter_start = i
            # 提取中文章节名
            page_text = ' '.join([t['text'] for t in extract_page_texts(slides[i])])
            cn_titles = re.findall(r'[\u4e00-\u9fa5]{2,10}', page_text)
            chapter_name = cn_titles[-1] if cn_titles else f'第{chapter_idx}章'
    
    # 补最后一个章节
    if chapter_start is not None:
        end_page = len(page_types)
        if 'end' in page_types:
            end_page = page_types.index('end')
        chapters.append({
            'key': f'chapter_{chapter_idx}',
            'start_page': chapter_start + 1,
            'end_page': end_page,
            'name': chapter_name
        })

    return chapters

def _hex_to_hsl(hex_color: str) -> Optional[tuple[float, float, float]]:
    """
    将 6 位 HEX 颜色转换为 HSL（H/S/L 均为 0-1 浮点）。
    解析失败返回 None。
    """
    if not hex_color or not isinstance(hex_color, str):
        return None
    hex_color = hex_color.strip().lstrip('#')
    if len(hex_color) != 6:
        return None
    try:
        r = int(hex_color[0:2], 16) / 255.0
        g = int(hex_color[2:4], 16) / 255.0
        b = int(hex_color[4:6], 16) / 255.0
    except ValueError:
        return None
    h, l, s = colorsys.rgb_to_hls(r, g, b)  # 注意 colorsys 返回顺序为 H, L, S
    return h, s, l


def _hue_to_color_family(h: float, s: float) -> str:
    """
    按色相 H（0-1）和饱和度 S（0-1）归类到色系：
      - 灰色系：S < 0.10
      - 蓝色系：H 在 200-240 度
      - 红色系：H 在 0-30 或 330-360 度
      - 绿色系：H 在 90-150 度
      - 其他单色（黄/橙/紫等）返回原始色相值字符串，由上层归并
    """
    if s < 0.10:
        return "灰色系"
    hue_deg = h * 360.0
    if 200 <= hue_deg <= 240:
        return "蓝色系"
    if 0 <= hue_deg <= 30 or 330 <= hue_deg <= 360:
        return "红色系"
    if 90 <= hue_deg <= 150:
        return "绿色系"
    return f"其他({int(hue_deg)}°)"


def detect_color_scheme(prs: Presentation) -> str:
    """
    从 PPT slide master 的主题色（<a:clrScheme>）中提取 accent1-6 的 RGB，
    按色相归类为色系字符串。

    规则：
      1. 取 accent1-6 的 srgbClr 或 sysClr/@lastClr 作为 RGB
      2. 计算 HSL，按色相归类到 蓝色系/红色系/绿色系/灰色系/黑白
      3. 若 accent1-6 色相分布跨多个不同色系（差异大），归为 "多彩"
      4. 全部 accent 饱和度极低（< 5%）或仅黑白，归为 "黑白"
      5. 兜底默认 "蓝色系"（商务模板主流）

    向后兼容：解析失败时返回 "蓝色系"。
    """
    default_scheme = "蓝色系"
    try:
        if not prs.slide_masters:
            return default_scheme
        master = prs.slide_masters[0]
        # 通过 master.part.rels 找到 theme 关系
        theme_part = None
        for rel in master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return default_scheme
        theme_el = etree.fromstring(theme_part.blob)
        scheme_list = theme_el.xpath('.//a:clrScheme', namespaces=_NS)
        if not scheme_list:
            return default_scheme
        scheme = scheme_list[0]

        # 提取 accent1-6 的 RGB
        accents: list[str] = []
        for i in range(1, 7):
            node = scheme.xpath(f'a:accent{i}', namespaces=_NS)
            if not node:
                continue
            srgb = node[0].xpath('a:srgbClr/@val', namespaces=_NS)
            if srgb:
                accents.append(srgb[0])
                continue
            sysc = node[0].xpath('a:sysClr/@lastClr', namespaces=_NS)
            if sysc:
                accents.append(sysc[0])

        if not accents:
            return default_scheme

        # 计算 HSL 并归类
        families: list[str] = []
        hues: list[float] = []
        all_low_sat = True
        for hex_c in accents:
            hsl = _hex_to_hsl(hex_c)
            if hsl is None:
                continue
            h, s, l = hsl
            if s >= 0.05:
                all_low_sat = False
            hues.append(h)
            families.append(_hue_to_color_family(h, s))

        # 全部极低饱和度 → 黑白
        if all_low_sat:
            return "黑白"

        # 色相差异大 → 多彩
        # 用色相极差（考虑环状）判定
        if len(hues) >= 3:
            hue_degs = [h * 360.0 for h in hues]
            # 计算每对色相的最小环距，用最大最小环距差近似
            sorted_hues = sorted(hue_degs)
            # 环状最大间距：max gap between adjacent (including wrap)
            gaps = []
            for i in range(len(sorted_hues)):
                nxt = sorted_hues[(i + 1) % len(sorted_hues)]
                gap = (nxt - sorted_hues[i]) % 360.0
                gaps.append(gap)
            max_gap = max(gaps)
            # 总圆周 - 最大空隙 ≈ 颜色覆盖范围
            coverage = 360.0 - max_gap
            if coverage > 180.0:
                # 同时多个不同色系成员
                unique_families = {f for f in families if not f.startswith("其他")}
                if len(unique_families) >= 3:
                    return "多彩"

        # 取 accent1 的归类作为主色系
        if families:
            main_family = families[0]
            if main_family.startswith("其他"):
                return default_scheme
            return main_family

        return default_scheme
    except Exception as e:
        logger.debug("色系识别失败: %s", e)
        return default_scheme


def generate_single_meta(pptx_path: Path, category: str) -> tuple[Optional[dict[str, Any]], Optional[str]]:
    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        return None, f"解析失败: {str(e)}"
    
    total_pages = len(prs.slides)
    page_types = []
    all_page_slots = {}
    page_meta = {}

    for idx, slide in enumerate(prs.slides):
        page_texts = extract_page_texts(slide)
        ptype = detect_page_type(page_texts, idx, total_pages)
        page_types.append(ptype)
        slots = extract_slots(page_texts, idx)

        # 收集页面非文本形状信息（chart/table/picture）
        shape_flags = detect_page_shapes(slide)

        # 追加图表/表格专用槽位（标识 chart/table 位置，供渲染器识别）
        if shape_flags["has_chart"]:
            slots.append({
                "slot": "chart_data",
                "shape_name": shape_flags.get("chart_shape_name", ""),
                "chart_type": shape_flags.get("chart_type"),
                "match_text": "",  # 图表形状无文本，不参与文本匹配
            })
        if shape_flags["has_table"]:
            slots.append({
                "slot": "table_data",
                "shape_name": shape_flags.get("table_shape_name", ""),
                "match_text": "",  # 表格形状无文本，不参与文本匹配
            })

        if slots:
            all_page_slots[str(idx + 1)] = slots

        page_key = str(idx + 1)
        page_entry: Optional[dict[str, Any]] = None
        if shape_flags["has_chart"] or shape_flags["has_table"] or shape_flags["has_picture"]:
            page_entry = {
                "has_chart": shape_flags["has_chart"],
                "has_table": shape_flags["has_table"],
                "has_picture": shape_flags["has_picture"],
                "chart_type": shape_flags["chart_type"],
                "picture_count": shape_flags["picture_count"],
                "chart_shape_name": shape_flags.get("chart_shape_name"),
                "table_shape_name": shape_flags.get("table_shape_name"),
            }

        # 元素分类告警（低置信度标注，来自 classify_page）
        try:
            classification = classify_page(slide)
            cls_warnings = classification.get("low_confidence_warnings", [])
            if cls_warnings:
                if page_entry is None:
                    page_entry = {}
                page_entry["classification_warnings"] = cls_warnings
        except Exception as e:
            logger.debug("页面分类失败 page %d: %s", idx + 1, e)

        if page_entry is not None:
            page_meta[page_key] = page_entry
    
    # 可删除版权页
    removable_pages = [i+1 for i, t in enumerate(page_types) if t == 'copyright']
    
    # 章节结构
    chapters = detect_chapters(prs.slides, page_types)

    # 补充固定页面（按页码升序合并，避免 cover 缺失时插入位置错乱）
    fixed_pages = []
    if 'cover' in page_types:
        fixed_pages.append({
            'key': 'cover',
            'page': page_types.index('cover') + 1,
            'name': '封面'
        })
    if 'catalog' in page_types:
        fixed_pages.append({
            'key': 'catalog',
            'page': page_types.index('catalog') + 1,
            'name': '目录'
        })
    if 'end' in page_types:
        fixed_pages.append({
            'key': 'end',
            'page': page_types.index('end') + 1,
            'name': '结束页'
        })
    chapters = fixed_pages + chapters
    chapters.sort(key=lambda c: c.get('start_page', c.get('page', 0)))
    
    # 生成模板ID（使用完整文件名确保唯一）
    template_id = f"{category}_{pptx_path.stem}".replace(' ', '_').replace('-', '_').lower()

    # 自动色系识别（基于 slide master 主题色）
    color_scheme = detect_color_scheme(prs)

    meta = {
        "template_id": template_id,
        "category": category,
        "style_tags": ["商务", "16:9"],
        "total_pages": total_pages,
        "removable_pages": removable_pages,
        "chapters": chapters,
        "page_slots": all_page_slots,
        "color_scheme": color_scheme
    }
    # 仅在有复合页面时才写入 page_meta，避免空字段污染 meta
    if page_meta:
        meta["page_meta"] = page_meta

    return meta, None

# ==================== 功能1：批量生成meta ====================
def cmd_generate(args):
    root = Path(args.dir)
    pptx_files = [p for p in root.rglob("*.pptx") if not p.name.startswith('~$')]
    logger.info("找到PPT文件: %d 个", len(pptx_files))
    
    success = 0
    failed = []
    
    for pptx_path in pptx_files:
        meta_path = pptx_path.with_suffix('.meta.json')
        if meta_path.exists() and not args.force:
            logger.info("跳过(已存在): %s", pptx_path.name)
            success += 1
            continue
            
        category = pptx_path.parent.name if pptx_path.parent.name != root.name else '通用'
        
        meta, err = generate_single_meta(pptx_path, category)
        if err:
            failed.append(f"{pptx_path.relative_to(root)}: {err}")
            logger.error("失败: %s - %s", pptx_path.name, err)
            continue
        
        with open(meta_path, 'w', encoding='utf-8') as f:
            json.dump(meta, f, ensure_ascii=False, indent=2)
        
        logger.info("生成成功: %s", pptx_path.name)
        success += 1
    
    logger.info("生成完成：成功 %d 个，失败 %d 个", success, len(failed))
    if failed:
        logger.warning("失败列表：")
        for f in failed:
            logger.warning("  - %s", f)

# ==================== P1-2.2.2 模板质量门禁：核心校验函数 ====================
def _find_pptx_for_meta(meta_path: Path) -> Optional[Path]:
    """根据 meta.json 文件路径推断对应的 pptx 文件路径

    支持两种命名约定：
      - xxx.pptx.meta.json → xxx.pptx
      - xxx.meta.json      → xxx.pptx
    """
    name = meta_path.name
    if name.endswith('.meta.json'):
        base = name[:-len('.meta.json')]
    elif name.endswith('.json'):
        base = name[:-len('.json')]
    else:
        base = name

    # 优先尝试 base 本身（适用于 xxx.pptx.meta.json）
    candidate = meta_path.parent / base
    if candidate.exists() and candidate.suffix.lower() == '.pptx':
        return candidate
    # 再尝试追加 .pptx 扩展名（适用于 xxx.meta.json）
    if not base.lower().endswith('.pptx'):
        candidate2 = meta_path.parent / (base + '.pptx')
        if candidate2.exists():
            return candidate2
    return None


def _check_meta_required(meta: dict[str, Any], rel_path: str = "") -> dict[str, Any]:
    """必填字段与基础识别质量校验（保留原有 4 项检查）

    包含：META_REQUIRED_FIELDS、章节数量、槽位数量、版权页识别
    """
    issues: list[str] = []
    prefix = f"{rel_path}: " if rel_path else ""

    for field in META_REQUIRED_FIELDS:
        if field not in meta:
            issues.append(f"{prefix}缺少必填字段 {field}")

    chapters = meta.get('chapters', [])
    if len(chapters) < 2:
        issues.append(f"{prefix}章节数量过少({len(chapters)}个)，识别可能不准确")

    slots_count = sum(len(v) for v in meta.get('page_slots', {}).values())
    if slots_count < 5:
        issues.append(f"{prefix}可替换槽位过少({slots_count}个)，识别可能不准确")

    if not meta.get('removable_pages'):
        end_page = next((c.get('page') for c in chapters if c.get('key') == 'end'), None)
        if end_page is not None and end_page != meta.get('total_pages'):
            issues.append(f"{prefix}未识别到版权页，需手动确认")

    return {"pass": len(issues) == 0, "issues": issues}


def _check_rendering_test(meta: dict[str, Any], pptx_path: Optional[Path]) -> dict[str, Any]:
    """渲染测试：用最小测试数据驱动 PptRenderer，验证模板可渲染且页数正确

    - 构造每个文本槽位填入"测试文本"
    - 调用 PptRenderer 渲染到临时目录（不污染 models/）
    - 校验：渲染不报错、输出文件存在、输出页数 = total_pages - len(removable_pages)
    - 失败时返回具体错误信息，不影响其他校验项继续执行
    """
    result: dict[str, Any] = {"pass": False, "issues": [], "output_pages": 0}

    if pptx_path is None or not pptx_path.exists():
        result["issues"].append("模板 pptx 文件不存在，跳过渲染测试")
        return result

    # 延迟导入，避免主流程对 ppt_renderer 的硬依赖
    try:
        from ppt_renderer import PptRenderer
    except Exception as e:
        result["issues"].append(f"PptRenderer 模块加载失败: {e}")
        return result

    # 构造最小测试数据：每个文本槽位填入"测试文本"，跳过 chart_data/table_data
    slot_data: dict[str, dict[str, str]] = {}
    for page_str, slots in meta.get('page_slots', {}).items():
        page_input: dict[str, str] = {}
        for slot in slots:
            slot_name = slot.get('slot', '')
            if slot_name in ('chart_data', 'table_data'):
                continue
            page_input[slot_name] = "测试文本"
        if page_input:
            slot_data[page_str] = page_input

    try:
        with tempfile.TemporaryDirectory(prefix="ppt_render_test_") as tmpdir:
            # PptRenderer 需要从文件加载 meta，写入临时 meta 文件
            tmp_meta_path = Path(tmpdir) / "tmp.meta.json"
            with open(tmp_meta_path, 'w', encoding='utf-8') as f:
                json.dump(meta, f, ensure_ascii=False)
            output_path = str(Path(tmpdir) / "test_output.pptx")

            renderer = PptRenderer(str(pptx_path), str(tmp_meta_path))
            renderer.render(
                slot_data,
                output_path,
                remove_copyright=True,
                auto_fit=True,
            )

            # 校验输出文件存在
            if not Path(output_path).exists():
                result["issues"].append("渲染完成但输出文件不存在")
                return result

            # 校验输出页数
            try:
                out_prs = Presentation(output_path)
                actual_pages = len(out_prs.slides)
            except Exception as e:
                result["issues"].append(f"读取渲染输出页数失败: {e}")
                return result

            expected_pages = meta.get('total_pages', 0) - len(meta.get('removable_pages', []))
            result["output_pages"] = actual_pages
            if actual_pages != expected_pages:
                result["issues"].append(
                    f"输出页数 {actual_pages} 与预期 {expected_pages} 不一致（total_pages - removable_pages）"
                )
                return result

            result["pass"] = True
            return result
    except Exception as e:
        # 渲染失败不影响其他校验项继续执行（异常被捕获后返回）
        result["issues"].append(f"渲染测试失败: {e}")
        return result


def _check_style(meta: dict[str, Any], pptx_path: Optional[Path]) -> dict[str, Any]:
    """样式检查：字体一致性、配色协调、字号合理性（仅产生 warning，不阻断）

    - 同种角色用同种字体（标题/正文）
    - 主色调不超过 3 种（基于文本颜色 RGB 去重）
    - 标题字号 ≥ 24pt
    - 正文字号 ≥ 14pt
    """
    warnings: list[str] = []
    if pptx_path is None or not pptx_path.exists():
        warnings.append("模板 pptx 文件不存在，跳过样式检查")
        return {"pass": True, "warnings": warnings}

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        warnings.append(f"加载 pptx 失败: {e}")
        return {"pass": True, "warnings": warnings}

    title_fonts: set[str] = set()
    body_fonts: set[str] = set()
    title_sizes: list[float] = []
    body_sizes: list[float] = []
    text_colors: set[str] = set()

    def _scan_shape(shape, role: str = 'body') -> None:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                # 段落级字体作为 run 字体的回退
                try:
                    pf_name = para.font.name
                except Exception:
                    pf_name = None
                for run in para.runs:
                    try:
                        font = run.font
                    except Exception:
                        continue
                    font_name = font.name or pf_name
                    if font_name:
                        if role == 'title':
                            title_fonts.add(font_name)
                        else:
                            body_fonts.add(font_name)
                    try:
                        if font.size is not None:
                            size_pt = font.size.pt
                            if role == 'title':
                                title_sizes.append(size_pt)
                            else:
                                body_sizes.append(size_pt)
                    except Exception:
                        pass
                    # 文本颜色（仅收集显式 RGB）
                    try:
                        if font.color is not None and font.color.type is not None:
                            rgb = font.color.rgb
                            if rgb is not None:
                                text_colors.add(str(rgb))
                    except Exception:
                        pass
        # GROUP 递归
        if shape.shape_type == 6:
            try:
                for sub in shape.shapes:
                    _scan_shape(sub, role)
            except Exception:
                pass

    for slide in prs.slides:
        for shape in slide.shapes:
            role = 'body'
            # 占位符 idx=0 通常是标题
            try:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    role = 'title'
            except Exception:
                pass
            shape_name = (shape.name or '')
            if 'title' in shape_name.lower() or '标题' in shape_name:
                role = 'title'
            try:
                _scan_shape(shape, role)
            except Exception:
                continue

    # 字体一致性（同种角色用同种字体）
    if len(title_fonts) > 1:
        warnings.append(f"标题字体不统一（{len(title_fonts)}种）: {sorted(title_fonts)}")
    if len(body_fonts) > 2:
        warnings.append(f"正文字体过多（{len(body_fonts)}种）: {sorted(body_fonts)}")

    # 配色协调（主色调不超过 3 种）
    if len(text_colors) > 3:
        warnings.append(f"主色调过多（{len(text_colors)}种）: {sorted(text_colors)}（建议不超过 3 种）")

    # 字号合理性
    if title_sizes and min(title_sizes) < 24:
        warnings.append(f"标题字号过小: 最小 {min(title_sizes)}pt（建议 ≥ 24pt）")
    if body_sizes and min(body_sizes) < 14:
        warnings.append(f"正文字号过小: 最小 {min(body_sizes)}pt（建议 ≥ 14pt）")

    return {"pass": True, "warnings": warnings}


def _check_meta_completeness(meta: dict[str, Any], rel_path: str = "",
                             pptx_path: Optional[Path] = None) -> dict[str, Any]:
    """元数据完整性校验

    - 每页必须含 page_type 或 pattern 字段（在 page_meta 或 page_slots 中查找）
    - 每页 slots 数组中每个槽位必须含 slot/match_text 字段
    - chapters 数组中每个章节必须含 key/name/page_range 字段
    - template_id 必须以「文件名（去扩展名）」为基础生成
    """
    issues: list[str] = []
    prefix = f"{rel_path}: " if rel_path else ""

    total_pages = meta.get('total_pages', 0)
    page_meta = meta.get('page_meta', {})
    page_slots = meta.get('page_slots', {})

    # 每页必须含 page_type 或 pattern 字段
    for page_num in range(1, total_pages + 1):
        page_key = str(page_num)
        pm = page_meta.get(page_key, {}) if isinstance(page_meta, dict) else {}
        has_page_type = 'page_type' in pm
        has_pattern = 'pattern' in pm
        # 也接受 page_slots 中含 page_type/pattern 的槽位
        if not has_page_type and not has_pattern:
            slots = page_slots.get(page_key, []) if isinstance(page_slots, dict) else []
            for s in slots:
                if isinstance(s, dict) and ('page_type' in s or 'pattern' in s):
                    has_page_type = True
                    break
        if not has_page_type and not has_pattern:
            issues.append(f"{prefix}第{page_num}页缺少 page_type 或 pattern 字段")

    # 每页 slots 数组中每个槽位必须含 slot/match_text 字段
    for page_key, slots in page_slots.items():
        if not isinstance(slots, list):
            continue
        for idx, slot in enumerate(slots):
            if not isinstance(slot, dict):
                continue
            if 'slot' not in slot:
                issues.append(f"{prefix}第{page_key}页槽位 #{idx} 缺少 slot 字段")
            if 'match_text' not in slot:
                issues.append(f"{prefix}第{page_key}页槽位 #{idx} 缺少 match_text 字段")

    # chapters 数组中每个章节必须含 key/name/page_range 字段
    for idx, ch in enumerate(meta.get('chapters', [])):
        if not isinstance(ch, dict):
            continue
        if 'key' not in ch:
            issues.append(f"{prefix}章节 #{idx} 缺少 key 字段")
        if 'name' not in ch:
            issues.append(f"{prefix}章节 #{idx} 缺少 name 字段")
        if 'page_range' not in ch:
            issues.append(f"{prefix}章节 #{idx} 缺少 page_range 字段")

    # template_id 必须以「文件名（去扩展名）」为基础生成
    template_id = meta.get('template_id', '')
    if pptx_path is not None and template_id:
        stem = pptx_path.stem
        # generate_single_meta 中：template_id = f"{category}_{pptx_path.stem}".replace(' ', '_').replace('-', '_').lower()
        normalized_stem = stem.replace(' ', '_').replace('-', '_').lower()
        normalized_tid = template_id.replace(' ', '_').replace('-', '_').lower()
        if normalized_stem and normalized_stem not in normalized_tid:
            issues.append(f"{prefix}template_id '{template_id}' 未基于文件名 '{stem}' 生成")

    return {"pass": len(issues) == 0, "issues": issues}


def run_quality_check(meta: dict[str, Any], pptx_path: Any = None,
                      rel_path: str = "") -> dict[str, Any]:
    """对单个模板 meta 执行完整质量校验，返回结构化结果

    :param meta: 模板元数据 dict
    :param pptx_path: 模板 pptx 文件路径（Path/str），用于渲染测试和样式检查
    :param rel_path: 用于日志展示的相对路径前缀
    :return: 结构化校验结果，包含 4 类校验明细

    返回结构：
      {
        "template_id": "...",
        "check_pass": bool,  # meta_required + rendering_test + meta_completeness 全部通过
        "checks": {
          "meta_required": {"pass": bool, "issues": [...]},
          "rendering_test": {"pass": bool, "issues": [...], "output_pages": int},
          "style_check": {"pass": bool, "warnings": [...]},
          "meta_completeness": {"pass": bool, "issues": [...]}
        },
        "total_issues": int,
        "total_warnings": int
      }
    """
    # 接受 str 或 Path，统一为 Path 或 None
    if pptx_path is not None and not isinstance(pptx_path, Path):
        pptx_path = Path(pptx_path)

    # 1. 必填字段与基础识别质量校验（保留原有 4 项检查）
    meta_required = _check_meta_required(meta, rel_path)

    # 2. 渲染测试（失败不影响其他校验项继续执行）
    try:
        rendering_test = _check_rendering_test(meta, pptx_path)
    except Exception as e:
        rendering_test = {"pass": False, "issues": [f"渲染测试异常: {e}"], "output_pages": 0}

    # 3. 样式检查（仅产生 warning，不阻断）
    try:
        style_check = _check_style(meta, pptx_path)
    except Exception as e:
        style_check = {"pass": True, "warnings": [f"样式检查异常: {e}"]}

    # 4. 元数据完整性校验
    try:
        meta_completeness = _check_meta_completeness(meta, rel_path, pptx_path)
    except Exception as e:
        meta_completeness = {"pass": False, "issues": [f"元数据完整性检查异常: {e}"]}

    # check_pass: meta_required + rendering_test + meta_completeness 全部通过
    # style_check 仅产生 warning，不影响 check_pass
    check_pass = (
        meta_required.get("pass", False)
        and rendering_test.get("pass", False)
        and meta_completeness.get("pass", False)
    )

    total_issues = (
        len(meta_required.get("issues", []))
        + len(rendering_test.get("issues", []))
        + len(meta_completeness.get("issues", []))
    )
    total_warnings = len(style_check.get("warnings", []))

    return {
        "template_id": meta.get("template_id", ""),
        "check_pass": check_pass,
        "checks": {
            "meta_required": meta_required,
            "rendering_test": rendering_test,
            "style_check": style_check,
            "meta_completeness": meta_completeness,
        },
        "total_issues": total_issues,
        "total_warnings": total_warnings,
    }


# ==================== 功能2：批量校验meta质量 ====================
def cmd_check(args):
    root = Path(args.dir)
    meta_files = list(root.rglob("*.meta.json"))
    logger.info("找到meta文件: %d 个", len(meta_files))

    issues = []  # 兼容旧逻辑：汇总所有 issues 供日志输出
    category_count = defaultdict(int)
    summary_lines: list[str] = []  # 每个模板的人类可读摘要行

    for meta_path in meta_files:
        try:
            with open(meta_path, 'r', encoding='utf-8') as f:
                meta = json.load(f)
        except Exception:
            issues.append(f"{meta_path.relative_to(root)}: 文件损坏无法读取")
            summary_lines.append(f"❌ {meta_path.name} - 文件损坏无法读取")
            continue

        rel_path = str(meta_path.relative_to(root))
        category = meta.get('category', '未知')
        category_count[category] += 1

        # 推断对应的 pptx 文件路径（用于渲染测试与样式检查）
        pptx_path = _find_pptx_for_meta(meta_path)

        # 执行完整质量校验（4 类校验项）
        result = run_quality_check(meta, pptx_path, rel_path)

        # 汇总 issues 供旧日志输出
        for chk_name, chk_result in result["checks"].items():
            for issue in chk_result.get("issues", []):
                issues.append(issue)

        # 人类可读：✅/❌ 标记每项校验结果
        checks = result["checks"]
        marker = "✅" if result["check_pass"] else "❌"
        line = f"{marker} {meta_path.name} (template_id={result['template_id']})"
        line += f"\n    meta_required: {'✅' if checks['meta_required']['pass'] else '❌'}"
        line += f"\n    rendering_test: {'✅' if checks['rendering_test']['pass'] else '❌'} (output_pages={checks['rendering_test'].get('output_pages', 0)})"
        line += f"\n    style_check: {'✅' if checks['style_check']['pass'] else '⚠️'} (warnings={len(checks['style_check'].get('warnings', []))})"
        line += f"\n    meta_completeness: {'✅' if checks['meta_completeness']['pass'] else '❌'}"
        line += f"\n    total_issues={result['total_issues']}, total_warnings={result['total_warnings']}"
        if not result["check_pass"]:
            # 输出具体问题，便于排查
            for chk_name, chk_result in checks.items():
                for issue in chk_result.get("issues", []):
                    line += f"\n      - [{chk_name}] {issue}"
        summary_lines.append(line)

    logger.info("分类统计：")
    for cat, cnt in category_count.items():
        logger.info("  %s: %d 套", cat, cnt)

    # 输出每个模板的校验摘要（人类可读 ✅/❌）
    for line in summary_lines:
        logger.info(line)

    logger.info("校验完成：共检查 %d 个meta，发现问题 %d 个", len(meta_files), len(issues))
    if issues:
        logger.warning("问题清单：")
        for issue in issues:
            logger.warning("  %s", issue)
    else:
        logger.info("全部meta校验通过！")

# ==================== 功能3：生成总索引文件 ====================
def cmd_index(args):
    root = Path(args.dir)
    meta_files = list(root.rglob("*.meta.json"))
    logger.info("找到meta文件: %d 个", len(meta_files))
    
    index = {
        "total": len(meta_files),
        "categories": {},
        "templates": []
    }
    
    category_map = defaultdict(list)
    
    for meta_path in meta_files:
        try:
            with open(meta_path, 'r', encoding='utf-8') as f:
                meta = json.load(f)
        except:
            continue
        
        # 精简索引信息
        template_info = {
            "template_id": meta['template_id'],
            "category": meta['category'],
            "name": meta_path.stem.replace('.meta', ''),
            "path": str(meta_path.relative_to(root)),
            "style_tags": meta.get('style_tags', []),
            "total_pages": meta['total_pages'],
            "chapter_count": len(meta.get('chapters', []))
        }
        
        index['templates'].append(template_info)
        category_map[meta['category']].append(template_info)
    
    index['categories'] = {k: len(v) for k, v in category_map.items()}
    
    output_path = root / 'templates_index.json'
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(index, f, ensure_ascii=False, indent=2)
    
    logger.info("总索引生成完成，共 %d 个模板", len(index['templates']))
    logger.info("索引文件路径: %s", output_path)
    logger.info("分类数量统计：")
    for cat, cnt in index['categories'].items():
        logger.info("  %s: %d 套", cat, cnt)

# ==================== 主入口 ====================
def _find_meta_by_template_id(template_id: str, models_root: Path = None) -> Optional[Path]:
    """根据 template_id 在 models 目录下查找对应的 meta.json 文件"""
    root = models_root or MODELS_ROOT
    for meta_path in root.rglob("*.meta.json"):
        try:
            with open(meta_path, 'r', encoding='utf-8') as f:
                meta = json.load(f)
            if meta.get('template_id') == template_id:
                return meta_path
        except Exception:
            continue
    return None


def cmd_info(args):
    """查看单个模板的元数据详情"""
    meta_path = _find_meta_by_template_id(args.template_id)
    if not meta_path:
        logger.error("未找到 template_id: %s", args.template_id)
        return

    with open(meta_path, 'r', encoding='utf-8') as f:
        meta = json.load(f)

    info = {
        "template_id": meta.get("template_id"),
        "category": meta.get("category"),
        "total_pages": meta.get("total_pages"),
        "removable_pages": meta.get("removable_pages", []),
        "chapters": meta.get("chapters", []),
        "page_slot_count": len(meta.get("page_slots", {})),
        "total_slots": sum(len(v) for v in meta.get("page_slots", {}).values()),
        "page_meta": meta.get("page_meta", {}),
        "meta_path": str(meta_path),
    }
    print(json.dumps(info, ensure_ascii=False, indent=2))


def cmd_check_one(args):
    """单个模板质量校验：输出结构化 JSON（包含 4 类校验明细 + 向后兼容字段）"""
    meta_path = _find_meta_by_template_id(args.template_id)
    if not meta_path:
        logger.error("未找到 template_id: %s", args.template_id)
        return

    with open(meta_path, 'r', encoding='utf-8') as f:
        meta = json.load(f)

    # 推断对应的 pptx 文件路径（用于渲染测试与样式检查）
    pptx_path = _find_pptx_for_meta(meta_path)

    # 执行完整质量校验（4 类校验项）
    result = run_quality_check(meta, pptx_path)

    # 向后兼容字段：保留原有顶层字段（total_pages/chapter_count/slots_count/issues/is_valid）
    # 聚合所有 issues（兼容旧字段）
    legacy_issues: list[str] = []
    for chk_result in result["checks"].values():
        legacy_issues.extend(chk_result.get("issues", []))

    result["total_pages"] = meta.get("total_pages")
    result["chapter_count"] = len(meta.get('chapters', []))
    result["slots_count"] = sum(len(v) for v in meta.get('page_slots', {}).values())
    result["issues"] = legacy_issues
    result["is_valid"] = result["check_pass"]

    print(json.dumps(result, ensure_ascii=False, indent=2))


def main():
    parser = argparse.ArgumentParser(description='PPT模板元数据批量处理工具')
    subparsers = parser.add_subparsers(dest='command', required=True, help='支持的命令')

    # generate 子命令
    gen_parser = subparsers.add_parser('generate', help='批量生成meta.json元数据')
    gen_parser.add_argument('--dir', required=True, help='模板根目录路径')
    gen_parser.add_argument('--force', action='store_true', help='强制覆盖已存在的meta文件')

    # check 子命令（批量）
    check_parser = subparsers.add_parser('check', help='批量校验meta质量')
    check_parser.add_argument('--dir', help='模板根目录路径（批量校验）')
    check_parser.add_argument('--template-id', help='单个模板 ID（单模板校验）')

    # info 子命令（单模板详情）
    info_parser = subparsers.add_parser('info', help='查看单个模板元数据详情')
    info_parser.add_argument('--template-id', required=True, help='模板 ID')

    # index 子命令
    index_parser = subparsers.add_parser('index', help='生成模板总索引文件')
    index_parser.add_argument('--dir', required=True, help='模板根目录路径')

    args = parser.parse_args()

    if args.command == 'generate':
        cmd_generate(args)
    elif args.command == 'check':
        if args.template_id:
            cmd_check_one(args)
        elif args.dir:
            cmd_check(args)
        else:
            parser.error("check 命令需要 --dir 或 --template-id 参数")
    elif args.command == 'info':
        cmd_info(args)
    elif args.command == 'index':
        cmd_index(args)

if __name__ == '__main__':
    main()