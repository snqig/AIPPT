"""T615 排版自检工具单元测试

覆盖 lint_pptx 主入口与 5 类检测项：
    1. 元素溢出（左/上/右/下、安全边距）
    2. 元素重叠（文本元素间）
    3. 间距异常（纵向相邻元素）
    4. 字号下限（< 10pt）
    5. 空文本框（大尺寸无内容）

并验证：
    - IssueReport 数据结构 / 序列化
    - 装饰元素（bg/card/divider/line/badge）跳过逻辑
    - format_report 人类可读输出
    - CLI 参数解析路径
"""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from aippt.layout.layout_lint import (
    LayoutIssue, IssueReport,
    lint_pptx, format_report,
    EMU_PER_INCH, MIN_FONT_PT, MIN_GAP_INCH,
)


# ==================== Fixture：构造 PPTX ====================

def _make_pptx(tmp_path, name="test.pptx"):
    """创建空白 16:9 PPTX"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, tmp_path / name


def _add_textbox(slide, left_in, top_in, w_in, h_in, name, text="", font_pt=None):
    """添加文本框，可选设置字号"""
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in),
    )
    tb.name = name
    if text:
        if font_pt is not None:
            run = tb.text_frame.paragraphs[0].add_run()
            run.text = text
            run.font.size = Pt(font_pt)
        else:
            tb.text_frame.text = text
    return tb


# ==================== 1. 数据结构测试 ====================

class TestIssueReport:
    """IssueReport / LayoutIssue 数据结构"""

    def test_issue_report_counts(self):
        report = IssueReport(file_path="x.pptx", total_pages=3)
        report.issues.append(LayoutIssue(1, "error", "overflow", "e1"))
        report.issues.append(LayoutIssue(1, "warning", "overlap", "w1"))
        report.issues.append(LayoutIssue(2, "info", "spacing", "i1"))
        report.issues.append(LayoutIssue(3, "error", "fontsize", "e2"))

        assert report.error_count == 2
        assert report.warning_count == 1
        assert not report.passed  # 有 error

    def test_issue_report_passed(self):
        report = IssueReport(file_path="x.pptx", total_pages=1)
        assert report.passed
        assert report.error_count == 0

    def test_to_dict_serialization(self):
        report = IssueReport(file_path="x.pptx", total_pages=2)
        report.issues.append(LayoutIssue(
            page=1, severity="warning", category="overlap",
            message="重叠", shape_name="s1",
            detail={"overlap_w_inch": 1.2},
        ))
        d = report.to_dict()
        assert d["file_path"] == "x.pptx"
        assert d["total_pages"] == 2
        assert d["passed"] is True
        assert d["warning_count"] == 1
        assert len(d["issues"]) == 1
        assert d["issues"][0]["shape_name"] == "s1"
        assert d["issues"][0]["detail"]["overlap_w_inch"] == 1.2

    def test_layout_issue_defaults(self):
        issue = LayoutIssue(page=1, severity="error", category="overflow", message="溢出")
        assert issue.shape_name is None
        assert issue.detail == {}


# ==================== 2. 元素溢出检测 ====================

class TestOverflowDetection:

    def test_left_overflow_error(self, tmp_path):
        """元素 left < 0 → error"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, -0.5, 1, 3, 0.5, "role:overflow_left", "text")
        prs.save(str(path))

        report = lint_pptx(str(path), check_overlap=False, check_spacing=False,
                           check_fontsize=False, check_empty=False)
        errors = [i for i in report.issues if i.severity == "error"]
        assert len(errors) == 1
        assert errors[0].category == "overflow"
        assert "左/上边界" in errors[0].message

    def test_right_overflow_error(self, tmp_path):
        """元素 right > 画布宽 → error"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 画布宽 13.333，元素 left=12, width=3 → right=15 超出
        _add_textbox(slide, 12, 1, 3, 0.5, "role:overflow_right", "text")
        prs.save(str(path))

        report = lint_pptx(str(path), check_overlap=False, check_spacing=False,
                           check_fontsize=False, check_empty=False)
        errors = [i for i in report.issues if i.severity == "error"]
        assert len(errors) == 1
        assert "右/下边界" in errors[0].message

    def test_safe_margin_warning(self, tmp_path):
        """文本元素进入安全边距 → warning"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 安全边距 0.5，left=0.3 进入安全区
        _add_textbox(slide, 0.3, 0.3, 3, 0.5, "role:in_safe", "text")
        prs.save(str(path))

        report = lint_pptx(str(path), check_overlap=False, check_spacing=False,
                           check_fontsize=False, check_empty=False)
        warns = [i for i in report.issues if i.severity == "warning" and i.category == "overflow"]
        assert len(warns) == 1

    def test_bg_shape_skipped(self, tmp_path):
        """bg 装饰元素不参与溢出检测"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 构造 bg 元素 left=0, width=20（明显溢出）
        _add_textbox(slide, 0, 0, 20, 7.5, "role:cover_bg", "background")
        prs.save(str(path))

        report = lint_pptx(str(path), check_overlap=False, check_spacing=False,
                           check_fontsize=False, check_empty=False)
        # 不应有 overflow error
        assert all(i.category != "overflow" for i in report.issues if i.severity == "error")


# ==================== 3. 元素重叠检测 ====================

class TestOverlapDetection:

    def test_text_overlap_warning(self, tmp_path):
        """两个文本元素重叠 → warning"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, 1, 2, 4, 0.8, "role:text_a", "A")
        _add_textbox(slide, 2, 2.2, 4, 0.8, "role:text_b", "B")
        prs.save(str(path))

        report = lint_pptx(str(path), check_overflow=False, check_spacing=False,
                           check_fontsize=False, check_empty=False)
        overlaps = [i for i in report.issues if i.category == "overlap"]
        assert len(overlaps) == 1
        assert "text_a" in overlaps[0].shape_name
        assert "text_b" in overlaps[0].shape_name

    def test_no_overlap_when_separated(self, tmp_path):
        """两个元素分离时不报告重叠"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, 1, 2, 3, 0.5, "role:text_a", "A")
        _add_textbox(slide, 6, 2, 3, 0.5, "role:text_b", "B")
        prs.save(str(path))

        report = lint_pptx(str(path), check_overflow=False, check_spacing=False,
                           check_fontsize=False, check_empty=False)
        assert all(i.category != "overlap" for i in report.issues)


# ==================== 4. 间距异常检测 ====================

class TestSpacingDetection:

    def test_small_gap_info(self, tmp_path):
        """纵向相邻元素间距过小 → info"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 两元素纵向紧贴，间距 0.01 inch < MIN_GAP_INCH(0.05)
        _add_textbox(slide, 1, 2, 3, 0.5, "role:t1", "A")
        _add_textbox(slide, 1, 2.51, 3, 0.5, "role:t2", "B")  # gap=0.01
        prs.save(str(path))

        report = lint_pptx(str(path), check_overflow=False, check_overlap=False,
                           check_fontsize=False, check_empty=False)
        spacings = [i for i in report.issues if i.category == "spacing"]
        assert len(spacings) == 1
        assert spacings[0].severity == "info"

    def test_normal_gap_no_issue(self, tmp_path):
        """正常间距无报告"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, 1, 2, 3, 0.5, "role:t1", "A")
        _add_textbox(slide, 1, 3, 3, 0.5, "role:t2", "B")  # gap=0.5
        prs.save(str(path))

        report = lint_pptx(str(path), check_overflow=False, check_overlap=False,
                           check_fontsize=False, check_empty=False)
        assert all(i.category != "spacing" for i in report.issues)


# ==================== 5. 字号下限检测 ====================

class TestFontSizeDetection:

    def test_small_font_warning(self, tmp_path):
        """字号 < 10pt → warning"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, 1, 1, 4, 0.5, "role:tiny", "tiny", font_pt=8)
        prs.save(str(path))

        report = lint_pptx(str(path), check_overflow=False, check_overlap=False,
                           check_spacing=False, check_empty=False)
        fs = [i for i in report.issues if i.category == "fontsize"]
        assert len(fs) == 1
        assert fs[0].detail["font_size_pt"] == 8.0
        assert fs[0].detail["min_font_pt"] == MIN_FONT_PT

    def test_normal_font_no_issue(self, tmp_path):
        """正常字号无报告"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, 1, 1, 4, 0.5, "role:normal", "normal", font_pt=16)
        prs.save(str(path))

        report = lint_pptx(str(path), check_overflow=False, check_overlap=False,
                           check_spacing=False, check_empty=False)
        assert all(i.category != "fontsize" for i in report.issues)


# ==================== 6. 空文本框检测 ====================

class TestEmptyTextboxDetection:

    def test_large_empty_textbox_info(self, tmp_path):
        """大尺寸空文本框 → info"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, 1, 1, 5, 1.5, "role:empty", "")  # 不填文本
        prs.save(str(path))

        report = lint_pptx(str(path), check_overflow=False, check_overlap=False,
                           check_spacing=False, check_fontsize=False)
        empties = [i for i in report.issues if i.category == "empty"]
        assert len(empties) == 1

    def test_card_shape_skipped(self, tmp_path):
        """card 装饰元素不参与空文本框检测"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, 1, 1, 5, 1.5, "role:card_1", "")  # 大尺寸空，但 name 含 card
        prs.save(str(path))

        report = lint_pptx(str(path), check_overflow=False, check_overlap=False,
                           check_spacing=False, check_fontsize=False)
        assert all(i.category != "empty" for i in report.issues)

    def test_small_empty_skipped(self, tmp_path):
        """小尺寸空文本框不报告（可能为装饰）"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # width=1.5 < 2 inch，不报告
        _add_textbox(slide, 1, 1, 1.5, 0.3, "role:small", "")
        prs.save(str(path))

        report = lint_pptx(str(path), check_overflow=False, check_overlap=False,
                           check_spacing=False, check_fontsize=False)
        assert all(i.category != "empty" for i in report.issues)


# ==================== 7. 主入口 & 报告输出 ====================

class TestLintPptxMain:

    def test_file_not_found(self):
        with pytest.raises(FileNotFoundError):
            lint_pptx("nonexistent.pptx")

    def test_disable_checks(self, tmp_path):
        """禁用全部检测项 → 空报告"""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, -1, 1, 3, 0.5, "role:overflow", "x", font_pt=8)
        prs.save(str(path))

        report = lint_pptx(
            str(path),
            check_overflow=False, check_overlap=False,
            check_spacing=False, check_fontsize=False, check_empty=False,
        )
        assert len(report.issues) == 0
        assert report.passed

    def test_format_report_basic(self):
        report = IssueReport(file_path="x.pptx", total_pages=2)
        report.issues.append(LayoutIssue(
            page=1, severity="error", category="overflow", message="测试错误",
        ))
        report.issues.append(LayoutIssue(
            page=2, severity="info", category="spacing", message="测试 info",
        ))
        text = format_report(report, verbose=False)
        assert "x.pptx" in text
        assert "测试错误" in text
        assert "测试 info" not in text  # 非 verbose 不输出 info

    def test_format_report_verbose(self):
        report = IssueReport(file_path="x.pptx", total_pages=1)
        report.issues.append(LayoutIssue(
            page=1, severity="info", category="spacing", message="测试 info",
        ))
        text = format_report(report, verbose=True)
        assert "测试 info" in text

    def test_format_report_with_detail(self):
        report = IssueReport(file_path="x.pptx", total_pages=1)
        report.issues.append(LayoutIssue(
            page=1, severity="warning", category="overlap", message="重叠",
            shape_name="s1", detail={"overlap_w_inch": 1.5},
        ))
        text = format_report(report)
        assert "s1" in text
        assert "overlap_w_inch=1.5" in text


# ==================== 8. 集成：完整检测流程 ====================

class TestIntegrationFullLint:
    """端到端：一份包含多类问题的 PPTX，验证报告完整性"""

    def test_mixed_issues_pptx(self, tmp_path):
        """一份包含 5 类问题的 PPTX，验证所有检测项触发"""
        prs, path = _make_pptx(tmp_path, name="mixed.pptx")
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 1. 溢出
        _add_textbox(slide, -0.3, 0.5, 3, 0.5, "role:overflow_left", "溢出")
        # 2. 重叠
        _add_textbox(slide, 1, 2, 4, 0.8, "role:text_a", "A")
        _add_textbox(slide, 2, 2.2, 4, 0.8, "role:text_b", "B")
        # 3. 字号过小
        _add_textbox(slide, 1, 4, 4, 0.5, "role:tiny", "小", font_pt=8)
        # 4. 空文本框
        _add_textbox(slide, 1, 5, 5, 1.5, "role:empty", "")

        prs.save(str(path))

        report = lint_pptx(str(path))
        categories = {i.category for i in report.issues}
        # 应至少触发 overflow/overlap/fontsize/empty
        assert "overflow" in categories
        assert "overlap" in categories
        assert "fontsize" in categories
        assert "empty" in categories
        assert not report.passed  # 有 error 级别
        assert report.error_count >= 1


# ==================== 9. CLI 接入验证 ====================

class TestCliLintLayout:
    """验证 aippt_outline.py lint-layout 子命令接入"""

    def test_cli_help_contains_lint_layout(self):
        """lint-layout 子命令在 --help 中可见"""
        import subprocess
        result = subprocess.run(
            ["python", "aippt_outline.py", "--help"],
            capture_output=True, text=True, cwd=str(Path(__file__).parent.parent),
        )
        assert result.returncode == 0
        assert "lint-layout" in result.stdout

    def test_cli_lint_layout_pass(self, tmp_path):
        """对通过自检的 PPTX，CLI 返回 0"""
        prs, path = _make_pptx(tmp_path, name="ok.pptx")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, 1, 1, 4, 0.5, "role:normal", "正常文本", font_pt=16)
        prs.save(str(path))

        import subprocess
        result = subprocess.run(
            ["python", "aippt_outline.py", "lint-layout", "--pptx", str(path)],
            capture_output=True, text=True, cwd=str(Path(__file__).parent.parent),
        )
        assert result.returncode == 0
        assert "通过" in result.stdout

    def test_cli_lint_layout_fail(self, tmp_path):
        """对有 error 的 PPTX，CLI 返回 1"""
        prs, path = _make_pptx(tmp_path, name="bad.pptx")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, -0.5, 1, 3, 0.5, "role:overflow", "溢出")
        prs.save(str(path))

        import subprocess
        result = subprocess.run(
            ["python", "aippt_outline.py", "lint-layout", "--pptx", str(path)],
            capture_output=True, text=True, cwd=str(Path(__file__).parent.parent),
        )
        assert result.returncode == 1
        assert "未通过" in result.stdout

    def test_cli_output_json(self, tmp_path):
        """--output 保存 JSON 报告"""
        prs, path = _make_pptx(tmp_path, name="ok.pptx")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, 1, 1, 4, 0.5, "role:normal", "文本", font_pt=14)
        prs.save(str(path))

        import json
        import subprocess
        out_json = tmp_path / "report.json"
        result = subprocess.run(
            ["python", "aippt_outline.py", "lint-layout",
             "--pptx", str(path), "--output", str(out_json)],
            capture_output=True, text=True, cwd=str(Path(__file__).parent.parent),
        )
        assert result.returncode == 0
        assert out_json.exists()
        data = json.loads(out_json.read_text(encoding="utf-8"))
        assert "file_path" in data
        assert "issues" in data
        assert "passed" in data
