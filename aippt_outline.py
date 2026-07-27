"""
AIPPT 大纲转换工具
功能：将四步工作流产出的 outline.json 转换为 business_data.json
依赖：ppt_scene_adapter.py 的 SCENE_SCHEMAS

大纲 JSON 格式（Step 2 产出）：
{
  "scene": "工作汇报",
  "cover": {"title": "Q3 项目复盘报告", "reporter": "张三", "period": "2024 Q3"},
  "sections": [
    {"key": "progress", "name": "工作进展", "items": [
      {"title": "Q3 里程碑", "desc": "完成核心功能开发与公测"}
    ]},
    {"key": "results", "name": "阶段成果", "items": [
      {"title": "用户增长", "desc": "DAU 同比增长 35%"}
    ]}
  ],
  "end": {"thanks": "感谢聆听"}
}

business_data JSON 格式（SceneAdapter 输入）：
{
  "cover": {"title": "...", "reporter": "...", "period": "..."},
  "sections": {"<key>": [{"title": "...", "desc": "..."}]},
  "end": {"thanks": "..."}
}
"""
import json
import argparse
from pathlib import Path
from typing import Any, Optional

from ppt_scene_adapter import SCENE_SCHEMAS, SceneAdapter
from aippt.config import SCENE_KEYWORDS
from aippt.logger import logger


def detect_scene(text: str) -> Optional[str]:
    text_lower = text.lower()
    for scene, keywords in SCENE_KEYWORDS.items():
        for kw in keywords:
            if kw.lower() in text_lower:
                logger.info("识别场景: %s (关键词: %s)", scene, kw)
                return scene
    return None


def list_scenes() -> None:
    logger.info("支持的 10 类场景：\n")
    for scene, schema in SCENE_SCHEMAS.items():
        logger.info("【%s】%s", scene, schema['name'])
        logger.info("  cover 字段: %s", list(schema['cover_fields'].keys()))
        logger.info("  章节结构:")
        for sec in schema["chapter_sections"]:
            logger.info("    - %s: %s（%s）", sec['key'], sec['name'], sec['desc'])
        logger.info("  end 字段: %s", list(schema['end_fields'].keys()))


def outline_to_business_data(outline):
    """
    将 outline.json 转换为 business_data.json

    :param outline: 大纲字典
    :return: business_data 字典
    """
    scene = outline.get("scene")
    if not scene:
        raise ValueError("outline 缺少 scene 字段")
    if scene not in SCENE_SCHEMAS:
        raise ValueError(f"不支持的场景: {scene}，支持: {list(SCENE_SCHEMAS.keys())}")

    schema = SCENE_SCHEMAS[scene]
    valid_keys = {s["key"] for s in schema["chapter_sections"]}

    # 转换 cover
    business_data = {"cover": dict(outline.get("cover", {}))}

    sections_in = outline.get("sections", [])
    sections_out: dict[str, list[dict[str, str]]] = {}
    for sec in sections_in:
        key = sec.get("key")
        if key not in valid_keys:
            logger.warning("跳过未知 section key: %s（场景 %s 不支持，有效: %s）", key, scene, valid_keys)
            continue
        items = sec.get("items", [])
        clean_items: list[dict[str, str]] = []
        for item in items:
            if not isinstance(item, dict):
                continue
            title = item.get("title", "").strip()
            desc = item.get("desc", "").strip()
            if not title and not desc:
                continue
            if desc and len(desc) > 80:
                logger.warning("section %s item [%s] desc 过长 (%d字)，建议30-60字", key, title, len(desc))
            clean_items.append({"title": title, "desc": desc})
        sections_out[key] = clean_items

    for sec_def in schema["chapter_sections"]:
        if sec_def["key"] not in sections_out:
            logger.info("section %s（%s）未提供，用空数组占位", sec_def['key'], sec_def['name'])
            sections_out[sec_def["key"]] = []

    business_data["sections"] = sections_out

    # 转换 end
    business_data["end"] = dict(outline.get("end", {"thanks": "感谢聆听"}))

    return business_data


def validate_outline(outline):
    """
    校验 outline 结构完整性

    :param outline: 大纲字典
    :return: (is_valid, issues 列表)
    """
    issues = []
    if not isinstance(outline, dict):
        return False, ["outline 应为对象"]

    scene = outline.get("scene")
    if not scene:
        issues.append("缺少 scene 字段")
    elif scene not in SCENE_SCHEMAS:
        issues.append(f"不支持的场景: {scene}")

    cover = outline.get("cover")
    if not isinstance(cover, dict):
        issues.append("cover 应为对象")

    sections = outline.get("sections")
    if not isinstance(sections, list):
        issues.append("sections 应为数组")
    else:
        for i, sec in enumerate(sections):
            if not isinstance(sec, dict):
                issues.append(f"sections[{i}] 应为对象")
                continue
            if "key" not in sec:
                issues.append(f"sections[{i}] 缺少 key 字段")
            if "items" not in sec:
                issues.append(f"sections[{i}] 缺少 items 字段")
            elif not isinstance(sec["items"], list):
                issues.append(f"sections[{i}].items 应为数组")

    return len(issues) == 0, issues


def generate_ppt(business_data, scene, template_id=None, output_path="final.pptx",
                 transitions="auto", animations="auto", templates_root="models",
                 animation_theme=None):
    """
    从 business_data 直接生成 PPT

    :param business_data: 业务数据字典
    :param scene: 场景名
    :param template_id: 模板 ID，None 时取该场景首个模板
    :param output_path: 输出路径
    :param transitions: 转场配置
    :param animations: 动画配置
    :param animation_theme: 动画预设主题名（business/tech/formal），None 表示不启用
    :return: 输出文件路径
    """
    from ppt_renderer import PptRenderer

    adapter = SceneAdapter(templates_root)

    # 模板选择
    if not template_id:
        templates = adapter.list_templates(category=scene)
        if not templates:
            raise ValueError(f"场景 {scene} 无可用模板")
        template_id = templates[0]["template_id"]
        print(f"ℹ️  未指定模板，使用场景首个模板: {template_id}")

    meta, meta_path = adapter.get_template_meta(template_id=template_id)
    # 直接从 meta_path 推导 pptx 路径（meta 文件名 xxx.meta.json → xxx.pptx）
    meta_path_obj = Path(meta_path)
    meta_name = meta_path_obj.name
    if meta_name.endswith(".meta.json"):
        pptx_name = meta_name[:-10] + ".pptx"
    else:
        pptx_name = meta_path_obj.stem + ".pptx"
    pptx_path = str(meta_path_obj.parent / pptx_name)

    is_valid, issues = adapter.validate_business_data(scene, business_data)
    if not is_valid:
        logger.warning("业务数据校验未通过:")
        for issue in issues:
            logger.warning("    - %s", issue)
        logger.info("将继续生成，但可能部分槽位无数据")

    slot_data = adapter.adapt(scene, business_data, meta)
    replaced = sum(len(v) for v in slot_data.values())
    logger.info("适配完成: %d 个槽位", replaced)

    renderer = PptRenderer(pptx_path, meta_path)
    renderer.render(slot_data, output_path,
                    remove_copyright=True, auto_fit=True,
                    transitions=transitions, animations=animations,
                    animation_theme=animation_theme)
    return output_path


def _extract_outline_page_effects(outline):
    """从 outline.pages 数组提取每页显式 transition/animation 配置

    仅提取 outline 中显式设置了 transition/animations 字段的页面，
    转换为渲染引擎所需的 per-page dict 格式。
    用于 --animation-theme 启用时，让单页显式配置优先于主题默认值。

    :param outline: outline dict
    :return: (transitions_dict, animations_dict) 两个 {page_num: spec} dict，
             仅包含显式设置的页
    """
    from aippt.animation_themes import build_page_transition_spec, build_page_animations_spec

    transitions_dict: dict[str, dict] = {}
    animations_dict: dict[str, list] = {}

    pages = outline.get("pages")
    if not pages or not isinstance(pages, list):
        return transitions_dict, animations_dict

    for page in pages:
        if not isinstance(page, dict):
            continue
        page_id = page.get("page_id")
        if page_id is None:
            continue
        page_type = page.get("page_type", "numbered_list")

        # transition：仅提取显式设置（None 表示未设置，不提取）
        explicit_t = page.get("transition")
        if explicit_t is not None and explicit_t != "none":
            t_spec = build_page_transition_spec(explicit_t)
            if t_spec:
                transitions_dict[str(page_id)] = t_spec

        # animations：仅提取显式设置
        explicit_a = page.get("animations")
        if explicit_a is not None:
            a_spec = build_page_animations_spec(page_type, explicit_a)
            if a_spec:
                animations_dict[str(page_id)] = a_spec

    return transitions_dict, animations_dict


def main():
    parser = argparse.ArgumentParser(
        description="AIPPT 大纲转换工具：outline.json → business_data.json → PPT"
    )
    sub = parser.add_subparsers(dest="cmd")

    # 子命令：list-scenes
    sub.add_parser("list-scenes", help="列出所有支持的场景")

    # 子命令：list-templates（查询模板列表）
    p_lt = sub.add_parser("list-templates", help="查询可用模板列表")
    p_lt.add_argument("--scene", default=None, help="按场景筛选（如 工作汇报）")
    p_lt.add_argument("--style", default=None, help="按风格标签筛选")
    p_lt.add_argument("--min-pages", type=int, default=None, help="最小页数")
    p_lt.add_argument("--max-pages", type=int, default=None, help="最大页数")
    p_lt.add_argument("--color-scheme", default=None,
                      help="按色系筛选（如 蓝色系/灰色系/红色系/绿色系/多彩/黑白）")
    p_lt.add_argument("--industry", default=None,
                      help="按适用行业筛选（如 通用/金融/教育/科技/制造/医疗/政府）")

    # 子命令：auto-generate（全链路一键生成）
    p_auto = sub.add_parser("auto-generate",
                            help="全链路一键生成：需求理解→大纲→模板匹配→渲染")
    p_auto.add_argument("--prompt", required=True, help="用户原始需求文本")
    p_auto.add_argument("--output", default="final.pptx", help="输出 PPT 路径")
    p_auto.add_argument("--template-id", default=None, help="指定模板 ID（不指定则自动匹配）")
    p_auto.add_argument("--transitions", default="auto", help="转场配置（auto/none）")
    p_auto.add_argument("--animations", default="auto", help="动画配置（auto/none）")
    p_auto.add_argument("--animation-theme", default=None,
                        help="动画预设主题（business/tech/formal），优先级低于单页配置，高于 --transitions/--animations")

    # 子命令：validate（格式校验）
    p_val = sub.add_parser("validate",
                           help="校验大纲格式合规性（六层防御体系）")
    p_val.add_argument("--outline", required=True, help="outline.json 路径")
    p_val.add_argument("--template-id", default=None, help="模板 ID（启用模板槽位匹配校验）")
    p_val.add_argument("--auto-fix", action="store_true", help="自动修复非原则性问题并回写文件")
    p_val.add_argument("--output", default=None, help="校验结果输出 JSON 路径（默认打印到 stdout）")

    # 子命令：convert（大纲 → business_data）
    p_conv = sub.add_parser("convert", help="将 outline.json 转换为 business_data.json")
    p_conv.add_argument("--outline", required=True, help="outline.json 路径")
    p_conv.add_argument("--output", required=True, help="business_data.json 输出路径")
    p_conv.add_argument("--validate", action="store_true", help="转换后校验")

    # 子命令：generate（大纲 → PPT，一步到位）
    p_gen = sub.add_parser("generate", help="从 outline.json 直接生成 PPT")
    p_gen.add_argument("--outline", required=True, help="outline.json 路径")
    p_gen.add_argument("--template-id", default=None, help="模板 ID（不指定则取场景首个）")
    p_gen.add_argument("--output", default="final.pptx", help="PPT 输出路径")
    p_gen.add_argument("--transitions", default="auto", help="转场配置（auto/none/dict）")
    p_gen.add_argument("--animations", default="auto", help="动画配置（auto/none/dict）")

    # 子命令：detect-scene（关键词推断场景）
    p_det = sub.add_parser("detect-scene", help="根据文本关键词推断场景")
    p_det.add_argument("--text", required=True, help="用户描述文本")

    # ============ 四步工作流子命令 ============
    # Step 1: 理解与拆解
    p_s1 = sub.add_parser("step1-understand",
                          help="Step 1: 理解用户意图，识别场景并生成澄清问题")
    p_s1.add_argument("--text", required=True, help="用户描述文本")
    p_s1.add_argument("--output", default=None, help="可选，将结果保存为 JSON 文件")

    # Step 2: 构建大纲
    p_s2 = sub.add_parser("step2-outline",
                          help="Step 2: 基于场景和用户回答生成大纲 JSON")
    p_s2.add_argument("--scene", required=True, help="场景名（Step 1 产出）")
    p_s2.add_argument("--purpose", required=True, help="演示用途/目标")
    p_s2.add_argument("--audience", default="", help="受众描述")
    p_s2.add_argument("--length", type=int, default=15, help="期望页数")
    p_s2.add_argument("--keys", default="", help="关键信息/数据点（逗号分隔）")
    p_s2.add_argument("--title", default="", help="封面主标题（不填则由用途推导）")
    p_s2.add_argument("--reporter", default="", help="汇报人/作者")
    p_s2.add_argument("--output", required=True, help="outline.json 输出路径")

    # Step 3: 视觉匹配
    p_s3 = sub.add_parser("step3-visuals",
                          help="Step 3: 模板推荐与版式匹配建议")
    p_s3.add_argument("--outline", required=True, help="outline.json 路径")
    p_s3.add_argument("--output", default=None, help="可选，将建议保存为 JSON 文件")

    # Step 4: 生成 PPT
    p_s4 = sub.add_parser("step4-generate",
                          help="Step 4: 渲染生成 PPT 成品")
    p_s4.add_argument("--outline", required=True, help="outline.json 路径")
    p_s4.add_argument("--template-id", default=None, help="模板 ID（Step 3 产出，不指定则取首个）")
    p_s4.add_argument("--output", default="final.pptx", help="PPT 输出路径")
    p_s4.add_argument("--transitions", default="auto", help="转场配置（auto/none）")
    p_s4.add_argument("--animations", default="auto", help="动画配置（auto/none）")
    p_s4.add_argument("--animation-theme", default=None,
                      help="动画预设主题（business/tech/formal），优先级低于单页配置，高于 --transitions/--animations")
    p_s4.add_argument("--trim-pages", default="", help="可选，要删除的页码（逗号分隔，如 6,10,11）")
    p_s4.add_argument("--insert-tables", action="store_true",
                      help="可选，自动插入表格页（对比表+报价表）")
    # T006：双引擎模式选择参数
    p_s4.add_argument("--mode", default="template", choices=["template", "auto"],
                      help="渲染模式：template（模板替换，默认）/ auto（无模板自动布局）")
    p_s4.add_argument("--theme", default=None,
                      help="视觉主题名（仅 auto 模式生效，如 商务蓝/极简灰/科技青）")

    args = parser.parse_args()

    if args.cmd == "list-scenes":
        list_scenes()

    elif args.cmd == "convert":
        outline = json.loads(Path(args.outline).read_text(encoding="utf-8"))
        is_valid, issues = validate_outline(outline)
        if not is_valid:
            print("❌ outline 校验失败:")
            for i in issues:
                print(f"    - {i}")
            return 1
        business_data = outline_to_business_data(outline)
        Path(args.output).write_text(
            json.dumps(business_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        print(f"✅ 转换完成: {args.output}")
        if args.validate:
            adapter = SceneAdapter("models")
            ok, v_issues = adapter.validate_business_data(outline["scene"], business_data)
            if ok:
                print("✅ 业务数据校验通过")
            else:
                print("⚠️  业务数据校验:")
                for i in v_issues:
                    print(f"    - {i}")
        return 0

    elif args.cmd == "generate":
        outline = json.loads(Path(args.outline).read_text(encoding="utf-8"))
        business_data = outline_to_business_data(outline)
        # 中间产物保存
        biz_path = Path(args.outline).with_suffix(".business.json")
        biz_path.write_text(
            json.dumps(business_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        print(f"ℹ️  business_data 已保存: {biz_path}")
        generate_ppt(business_data, outline["scene"],
                     template_id=args.template_id, output_path=args.output,
                     transitions=args.transitions, animations=args.animations)
        return 0

    elif args.cmd == "detect-scene":
        scene = detect_scene(args.text)
        if scene:
            print(f"识别场景: {scene}")
            print(f"场景定义: {SCENE_SCHEMAS[scene]['name']}")
            print(f"章节结构:")
            for sec in SCENE_SCHEMAS[scene]["chapter_sections"]:
                print(f"  - {sec['key']}: {sec['name']}")
        else:
            print("未匹配到场景，请明确指定")
            print("支持的场景:", list(SCENE_SCHEMAS.keys()))
        return 0

    elif args.cmd == "list-templates":
        adapter = SceneAdapter("models")
        templates = adapter.list_templates(
            category=args.scene, style_tag=args.style,
            min_pages=args.min_pages, max_pages=args.max_pages,
            color_scheme=args.color_scheme, industry=args.industry,
        )
        print(json.dumps(templates, ensure_ascii=False, indent=2))
        return 0

    elif args.cmd == "validate":
        # 六层防御体系 · 格式校验入口
        from aippt.validators import validate_all
        outline_path = Path(args.outline)
        try:
            outline = json.loads(outline_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as e:
            result_dict = {
                "validate_pass": False,
                "error_count": 1,
                "warning_count": 0,
                "fixed_count": 0,
                "errors": [{
                    "code": "F005", "level": "error",
                    "path": "(root)", "message": f"JSON 解析失败: {e}",
                    "suggestion": "请检查 JSON 语法（括号配对、逗号、引号）",
                }],
                "warnings": [], "fixes": [],
            }
            print(json.dumps(result_dict, ensure_ascii=False, indent=2))
            return 1

        meta = None
        if args.template_id:
            adapter = SceneAdapter("models")
            meta, _ = adapter.get_template_meta(template_id=args.template_id)

        fixed_outline, result = validate_all(outline, meta=meta, auto_fix=args.auto_fix)
        result_dict = result.to_dict()

        if args.auto_fix and (result.fixed or not result.is_valid):
            outline_path.write_text(
                json.dumps(fixed_outline, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            if result.fixed:
                logger.info("已自动修复 %d 项并回写: %s", len(result.fixed), outline_path)

        if args.output:
            Path(args.output).write_text(
                json.dumps(result_dict, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            print(f"校验结果已保存: {args.output}")
        else:
            print(json.dumps(result_dict, ensure_ascii=False, indent=2))

        return 0 if result.is_valid else 1

    elif args.cmd == "auto-generate":
        # 全链路一键生成：需求理解 → 大纲 → 模板匹配 → 渲染
        import time as _time
        start_ts = _time.time()

        # Step 1: 场景识别
        scene = detect_scene(args.prompt)
        if not scene:
            scene = "工作汇报"
            logger.warning("未识别到场景，使用默认: %s", scene)
        logger.info("[1/4] 识别场景: %s", scene)

        # Step 2: 自动生成大纲（基于 prompt 生成结构化 outline）
        # 注意：真正的大纲内容需由宿主大模型生成，这里提供骨架并填入 prompt 作为标题来源
        schema = SCENE_SCHEMAS[scene]
        first_line = args.prompt.strip().split("\n")[0][:40]
        outline = {
            "scene": scene,
            "cover": {
                "title": first_line,
                "reporter": "",
                "period": "",
            },
            "sections": [
                {"key": sec["key"], "name": sec["name"], "items": []}
                for sec in schema["chapter_sections"]
            ],
            "end": {"thanks": "感谢聆听"},
        }
        # 保存中间大纲文件
        outline_path = Path(args.output).with_suffix(".outline.json")
        outline_path.write_text(
            json.dumps(outline, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        logger.info("[2/4] 大纲骨架已生成: %s（请宿主大模型填充 sections 内容后重跑 step4-generate）",
                    outline_path)

        # Step 3: 模板匹配
        template_id = args.template_id
        if not template_id:
            templates = SceneAdapter("models").list_templates(category=scene)
            if templates:
                template_id = templates[0]["template_id"]
                logger.info("[3/4] 自动选择模板: %s（共 %d 个可选）",
                            template_id, len(templates))
            else:
                logger.error("场景 %s 无可用模板，无法生成", scene)
                return 1
        else:
            logger.info("[3/4] 使用指定模板: %s", template_id)

        # Step 4: 渲染
        business_data = outline_to_business_data(outline)
        biz_path = Path(args.output).with_suffix(".business.json")
        biz_path.write_text(
            json.dumps(business_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        # 动画主题：提取 outline 单页显式配置（优先级高于主题），主题默认由渲染引擎注入
        anim_theme = args.animation_theme
        t_cfg = args.transitions
        a_cfg = args.animations
        if anim_theme:
            from aippt.animation_themes import get_theme, list_themes
            if anim_theme not in list_themes():
                logger.error("未知动画主题: %s（可用: %s）", anim_theme, list_themes())
                return 1
            logger.info("启用动画主题: %s", anim_theme)
            t_dict, a_dict = _extract_outline_page_effects(outline)
            # 单页显式配置以 dict 形式传入，覆盖主题默认值
            if t_dict:
                t_cfg = t_dict
            if a_dict:
                a_cfg = a_dict
        generate_ppt(business_data, scene, template_id=template_id,
                     output_path=args.output,
                     transitions=t_cfg, animations=a_cfg,
                     animation_theme=anim_theme)
        elapsed = _time.time() - start_ts
        logger.info("[4/4] 生成完成，耗时 %.1fs", elapsed)
        logger.info("成品 PPT: %s", args.output)
        logger.info("⚠️  auto-generate 生成的是骨架 PPT，建议用 step2-outline 填充真实内容后用 step4-generate 重渲染")
        return 0

    # ============ 四步工作流处理 ============
    elif args.cmd == "step1-understand":
        # Step 1: 理解与拆解
        scene = detect_scene(args.text)
        result = {
            "step": 1,
            "user_input": args.text,
            "detected_scene": scene,
            "scene_name": SCENE_SCHEMAS[scene]["name"] if scene else None,
            "clarification_questions": [
                "这次 PPT 的主要用途是什么？受众是谁？",
                "有偏好的风格吗（学术/商务/创意）？",
                "大概期望多少页？",
                "有必须包含的关键信息或数据吗？",
            ],
            "next_step": "回答上述问题后，使用 step2-outline 命令生成大纲",
        }
        if scene:
            result["scene_chapters"] = [
                {"key": s["key"], "name": s["name"]}
                for s in SCENE_SCHEMAS[scene]["chapter_sections"]
            ]
        logger.info(json.dumps(result, ensure_ascii=False, indent=2))
        if args.output:
            Path(args.output).write_text(
                json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
            )
            logger.info("结果已保存: %s", args.output)
        logger.info("--- 确认 gate ---")
        logger.info("请回答澄清问题，确认后进入 Step 2 (step2-outline)")
        return 0

    elif args.cmd == "step2-outline":
        # Step 2: 构建大纲
        scene = args.scene
        if scene not in SCENE_SCHEMAS:
            print(f"❌ 不支持的场景: {scene}")
            print("支持的场景:", list(SCENE_SCHEMAS.keys()))
            return 1
        schema = SCENE_SCHEMAS[scene]
        chapter_sections = schema["chapter_sections"]

        # 由用途推导封面标题
        cover_title = args.title or args.purpose

        # 关键信息拆解
        keys = [k.strip() for k in args.keys.split(",") if k.strip()] if args.keys else []

        # 构建大纲：每个 section 放入关键信息或默认占位
        sections = []
        for sec_def in chapter_sections:
            items = []
            # 把关键信息分配到各 section（按顺序）
            for i in range(2):  # 每 section 默认 2 条占位
                key_idx = len(sections) * 2 + i
                if key_idx < len(keys):
                    items.append({"title": keys[key_idx], "desc": f"关于 {keys[key_idx]} 的详细说明"})
                else:
                    items.append({"title": f"{sec_def['name']}要点{i+1}",
                                  "desc": f"{sec_def['desc']}，待补充具体内容"})
            sections.append({"key": sec_def["key"], "name": sec_def["name"], "items": items})

        outline = {
            "scene": scene,
            "cover": {
                "title": cover_title,
                "reporter": args.reporter,
                "period": "",
            },
            "sections": sections,
            "end": {"thanks": "感谢聆听"},
        }

        Path(args.output).write_text(
            json.dumps(outline, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        logger.info("大纲已生成: %s", args.output)
        logger.info("--- 大纲预览 ---")
        logger.info("场景: %s（%s）", scene, schema['name'])
        logger.info("封面: %s", cover_title)
        logger.info("章节: %d 个", len(sections))
        for sec in sections:
            logger.info("  [%s] %s: %d 条要点", sec['key'], sec['name'], len(sec['items']))
        logger.info("--- 确认 gate ---")
        logger.info("请审阅大纲结构，如需调整请编辑 outline.json")
        logger.info("确认后进入 Step 3 (step3-visuals)")
        return 0

    elif args.cmd == "step3-visuals":
        # Step 3: 视觉匹配
        outline = json.loads(Path(args.outline).read_text(encoding="utf-8"))
        scene = outline.get("scene")
        if not scene:
            print("❌ outline 缺少 scene 字段")
            return 1

        adapter = SceneAdapter("models")
        templates = adapter.list_templates(category=scene)

        # 版式匹配建议（基于大纲内容特征）
        section_count = len(outline.get("sections", []))
        total_items = sum(len(s.get("items", [])) for s in outline.get("sections", []))

        recommendations = []
        for t in templates:
            # 读取 meta 判断页面特征
            try:
                meta, _ = adapter.get_template_meta(template_id=t["template_id"])
                page_meta = meta.get("page_meta", {})
                has_chart = any(pm.get("has_chart") for pm in page_meta.values())
                has_picture = any(pm.get("has_picture") for pm in page_meta.values())
            except Exception:
                has_chart = has_picture = False

            # 查找截图路径（与 pptx 同名 .png）
            import os
            pptx_path = os.path.join("models", t["path"])
            png_path = pptx_path.replace(".pptx", ".png")
            screenshot = t["path"].replace(".pptx", ".png") if os.path.exists(png_path) else None

            score = 0
            reasons = []
            if has_chart:
                score += 2
                reasons.append("含图表页")
            if has_picture:
                score += 1
                reasons.append("含图片位")
            if t["total_pages"] >= total_items + section_count + 2:
                score += 2
                reasons.append(f"页数充足({t['total_pages']}页)")
            recommendations.append({
                "template_id": t["template_id"],
                "name": t["name"],
                "total_pages": t["total_pages"],
                "score": score,
                "reasons": reasons,
                "screenshot": screenshot,
            })

        # 按得分降序
        recommendations.sort(key=lambda x: x["score"], reverse=True)

        # 版式匹配建议表
        layout_advice = [
            {"page_type": "封面页", "layout": "居中大标题+底部信息"},
            {"page_type": "目录页", "layout": "四宫格编号列表"},
            {"page_type": "章节分隔页", "layout": "左侧大数字+右侧标题"},
            {"page_type": "要点列表页", "layout": "标题+描述双栏"},
            {"page_type": "数据展示页", "layout": "编号步骤+描述"},
            {"page_type": "结束页", "layout": "居中大字+THANK YOU"},
        ]

        result = {
            "step": 3,
            "scene": scene,
            "template_recommendations": recommendations,
            "layout_advice": layout_advice,
            "style_decision": {
                "主色调": "沿用模板原商务蓝调",
                "字体": "模板自带无衬线字体",
                "动画": "淡入为主，章节分隔页加擦除",
                "转场": "全页 fade 淡入转场（speed=med）",
            },
            "next_step": "选定 template_id 后，使用 step4-generate 命令生成 PPT",
        }
        logger.info(json.dumps(result, ensure_ascii=False, indent=2))
        if args.output:
            Path(args.output).write_text(
                json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
            )
            logger.info("建议已保存: %s", args.output)
        logger.info("--- 确认 gate ---")
        logger.info("请确认模板选择和动画/转场配置，确认后进入 Step 4 (step4-generate)")
        return 0

    elif args.cmd == "step4-generate":
        # Step 4: 生成 PPT
        outline = json.loads(Path(args.outline).read_text(encoding="utf-8"))

        # T006：双引擎模式分支
        # - mode=auto：无模板自动布局，跳过模板相关校验，直接走 AutoLayoutRenderer
        # - mode=template：原有模板替换逻辑（100% 向后兼容）
        if args.mode == "auto":
            from aippt.render import AutoLayoutRenderer, RenderArgs
            from aippt.theme_loader import list_themes
            from aippt.validators import validate_all as _validate_all_auto

            # auto 模式仅校验 outline 基础结构（无模板 meta）
            fixed_outline, pre_check = _validate_all_auto(outline, meta=None, auto_fix=True)
            if pre_check.fixed:
                logger.info("渲染前自动修复 %d 项", len(pre_check.fixed))
                for fix in pre_check.fixed:
                    logger.info("  ✓ %s", fix)
                Path(args.outline).write_text(
                    json.dumps(fixed_outline, ensure_ascii=False, indent=2), encoding="utf-8"
                )
            if pre_check.errors:
                logger.error("渲染前终检发现 %d 项错误，已拦截:", len(pre_check.errors))
                for err in pre_check.errors:
                    logger.error("  [%s] %s: %s", err.code, err.path, err.message)
                print(json.dumps(pre_check.to_dict(), ensure_ascii=False, indent=2))
                return 1

            # 主题校验（未指定时使用商务蓝）
            theme_name = args.theme or "商务蓝"
            available_themes = list_themes()
            if available_themes and theme_name not in available_themes:
                logger.warning("主题 %s 未找到，可用: %s，降级到默认主题",
                               theme_name, available_themes)

            renderer = AutoLayoutRenderer(theme_name=theme_name)
            render_args = RenderArgs(
                transitions=args.transitions,
                animations=args.animations,
                animation_theme=args.animation_theme,
                theme=theme_name,
            )
            result = renderer.render_outline(fixed_outline, args.output, render_args)

            logger.info("--- 质量校验 ---")
            logger.info("总页数: %d", result.total_pages)
            logger.info("元素数: %d", result.replaced)
            logger.info("渲染模式: %s, 主题: %s", result.mode, theme_name)
            logger.info("成品 PPT: %s", args.output)
            print(json.dumps({
                "mode": result.mode,
                "theme": theme_name,
                "output_path": result.output_path,
                "total_pages": result.total_pages,
                "element_count": result.replaced,
            }, ensure_ascii=False, indent=2))
            return 0

        # === 渲染前终检（六层防御体系 Layer 3-5）===
        from aippt.validators import validate_all
        meta = None
        if args.template_id:
            try:
                adapter = SceneAdapter("models")
                meta, _ = adapter.get_template_meta(template_id=args.template_id)
            except Exception:
                pass

        fixed_outline, pre_check = validate_all(outline, meta=meta, auto_fix=True)
        if pre_check.fixed:
            logger.info("渲染前自动修复 %d 项", len(pre_check.fixed))
            for fix in pre_check.fixed:
                logger.info("  ✓ %s", fix)
            Path(args.outline).write_text(
                json.dumps(fixed_outline, ensure_ascii=False, indent=2), encoding="utf-8"
            )

        if pre_check.errors:
            logger.error("渲染前终检发现 %d 项错误，已拦截:", len(pre_check.errors))
            for err in pre_check.errors:
                logger.error("  [%s] %s: %s", err.code, err.path, err.message)
                if err.suggestion:
                    logger.error("    → %s", err.suggestion)
            print(json.dumps(pre_check.to_dict(), ensure_ascii=False, indent=2))
            return 1

        if pre_check.warnings:
            logger.warning("渲染前终检 %d 项警告（可继续渲染）:", len(pre_check.warnings))
            for warn in pre_check.warnings:
                logger.warning("  [%s] %s: %s", warn.code, warn.path, warn.message)

        outline = fixed_outline  # 使用修复后的大纲
        business_data = outline_to_business_data(outline)

        # 保存 business_data 中间产物
        biz_path = Path(args.outline).with_suffix(".business.json")
        biz_path.write_text(
            json.dumps(business_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        logger.info("business_data 已保存: %s", biz_path)

        # 动画主题：提取 outline 单页显式配置（优先级高于主题），主题默认由渲染引擎注入
        anim_theme = args.animation_theme
        t_cfg = args.transitions
        a_cfg = args.animations
        if anim_theme:
            from aippt.animation_themes import list_themes
            if anim_theme not in list_themes():
                logger.error("未知动画主题: %s（可用: %s）", anim_theme, list_themes())
                return 1
            logger.info("启用动画主题: %s", anim_theme)
            t_dict, a_dict = _extract_outline_page_effects(outline)
            # 单页显式配置以 dict 形式传入，覆盖主题默认值
            if t_dict:
                t_cfg = t_dict
            if a_dict:
                a_cfg = a_dict
        generate_ppt(business_data, outline["scene"],
                     template_id=args.template_id, output_path=args.output,
                     transitions=t_cfg, animations=a_cfg,
                     animation_theme=anim_theme)
        if args.trim_pages:
            pages = [int(x.strip()) for x in args.trim_pages.split(",") if x.strip()]
            from trim_ppt import trim_slides
            tmp_path = args.output + ".tmp"
            trim_slides(args.output, tmp_path, pages)
            import os
            os.replace(tmp_path, args.output)
            logger.info("已裁剪页面: %s", pages)

        if args.insert_tables:
            from insert_tables import main as insert_main
            tmp_path = args.output + ".tmp"
            import sys as _sys
            _sys.argv = ["insert_tables.py", args.output, tmp_path]
            insert_main()
            import os
            os.replace(tmp_path, args.output)

        from pptx import Presentation
        from ppt_meta_tool import extract_page_texts
        from aippt.constants import SLOT_MATCH_KEYWORDS
        prs = Presentation(args.output)
        issues = 0
        for slide in prs.slides:
            for t in extract_page_texts(slide):
                for kw in SLOT_MATCH_KEYWORDS:
                    if kw in t['text']:
                        issues += 1
                        break

        logger.info("--- 质量校验 ---")
        logger.info("总页数: %d", len(prs.slides))
        logger.info("残留占位文本: %d 处", issues)
        logger.info("成品 PPT: %s", args.output)
        logger.info("--- 确认 gate ---")
        logger.info("请打开 PPT 验收，如需迭代修改请告知具体调整")
        return 0

    else:
        parser.print_help()
        return 1


if __name__ == "__main__":
    import sys
    sys.exit(main() or 0)
