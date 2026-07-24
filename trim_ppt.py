"""
PPT 页面裁剪工具
功能：删除指定页码的 slide，从后往前删避免页码偏移
用法：python trim_ppt.py <input.pptx> <output.pptx> <page1> <page2> ...
示例：python trim_ppt.py vnerp_demo.pptx vnerp_trim.pptx 6 10 11 15 17 19
"""
import sys
from pptx import Presentation


def trim_slides(input_path, output_path, pages_to_delete):
    """
    删除指定页码的 slide（1-based）

    :param input_path: 输入 PPT 路径
    :param output_path: 输出 PPT 路径
    :param pages_to_delete: 要删除的页码列表（1-based）
    :return: 删除后的页数
    """
    prs = Presentation(input_path)
    total = len(prs.slides)
    print(f"原始页数: {total}")

    # 从后往前删，避免页码偏移
    sorted_pages = sorted(set(pages_to_delete), reverse=True)
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)

    deleted = []
    for page in sorted_pages:
        if page < 1 or page > total:
            print(f"⚠️  页码 {page} 超出范围（1-{total}），跳过")
            continue
        idx = page - 1
        xml_slides.remove(slides_list[idx])
        deleted.append(page)

    prs.save(output_path)
    remaining = total - len(deleted)
    print(f"删除页码: {sorted(deleted, reverse=True)}")
    print(f"剩余页数: {remaining}")
    print(f"✅ 已保存: {output_path}")
    return remaining


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("用法: python trim_ppt.py <input.pptx> <output.pptx> <page1> <page2> ...")
        sys.exit(1)
    input_pptx = sys.argv[1]
    output_pptx = sys.argv[2]
    pages = [int(x) for x in sys.argv[3:]]
    trim_slides(input_pptx, output_pptx, pages)
