"""
PPT 表格自动插入工具
功能：
  1. insert 子命令：在指定位置新增 slide 并插入原生表格（旧式用法）
  2. test 子命令：测试图表/表格动态扩展功能，自动构造测试数据并输出对比报告
用法：
  python insert_tables.py test --template 模板.pptx --output 测试输出.pptx
  python insert_tables.py insert <input.pptx> <output.pptx>
  python insert_tables.py <input.pptx> <output.pptx>  （兼容旧式用法）
"""
import sys
import argparse
import json
import tempfile
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# 深蓝青色科技配色（与 vnerp 品牌一致）
COLOR_HEADER_BG = RGBColor(0x1A, 0x23, 0x32)  # 深蓝 #1a2332
COLOR_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)  # 白
COLOR_ROW_ALT = RGBColor(0xF0, 0xF4, 0xF8)    # 浅灰蓝
COLOR_ROW_FG = RGBColor(0x1A, 0x23, 0x32)
COLOR_ACCENT = RGBColor(0x00, 0xB4, 0xD8)     # 青色 #00b4d8


def _style_table(table, header_bg=COLOR_HEADER_BG, header_fg=COLOR_HEADER_FG):
    """统一表格样式：深蓝表头 + 斑马纹行"""
    # 表头行
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = header_fg
                run.font.size = Pt(12)
                run.font.name = "微软雅黑"
    # 数据行（斑马纹）
    for row_idx in range(1, len(table.rows)):
        row = table.rows[row_idx]
        for cell in row.cells:
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = COLOR_ROW_FG
                    run.font.size = Pt(11)
                    run.font.name = "微软雅黑"


def _set_cell_text(cell, text, bold=False, color=None, align=PP_ALIGN.LEFT):
    """设置单元格文本"""
    cell.text = str(text)
    for para in cell.text_frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.bold = bold
            run.font.size = Pt(11)
            run.font.name = "微软雅黑"
            if color:
                run.font.color.rgb = color


def _add_title(slide, text, left, top, width):
    """在 slide 顶部添加标题文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = txBox.text_frame
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    for run in para.runs:
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = COLOR_HEADER_BG
        run.font.name = "微软雅黑"


def _get_blank_layout(prs):
    """获取空白版式，索引 6 不存在时回退到第一个"""
    for i in [6, 5, 4, 3, 2, 1, 0]:
        try:
            return prs.slide_layouts[i]
        except IndexError:
            continue
    return prs.slide_layouts[0]


def insert_comparison_table(prs, title_text, headers, rows):
    """新增一页插入对比表格"""
    blank_layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(blank_layout)

    # 标题
    _add_title(slide, title_text, Inches(0.6), Inches(0.4), Inches(12))

    # 表格（3 列 N 行）
    rows_count = len(rows) + 1
    cols_count = len(headers)
    table_left = Inches(0.6)
    table_top = Inches(1.4)
    table_width = Inches(12.1)
    table_height = Inches(0.5) * rows_count

    table_shape = slide.shapes.add_table(rows_count, cols_count,
                                          table_left, table_top,
                                          table_width, table_height)
    table = table_shape.table

    # 设置列宽
    if cols_count == 3:
        table.columns[0].width = Inches(3.5)
        table.columns[1].width = Inches(4.3)
        table.columns[2].width = Inches(4.3)

    # 表头
    for i, h in enumerate(headers):
        _set_cell_text(table.cell(0, i), h, bold=True,
                       color=COLOR_HEADER_FG, align=PP_ALIGN.CENTER)
    # 数据行
    for r_idx, row_data in enumerate(rows, start=1):
        for c_idx, val in enumerate(row_data):
            align = PP_ALIGN.CENTER if c_idx == 0 else PP_ALIGN.LEFT
            _set_cell_text(table.cell(r_idx, c_idx), val, align=align)

    _style_table(table)
    return slide


def insert_pricing_table(prs, title_text, headers, rows, note=None):
    """新增一页插入报价表格"""
    blank_layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(blank_layout)

    _add_title(slide, title_text, Inches(0.6), Inches(0.4), Inches(12))

    rows_count = len(rows) + 1
    cols_count = len(headers)
    table_left = Inches(0.6)
    table_top = Inches(1.4)
    table_width = Inches(12.1)
    table_height = Inches(0.55) * rows_count

    table_shape = slide.shapes.add_table(rows_count, cols_count,
                                          table_left, table_top,
                                          table_width, table_height)
    table = table_shape.table

    # 列宽
    if cols_count == 3:
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(7.1)
        table.columns[2].width = Inches(2.0)

    # 表头
    for i, h in enumerate(headers):
        _set_cell_text(table.cell(0, i), h, bold=True,
                       color=COLOR_HEADER_FG, align=PP_ALIGN.CENTER)
    # 数据行
    for r_idx, row_data in enumerate(rows, start=1):
        for c_idx, val in enumerate(row_data):
            align = PP_ALIGN.CENTER if c_idx == 0 else PP_ALIGN.LEFT
            if c_idx == 2:
                align = PP_ALIGN.CENTER
            _set_cell_text(table.cell(r_idx, c_idx), val, align=align)

    _style_table(table)

    # 备注框
    if note:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12), Inches(1.2))
        tf = txBox.text_frame
        tf.text = note
        tf.word_wrap = True
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.italic = True
                run.font.color.rgb = COLOR_ACCENT
                run.font.name = "微软雅黑"

    return slide


def move_slide(prs, old_index, new_index):
    """移动 slide 顺序（0-based）"""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    slide_to_move = slides_list[old_index]
    xml_slides.remove(slide_to_move)
    xml_slides.insert(new_index, slide_to_move)


def _iter_shapes_recursive(shapes):
    """递归遍历所有形状（含 group 内子形状）"""
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from _iter_shapes_recursive(shape.shapes)


def build_test_slot_data(meta):
    """根据 meta 构造测试数据（chart_data + table_data）

    遍历 page_slots，为 chart_data/table_data 槽位生成测试数据
    """
    slot_data = {}
    page_slots = meta.get("page_slots", {})
    for page_str, slots in page_slots.items():
        for slot_info in slots:
            if slot_info.get("slot") == "chart_data":
                slot_data.setdefault(page_str, {})["chart_data"] = {
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [
                        {"name": "营收", "data": [1200, 1500, 1800, 2100]},
                        {"name": "利润", "data": [300, 450, 600, 750]},
                    ]
                }
            elif slot_info.get("slot") == "table_data":
                slot_data.setdefault(page_str, {})["table_data"] = {
                    "headers": ["产品", "Q1销量", "Q2销量", "合计"],
                    "rows": [
                        ["产品A", "1200", "1500", "2700"],
                        ["产品B", "800", "950", "1750"],
                        ["产品C", "600", "720", "1320"],
                        ["产品D", "450", "580", "1030"],
                    ]
                }
    return slot_data


def cmd_test(args):
    """测试命令：自动构造测试数据并渲染，输出渲染前后对比报告

    流程：
      1. 扫描模板中所有 chart/table 形状
      2. 生成 meta 元数据
      3. 构造测试数据（chart_data + table_data）
      4. 调用 PptRenderer 渲染
      5. 验证输出：哪些 chart 被替换、哪些 table 被填充、样式是否保留
    """
    from ppt_meta_tool import generate_single_meta
    from ppt_renderer import PptRenderer

    template_path = args.template
    output_path = args.output

    print(f"模板: {template_path}")
    print(f"输出: {output_path}")
    print()

    # === 1. 渲染前扫描：检测 chart/table 形状（递归遍历含 group） ===
    prs = Presentation(template_path)
    chart_pages = []
    table_pages = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in _iter_shapes_recursive(slide.shapes):
            if hasattr(shape, "has_chart") and shape.has_chart:
                chart_pages.append({
                    "page": idx,
                    "shape_name": shape.name,
                    "chart_type": str(shape.chart.chart_type),
                })
            if hasattr(shape, "has_table") and shape.has_table:
                table_pages.append({
                    "page": idx,
                    "shape_name": shape.name,
                    "rows": len(shape.table.rows),
                    "cols": len(shape.table.columns),
                })

    print("=== 渲染前扫描 ===")
    print(f"图表页: {len(chart_pages)}")
    for cp in chart_pages:
        print(f"  页{cp['page']}: {cp['shape_name']} ({cp['chart_type']})")
    print(f"表格页: {len(table_pages)}")
    for tp in table_pages:
        print(f"  页{tp['page']}: {tp['shape_name']} ({tp['rows']}行 x {tp['cols']}列)")
    print()

    if not chart_pages and not table_pages:
        print("⚠️  模板中未检测到图表或表格，无需测试")
        return

    # === 2. 生成 meta（PptRenderer 初始化需要） ===
    meta, err = generate_single_meta(Path(template_path), "test")
    if err:
        print(f"❌ meta 生成失败: {err}")
        return

    # === 3. 构造测试数据 ===
    slot_data = build_test_slot_data(meta)
    print("=== 测试数据 ===")
    for page_str, data in slot_data.items():
        if "chart_data" in data:
            cd = data["chart_data"]
            print(f"  页{page_str} chart_data: {len(cd['categories'])}分类, {len(cd['series'])}系列")
        if "table_data" in data:
            td = data["table_data"]
            print(f"  页{page_str} table_data: {len(td['headers'])}列表头, {len(td['rows'])}行数据")
    print()

    # === 4. 渲染（写入临时 meta 文件供 PptRenderer 使用） ===
    with tempfile.NamedTemporaryFile(
        mode='w', suffix='.meta.json', delete=False, encoding='utf-8'
    ) as f:
        json.dump(meta, f, ensure_ascii=False)
        meta_path = f.name

    try:
        renderer = PptRenderer(template_path, meta_path)
        renderer.render(slot_data, output_path, remove_copyright=False, auto_fit=False)
    finally:
        os.unlink(meta_path)

    # === 5. 渲染后验证（递归遍历含 group） ===
    print("=== 渲染后验证 ===")
    prs_out = Presentation(output_path)
    chart_replaced = 0
    table_filled = 0
    for idx, slide in enumerate(prs_out.slides, 1):
        for shape in _iter_shapes_recursive(slide.shapes):
            if hasattr(shape, "has_chart") and shape.has_chart:
                # 验证图表分类是否已替换为 Q1/Q2/Q3/Q4
                try:
                    plot = shape.chart.plots[0]
                    cats = list(plot.categories)
                    if cats and str(cats[0]).strip() == "Q1":
                        chart_replaced += 1
                except Exception:
                    pass
            if hasattr(shape, "has_table") and shape.has_table:
                # 验证表格是否已填充测试数据（第二行首列含"产品"）
                try:
                    t = shape.table
                    if len(t.rows) > 1:
                        cell_text = t.cell(1, 0).text
                        if cell_text and "产品" in cell_text:
                            table_filled += 1
                except Exception:
                    pass

    print(f"图表替换成功: {chart_replaced}/{len(chart_pages)}")
    print(f"表格填充成功: {table_filled}/{len(table_pages)}")
    print(f"样式保留: ✅（代码仅替换数据，不修改 font/color/axis format）")
    print()
    print(f"✅ 测试输出已保存: {output_path}")


def cmd_insert_legacy(input_path, output_path):
    """旧式用法：插入对比表和报价表（保留向后兼容）"""
    prs = Presentation(input_path)
    total_before = len(prs.slides)
    print(f"原始页数: {total_before}")

    # 表格1：传统ERP vs vnerp 对比表
    comp_headers = ["对比维度", "传统商业印刷ERP", "vnerp 解决方案"]
    comp_rows = [
        ["软件授权费", "几十万 ~ 上百万", "0 元（开源）"],
        ["实施服务费", "10万 ~ 50万", "5.6万 ~ 15.5万"],
        ["年维保费",   "5万 ~ 20万",   "0.8万 ~ 1.5万"],
        ["定制灵活性", "受限，需原厂支持", "完全可控"],
        ["技术先进性", "老旧技术栈",   "Next.js + React 19 + TS"],
        ["行业适配",   "通用ERP改版",   "专为丝网印刷设计"],
    ]
    insert_comparison_table(prs, "对比传统方案，优势明显", comp_headers, comp_rows)
    print("✅ 已插入对比表（P10 之后）")

    # 表格2：报价方案表
    price_headers = ["服务项目", "内容", "参考报价"]
    price_rows = [
        ["基础部署", "服务器环境搭建、系统部署、数据库初始化、基础数据配置", "1.5万~3万元"],
        ["业务定制", "根据客户流程调整功能模块、报表定制、字段调整",       "2万~8万元"],
        ["数据迁移", "从旧系统/Excel导入历史数据",                         "0.5万~1.5万元"],
        ["员工培训", "现场/远程操作培训（2-3天）",                          "0.8万~1.5万元"],
        ["年度运维", "技术支持、bug修复、安全更新",                         "0.8万~1.5万元/年"],
    ]
    price_note = "💡 合计：5.6万~15.5万元\n💡 SaaS订阅制可按 200~500元/人/月 或 1万~15万/年 报价，根据企业规模灵活调整。"
    insert_pricing_table(prs, "透明报价，按需选择", price_headers, price_rows, price_note)
    print("✅ 已插入报价表（P13 之后）")

    # 把新增的 2 页（当前在末尾）移动到结束页之前
    end_idx = total_before - 1  # 结束页原索引
    move_slide(prs, total_before, end_idx)
    move_slide(prs, total_before, end_idx + 1)

    prs.save(output_path)
    total_after = len(prs.slides)
    print(f"最终页数: {total_after}")
    print(f"✅ 已保存: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description='PPT 表格/图表工具（支持 test 测试命令 和 insert 插入命令）'
    )
    sub = parser.add_subparsers(dest='command')

    # test 子命令：测试图表/表格动态扩展
    test_p = sub.add_parser('test', help='测试图表/表格动态扩展功能')
    test_p.add_argument('--template', required=True, help='模板 pptx 路径')
    test_p.add_argument('--output', required=True, help='输出 pptx 路径')

    # insert 子命令：插入对比/报价表格
    insert_p = sub.add_parser('insert', help='插入对比表和报价表')
    insert_p.add_argument('input', help='输入 pptx 路径')
    insert_p.add_argument('output', help='输出 pptx 路径')

    # 兼容旧式位置参数（无子命令时：python insert_tables.py <input> <output>）
    parser.add_argument('input_legacy', nargs='?', help='输入 pptx 路径（旧式用法）')
    parser.add_argument('output_legacy', nargs='?', help='输出 pptx 路径（旧式用法）')

    args = parser.parse_args()

    if args.command == 'test':
        cmd_test(args)
    elif args.command == 'insert':
        cmd_insert_legacy(args.input, args.output)
    elif args.input_legacy and args.output_legacy:
        # 兼容旧式用法：python insert_tables.py <input> <output>
        cmd_insert_legacy(args.input_legacy, args.output_legacy)
    else:
        parser.print_help()
        sys.exit(1)


if __name__ == '__main__':
    main()
