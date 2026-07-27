"""
PPT 模板批量导入工具
功能：
  1. 扫描源目录下所有 .pptx 模板
  2. 根据首页文本关键词自动分类到 10 大类目
  3. 复制到 models/<分类>/ 并重命名
  4. 生成 meta.json 元数据
  5. 用 PowerPoint COM 生成第一页截图 PNG
  6. 更新 templates_index.json 总索引
  7. 输出 preview_manifest.json 供 SKILL 展示选择

用法：
  python import_templates.py --src "1004 商务风" --prefix 商务风
  python import_templates.py --src "1004 商务风" --prefix 商务风 --force  # 覆盖已存在
  python import_templates.py --src "新模板" --prefix 新模板 --no-screenshot  # 跳过截图（无 PowerPoint 时）

依赖：
  - python-pptx
  - pywin32（仅 Windows + PowerPoint 时需要，用于截图）
"""
import argparse
import json
import shutil
import sys
import time
from pathlib import Path
from collections import defaultdict
from typing import Any, Optional

from pptx import Presentation

from ppt_meta_tool import generate_single_meta, META_REQUIRED_FIELDS
from aippt.logger import logger


# ==================== 分类规则 ====================
# 10 大类目关键词，按优先级从上到下匹配（先匹配先归）
CATEGORY_KEYWORDS = [
    ("年终总结", ["年终总结", "年度总结", "年终报告", "year end", "annual summary"]),
    ("工作总结", ["工作总结", "总结报告", "工作回顾", "阶段总结", "work summary"]),
    ("述职报告", ["述职", "述廉", "duty report"]),
    ("工作汇报", ["工作汇报", "项目汇报", "汇报材料", "work report"]),
    ("工作计划", ["工作计划", "计划书", "工作规划", "work plan"]),
    ("个人简历", ["简历", "竞聘", "求职", "resume", "cv", "个人简历"]),
    ("自我介绍", ["自我介绍", "self introduction", "personal intro"]),
    ("开题报告", ["开题", "答辩", "thesis proposal", "opening report"]),
    ("公司简介", ["公司介绍", "企业介绍", "公司简介", "company profile", "company intro"]),
    ("职业规划", ["职业规划", "生涯规划", "career plan", "career planning"]),
    ("安全教育", ["安全教育", "安全培训", "安全生产", "安全知识", "消防", "emergency", "safety"]),
]

# 默认分类（无法识别时归入）
DEFAULT_CATEGORY = "工作汇报"


def detect_category(pptx_path: Path) -> tuple[str, Optional[str]]:
    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        return DEFAULT_CATEGORY, None

    # 收集前 3 页文本
    texts = []
    for i, slide in enumerate(prs.slides):
        if i >= 3:
            break
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    texts.append(txt)

    combined = "\n".join(texts).lower()

    # 按关键词优先级匹配
    for category, keywords in CATEGORY_KEYWORDS:
        for kw in keywords:
            if kw.lower() in combined:
                return category, kw

    return DEFAULT_CATEGORY, None


# ==================== 截图生成 ====================
def export_screenshot_powerpoint(
    pptx_path: Path,
    output_png: Path,
    slide_index: int = 0,
    width: int = 800,
    height: int = 600,
) -> bool:
    try:
        import win32com.client
    except ImportError:
        logger.warning("pywin32 未安装，跳过截图: %s", pptx_path.name)
        return False

    ppt_app = None
    pres = None
    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        # WithWindow=False 隐藏 PowerPoint 窗口
        pres = ppt_app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        if slide_index >= pres.Slides.Count:
            slide_index = 0
        slide = pres.Slides.Item(slide_index + 1)  # COM 是 1-based
        slide.Export(str(output_png.resolve()), "PNG", width, height)
        return True
    except Exception as e:
        logger.warning("截图失败 %s: %s", pptx_path.name, e)
        return False
    finally:
        if pres is not None:
            try:
                pres.Close()
            except Exception:
                pass
        if ppt_app is not None:
            try:
                ppt_app.Quit()
            except Exception:
                pass


def export_screenshot(pptx_path: Path, output_png: Path, slide_index: int = 0, use_powerpoint: bool = True) -> bool:
    if use_powerpoint:
        return export_screenshot_powerpoint(pptx_path, output_png, slide_index)
    return False


# ==================== 主流程 ====================
def import_templates(
    src_dir: str,
    models_dir: str = "models",
    prefix: str = "模板",
    force: bool = False,
    use_screenshot: bool = True,
    start_index: int = 1,
    removable_tail: int = 0,
) -> Optional[dict[str, Any]]:
    src_path = Path(src_dir)
    models_path = Path(models_dir)

    if not src_path.exists():
        logger.error("源目录不存在: %s", src_path)
        return None

    pptx_files = sorted(src_path.glob("*.pptx"))
    if not pptx_files:
        logger.error("源目录无 .pptx 文件: %s", src_path)
        return None

    logger.info("扫描到 %d 个 pptx 文件", len(pptx_files))
    logger.info("目标目录: %s", models_path)
    logger.info("命名前缀: %s", prefix)
    logger.info("截图生成: %s", "是" if use_screenshot else "否")
    if removable_tail > 0:
        logger.info("末尾标记删除: %d 页", removable_tail)

    category_stats: dict[str, int] = defaultdict(int)
    preview_manifest: list[dict[str, Any]] = []
    failures: list[dict[str, str]] = []
    start = start_index
    total = len(pptx_files)

    for idx, src_pptx in enumerate(pptx_files, start=start):
        logger.info("[%d/%d] 处理: %s", idx, total + start - 1, src_pptx.name)

        category, matched_kw = detect_category(src_pptx)
        category_stats[category] += 1
        kw_info = f"（{matched_kw}）" if matched_kw else "（默认）"
        logger.info("  分类: %s %s", category, kw_info)

        target_dir = models_path / category
        target_dir.mkdir(parents=True, exist_ok=True)
        target_name = f"{prefix}_{idx:03d}.pptx"
        target_pptx = target_dir / target_name
        target_meta = target_dir / f"{target_name}.meta.json"
        target_png = target_dir / f"{target_name}.png"

        if target_pptx.exists() and not force:
            logger.info("  已存在，跳过: %s", target_pptx.name)
            continue

        try:
            shutil.copy2(src_pptx, target_pptx)
            logger.info("  复制: %s", target_pptx.relative_to(models_path))
        except Exception as e:
            logger.error("  复制失败: %s", e)
            failures.append({"file": src_pptx.name, "error": str(e)})
            continue

        try:
            meta_result = generate_single_meta(target_pptx, category)
            if isinstance(meta_result, tuple):
                meta, meta_error = meta_result
            else:
                meta, meta_error = meta_result, None
            if meta is None:
                logger.error("  meta 生成失败: %s", meta_error)
                failures.append({"file": src_pptx.name, "error": f"meta: {meta_error}"})
                continue

            if removable_tail > 0:
                total_p = meta.get('total_pages', 0)
                existing = meta.get('removable_pages', [])
                tail = list(range(total_p - removable_tail + 1, total_p + 1))
                meta['removable_pages'] = sorted(set(existing + tail))
                logger.info("  标记末尾 %d 页可删除: %s", removable_tail, tail)

            with open(target_meta, 'w', encoding='utf-8') as f:
                json.dump(meta, f, ensure_ascii=False, indent=2)
            logger.info("  meta: %s（%s页, %d页有槽位）",
                        target_meta.name, meta.get('total_pages', '?'),
                        len(meta.get('page_slots', {})))
            template_id = meta.get('template_id', '')
            total_pages = meta.get('total_pages', 0)
        except Exception as e:
            logger.error("  meta 生成异常: %s", e)
            failures.append({"file": src_pptx.name, "error": str(e)})
            template_id = ""
            total_pages = 0

        screenshot_ok = False
        if use_screenshot:
            time.sleep(0.3)
            screenshot_ok = export_screenshot(target_pptx, target_png)
            if screenshot_ok:
                logger.info("  截图: %s", target_png.name)
            else:
                logger.warning("  截图未生成")

        preview_manifest.append({
            "template_id": template_id,
            "category": category,
            "name": target_name.replace(".pptx", ""),
            "source_file": src_pptx.name,
            "pptx_path": str(target_pptx.relative_to(models_path)),
            "meta_path": str(target_meta.relative_to(models_path)),
            "screenshot": str(target_png.relative_to(models_path)) if screenshot_ok else None,
            "total_pages": total_pages,
            "matched_keyword": matched_kw,
        })

    logger.info("更新总索引...")
    update_templates_index(models_path)

    manifest_path = models_path / "preview_manifest.json"
    with open(manifest_path, 'w', encoding='utf-8') as f:
        json.dump(preview_manifest, f, ensure_ascii=False, indent=2)
    logger.info("预览清单已保存: %s", manifest_path)

    logger.info("导入统计：")
    for cat, cnt in sorted(category_stats.items()):
        logger.info("  %s: %d 套", cat, cnt)
    logger.info("  总计: %d 套", sum(category_stats.values()))
    if failures:
        logger.warning("失败 %d 个:", len(failures))
        for fail in failures:
            logger.warning("  - %s: %s", fail['file'], fail['error'])

    return {
        "total": sum(category_stats.values()),
        "categories": dict(category_stats),
        "failures": failures,
        "manifest": preview_manifest,
    }


def update_templates_index(models_dir: str) -> bool:
    models_path = Path(models_dir)
    meta_files = list(models_path.rglob("*.meta.json"))
    if not meta_files:
        logger.warning("未找到 meta 文件")
        return False

    index = {
        "total": 0,
        "categories": {},
        "templates": []
    }
    category_map = defaultdict(list)

    for meta_path in meta_files:
        try:
            with open(meta_path, 'r', encoding='utf-8') as f:
                meta = json.load(f)
        except Exception:
            continue

        template_info = {
            "template_id": meta.get('template_id', ''),
            "category": meta.get('category', ''),
            "name": meta_path.stem.replace('.meta', ''),
            "path": str(meta_path.relative_to(models_path)).replace('\\', '/'),
            "style_tags": meta.get('style_tags', []),
            "total_pages": meta.get('total_pages', 0),
            "chapter_count": len(meta.get('chapters', []))
        }
        index['templates'].append(template_info)
        category_map[meta.get('category', '未分类')].append(template_info)

    index['categories'] = {k: len(v) for k, v in category_map.items()}
    index['total'] = len(index['templates'])

    output_path = models_path / 'templates_index.json'
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(index, f, ensure_ascii=False, indent=2)

    logger.info("总索引更新完成，共 %d 个模板", index['total'])
    logger.info("索引文件: %s", output_path)
    for cat, cnt in sorted(index['categories'].items()):
        logger.info("  %s: %d 套", cat, cnt)
    return True


def auto_annotate(pptx_path: Path, scene: str, output_dir: Path) -> dict[str, Any]:
    """
    自定义模板自动标注：解析模板结构、识别槽位、分类页面类型，生成标准 .meta.json

    :param pptx_path: 模板 pptx 路径
    :param scene: 所属场景
    :param output_dir: 元数据输出目录
    :return: 标注结果字典（含 warnings / layout_profile 新增字段，向后兼容原有字段）
    """
    from ppt_meta_tool import generate_single_meta
    from aippt.profile_layouts import profile_layouts
    from aippt.ppt_element_classifier import classify_page

    if not pptx_path.exists():
        return {"ok": False, "error": f"文件不存在: {pptx_path}"}

    output_dir.mkdir(parents=True, exist_ok=True)
    meta_result = generate_single_meta(pptx_path, scene)
    if isinstance(meta_result, tuple):
        meta, err = meta_result
    else:
        meta, err = meta_result, None

    if meta is None:
        return {"ok": False, "error": f"解析失败: {err}"}

    meta_path = output_dir / f"{pptx_path.name}.meta.json"
    with open(meta_path, 'w', encoding='utf-8') as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

    slots_count = sum(len(v) for v in meta.get('page_slots', {}).values())

    # 采集 layout_profile（母版布局画像）与逐页低置信度告警
    layout_profile: dict[str, Any] = {}
    all_warnings: list[dict[str, Any]] = []
    try:
        prs = Presentation(str(pptx_path))
        layout_profile = profile_layouts(prs)
        for idx, slide in enumerate(prs.slides):
            try:
                classification = classify_page(slide)
                for w in classification.get("low_confidence_warnings", []):
                    all_warnings.append({"page": idx + 1, **w})
            except Exception as e:
                logger.debug("页面分类失败 page %d: %s", idx + 1, e)
    except Exception as e:
        logger.warning("layout_profile/分类告警采集失败: %s", e)

    # 更新全局索引（参考主流程 update_templates_index(models_path)）
    try:
        update_templates_index(output_dir.parent)
    except Exception as e:
        logger.warning("更新全局索引失败: %s", e)

    return {
        "ok": True,
        "template_id": meta.get("template_id"),
        "category": scene,
        "total_pages": meta.get("total_pages"),
        "chapter_count": len(meta.get("chapters", [])),
        "slots_count": slots_count,
        "removable_pages": meta.get("removable_pages", []),
        "meta_path": str(meta_path),
        "warnings": all_warnings,
        "layout_profile": layout_profile,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description='PPT 模板批量导入工具')
    sub = parser.add_subparsers(dest="cmd")

    # 默认子命令：import（保持向后兼容）
    p_import = sub.add_parser("import", help="批量导入模板（自动分类+复制+meta+截图）")
    p_import.add_argument('--src', required=True, help='源模板目录路径')
    p_import.add_argument('--models-dir', default='models', help='models 根目录（默认: models）')
    p_import.add_argument('--prefix', default='模板', help='模板重命名前缀（如"商务风"）')
    p_import.add_argument('--start', type=int, default=1, help='起始编号（默认: 1）')
    p_import.add_argument('--force', action='store_true', help='覆盖已存在的模板')
    p_import.add_argument('--no-screenshot', action='store_true', help='跳过截图生成（无 PowerPoint 时）')
    p_import.add_argument('--removable-tail', type=int, default=0,
                          help='标记每个模板末尾 N 页为可删除（版权页/致谢页），渲染时自动删除')

    # auto-annotate 子命令：自定义模板自动标注
    p_anno = sub.add_parser("auto-annotate", help="自定义模板自动标注，生成 .meta.json 元数据")
    p_anno.add_argument("--input", required=True, help="模板文件路径 .pptx")
    p_anno.add_argument("--scene", required=True, help="所属场景（如 工作汇报）")
    p_anno.add_argument("--output", default="models/_annotated", help="元数据输出目录")

    # rebuild-index 子命令：全量模板索引更新
    p_idx = sub.add_parser("rebuild-index", help="重新生成 templates_index.json 全局模板索引")
    p_idx.add_argument("--models-dir", default="models", help="models 根目录")

    args = parser.parse_args()

    # 向后兼容：无子命令时按旧方式处理（将所有参数视为 import）
    if args.cmd is None:
        # 旧模式兼容：重新解析为 import
        parser.parse_args(["import"] + sys.argv[1:])
        return

    if args.cmd == "import":
        import_templates(
            src_dir=args.src,
            models_dir=args.models_dir,
            prefix=args.prefix,
            force=args.force,
            use_screenshot=not args.no_screenshot,
            start_index=args.start,
            removable_tail=args.removable_tail,
        )
    elif args.cmd == "auto-annotate":
        result = auto_annotate(Path(args.input), args.scene, Path(args.output))
        print(json.dumps(result, ensure_ascii=False, indent=2))
    elif args.cmd == "rebuild-index":
        update_templates_index(args.models_dir)


if __name__ == "__main__":
    main()
