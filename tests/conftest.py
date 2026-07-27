"""pytest 共享 fixtures

提供跨测试文件复用的测试夹具：
  - tmp_pptx：生成临时空白 pptx 用于测试
  - sample_outline：返回标准 outline.json 字典（pages 数组格式）
  - sample_template_id：返回一个真实可用的模板 ID（从 models/ 中取首个）
"""
import json
import sys
from pathlib import Path

import pytest
from pptx import Presentation

from aippt.logger import logger

# 将项目根目录加入 sys.path，确保可导入根目录下的模块（ppt_scene_adapter 等）
ROOT = Path(__file__).parent.parent
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

MODELS_DIR = ROOT / "models"


@pytest.fixture
def tmp_pptx(tmp_path) -> str:
    """生成临时空白 pptx 用于测试

    使用 python-pptx 创建一个包含单张空白幻灯片的 pptx 文件，
    返回文件路径字符串。
    """
    prs = Presentation()
    # 使用空白布局（layout 索引 6 通常为空白）
    blank_layout = prs.slide_layouts[6]
    prs.slides.add_slide(blank_layout)
    out_path = tmp_path / "blank.pptx"
    prs.save(str(out_path))
    logger.debug(f"已生成临时 pptx: {out_path}")
    return str(out_path)


@pytest.fixture
def sample_outline() -> dict:
    """返回标准 outline.json 字典（pages 数组格式，5 页）

    该 outline 通过 schemas/outline.schema.json 校验，
    且通过 aippt.validators.validate_outline 无 error。
    """
    return {
        "scene": "工作总结",
        "purpose": "回顾本期工作成果与不足",
        "audience": "部门同事",
        "total_pages": 5,
        "pages": [
            {
                "page_id": 1,
                "page_type": "cover",
                "title": "本期工作总结",
                "subtitle": "驭势而上 砺行致远",
            },
            {
                "page_id": 2,
                "page_type": "numbered_list",
                "title": "核心工作成果",
                "items": ["产品迭代升级", "用户增长突破", "供应链优化"],
            },
            {
                "page_id": 3,
                "page_type": "kpi",
                "title": "关键指标",
                "kpi_items": [
                    {"label": "活跃用户", "value": "120w", "trend": "+35%"},
                    {"label": "采购成本", "value": "降15%", "trend": "同比下降"},
                ],
            },
            {
                "page_id": 4,
                "page_type": "timeline",
                "title": "项目里程碑",
                "timeline_items": [
                    {"time": "2026-Q1", "event": "V2.0 启动"},
                    {"time": "2026-Q2", "event": "全量发布"},
                    {"time": "2026-Q3", "event": "复盘优化"},
                ],
            },
            {
                "page_id": 5,
                "page_type": "ending",
                "title": "感谢聆听",
            },
        ],
    }


@pytest.fixture
def sample_template_id() -> str:
    """返回一个真实可用的模板 ID（从 models/ 中取首个）

    遍历 models/ 目录下的 *.meta.json 文件，返回首个包含 template_id 的元数据。
    """
    for meta_path in sorted(MODELS_DIR.rglob("*.meta.json")):
        try:
            with open(meta_path, "r", encoding="utf-8") as f:
                meta = json.load(f)
            tid = meta.get("template_id")
            if tid:
                logger.debug(f"sample_template_id 选用: {tid} (来自 {meta_path.name})")
                return tid
        except (json.JSONDecodeError, OSError):
            continue
    # 兜底：返回已知存在的模板 ID
    logger.warning("未能从 models/ 扫描到 template_id，使用兜底值")
    return "工作总结_工作总结"
