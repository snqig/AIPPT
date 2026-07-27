"""
AutoLayoutRenderer 无模板自动布局渲染器（T005）

继承 BaseRenderer，实现无模板自动生成 PPT 的渲染引擎。

执行链路：
    1. 初始化 Presentation 16:9 画布
    2. 遍历 outline.pages，调用 ppt_auto_layout 生成每页元素
    3. 收集页面内所有元素 shape_id + role 列表
    4. 调用公共动画模块执行动画、转场注入（T007 实现）
    5. 执行文本 auto_fit 自适应逻辑（T008 实现）

设计约束：
    - 不依赖任何模板文件，从零生成原生可编辑 PPTX 元素
    - 所有样式从主题令牌读取，禁止硬编码
    - 自动生成元素必须附加 role + shape_id
    - 与 PptRenderer 共用 BaseRenderer 抽象接口
"""
from __future__ import annotations

from typing import Any, Optional

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Pt

from aippt.logger import logger
from aippt.render import BaseRenderer, RenderArgs, RenderResult
from aippt.theme_loader import load_theme
from aippt.layout import (
    create_presentation, add_blank_slide,
    dispatch_page_layout, LayoutContext, PAGE_LAYOUT_REGISTRY,
    safe_area, column_x, row_y, GridRect, get_token,
    add_image_box,
)


class AutoLayoutRenderer(BaseRenderer):
    """无模板自动布局渲染器

    继承 BaseRenderer，MODE = "auto"。
    通过 ppt_auto_layout 引擎从零生成 PPT，无需模板文件。

    与 PptRenderer 区别：
        - PptRenderer：基于模板 + meta 槽位替换，100% 保留模板样式
        - AutoLayoutRenderer：无模板，按主题令牌自动布局生成原生元素
    """

    #: 引擎模式标识
    MODE: str = "auto"

    def __init__(self, theme_name: Optional[str] = None) -> None:
        """初始化自动布局渲染器

        :param theme_name: 主题名（如 "商务蓝"），None 时使用默认主题
        """
        self.theme_name = theme_name
        self.theme = load_theme(theme_name) if theme_name else load_theme("商务蓝")
        if theme_name and self.theme.get("name") == "默认主题" and theme_name != "默认主题":
            logger.warning("主题 %s 未找到，降级到默认主题", theme_name)

    def render_outline(
        self,
        outline_data: dict[str, Any],
        output_path: str,
        render_args: Optional[RenderArgs] = None,
    ) -> RenderResult:
        """BaseRenderer 统一接口实现：outline → 自动布局 → 渲染

        执行链路：
            1. 创建 16:9 空白 Presentation
            2. 遍历 outline.pages，按 page_type 分发到布局函数
            3. 收集每页元素 ElementMeta 列表（含 shape_id + role）
            4. 调用动画/转场注入（T007 实现后启用）
            5. 保存 PPTX

        :param outline_data: outline.json 原始字典，必须含 pages 数组
        :param output_path: 输出 PPTX 路径
        :param render_args: 渲染参数容器，None 使用默认值
        :return: RenderResult 渲染结果
        :raises ValueError: outline_data 缺少 pages 数组
        """
        args = self.normalize_args(render_args)

        # T613：从 args.extra 读取全局 variant_override 映射
        # 结构：{"kpi": "grid_2x2", "numbered_list": "single_column"}
        variant_override = args.extra.get("variant_override") if args.extra else None

        # 支持两种 outline 格式：pages 数组 或 cover/sections/end
        pages = outline_data.get("pages")
        if not pages:
            # 转换 cover/sections/end 格式为 pages 数组
            pages = self._convert_outline_to_pages(outline_data)
            if not pages:
                raise ValueError("outline_data 缺少 pages 数组，且无法从 cover/sections/end 转换")

        # T801：资产获取（图片/图标）
        page_assets: dict[int, dict[str, str]] = {}
        if args.enable_assets:
            page_assets = self._fetch_page_assets(pages, args)

        # 创建 16:9 画布
        prs = create_presentation()

        # 遍历页面分发到布局函数
        all_elements: list[dict] = []
        for page_data in pages:
            page_num = page_data.get("page_id", len(all_elements) + 1)
            page_type = page_data.get("page_type", "numbered_list")

            # 检查 page_type 是否已注册
            if page_type not in PAGE_LAYOUT_REGISTRY:
                logger.warning("页%s 未注册的 page_type: %s，降级为 numbered_list",
                               page_num, page_type)
                page_data = {**page_data, "page_type": "numbered_list"}

            slide = add_blank_slide(prs)
            ctx = LayoutContext(page_num=page_num)
            elements = dispatch_page_layout(
                slide, page_data, self.theme, ctx,
                variant_override=variant_override,
            )

            # T801: 资产注入（在布局元素之后放置图片/图标）
            page_asset_map = page_assets.get(page_num, {})
            if page_asset_map:
                self._apply_page_assets(slide, page_data, page_asset_map, self.theme, ctx)

            # 收集元素元数据（供动画模块使用）
            for el in elements:
                all_elements.append({
                    "page_num": page_num,
                    "shape_id": el.shape_id,
                    "role": el.role,
                    "page_type": page_type,
                })

        # T007：动画/转场注入（接入现有 ppt_animations / ppt_transitions 模块）
        if args.transitions or args.animations:
            self._inject_effects(prs, args, all_elements, pages)

        # T008：auto_fit 自适应（复用现有 aippt.text_replacer.auto_fit）
        # 对所有有文本的 shape 执行字号自适应，下限 10pt
        if args.auto_fit:
            self._apply_autofit(prs, all_elements)

        prs.save(str(output_path))
        total_pages = len(prs.slides)
        logger.info("自动布局渲染完成: %s", output_path)
        logger.info("总页数: %d, 总元素数: %d", total_pages, len(all_elements))

        return RenderResult(
            output_path=str(output_path),
            mode=self.MODE,
            total_pages=total_pages,
            replaced=len(all_elements),
            missed=0,
            skipped=0,
            removed_pages=[],
            warnings=[],
            meta={
                "theme": self.theme_name or self.theme.get("name", "默认主题"),
                "elements": all_elements,
            },
        )

    def _convert_outline_to_pages(self, outline: dict) -> list[dict]:
        """将 cover/sections/end 格式 outline 转换为 pages 数组

        用于兼容现有 cover/sections/end 格式的 outline.json。
        转换规则：
            - cover → cover 页
            - sections[].items → numbered_list 页（每个 section 一页）
            - end → ending 页（暂用 divider 布局兜底）

        :param outline: outline dict
        :return: pages 数组（可能为空）
        """
        pages: list[dict] = []
        page_id = 1

        # cover
        cover = outline.get("cover", {})
        if cover:
            pages.append({
                "page_id": page_id,
                "page_type": "cover",
                "title": cover.get("title", ""),
                "subtitle": cover.get("reporter", "") or cover.get("period", ""),
            })
            page_id += 1

        # sections → numbered_list
        sections = outline.get("sections", {})
        if isinstance(sections, dict):
            for sec_key, items in sections.items():
                if not isinstance(items, list) or not items:
                    continue
                # 提取标题列表
                titles = []
                for item in items:
                    if isinstance(item, dict):
                        titles.append(item.get("title", ""))
                    else:
                        titles.append(str(item))
                pages.append({
                    "page_id": page_id,
                    "page_type": "numbered_list",
                    "title": sec_key,
                    "items": titles,
                })
                page_id += 1

        return pages

    def _inject_effects(
        self,
        prs: Presentation,
        args: RenderArgs,
        elements: list[dict],
        outline_pages: Optional[list[dict]] = None,
    ) -> None:
        """T007：动画/转场注入（调度层对接）

        实现链路：
            1. 按页码分组元素，获取 page_type
            2. 调度层接管规则匹配、时序计算、三层层注入

        优先级（高 → 低）：
            1. outline.json 单页显式 transition/animations 字段（outline_pages）
            2. CLI args.transitions / args.animations（dict[page_num] 形式）
            3. --animation-theme 主题默认

        :param prs: Presentation 对象
        :param args: 渲染参数
        :param elements: 全部元素元数据列表
        :param outline_pages: outline.json 的 pages 数组，用于读取单页显式配置
        """
        try:
            from aippt.animation_scheduler import inject_page_effects, list_animation_themes
        except Exception as e:
            logger.warning("动画调度层模块加载失败，跳过注入: %s", e)
            return

        theme_name = args.animation_theme
        if theme_name and theme_name not in list_animation_themes():
            logger.warning("未知动画主题 %s，跳过主题注入", theme_name)
            theme_name = None

        # 构建 page_num → outline page_data 映射，读取单页显式 transition/animations
        outline_by_page: dict[int, dict] = {}
        if outline_pages:
            for pd in outline_pages:
                pnum = pd.get("page_id")
                if pnum is not None:
                    outline_by_page[int(pnum)] = pd

        page_elements: dict[int, list[dict]] = {}
        for el in elements:
            page_num = el.get("page_num", 0)
            page_elements.setdefault(page_num, []).append(el)

        injected_pages = 0
        for page_num, slide in enumerate(prs.slides, 1):
            page_els = page_elements.get(page_num, [])
            if not page_els:
                continue
            page_type = page_els[0].get("page_type", "numbered_list")

            # 优先级 1：outline 单页显式 transition/animations
            outline_pd = outline_by_page.get(page_num, {})
            t_name = outline_pd.get("transition")
            a_spec = outline_pd.get("animations")

            # 优先级 2：CLI args.transitions / args.animations（outline 未显式时使用）
            if t_name is None and isinstance(args.transitions, dict) and str(page_num) in args.transitions:
                t_spec = args.transitions[str(page_num)]
                if isinstance(t_spec, dict):
                    t_name = t_spec.get("type")
            if a_spec is None and isinstance(args.animations, dict) and str(page_num) in args.animations:
                a_spec = args.animations[str(page_num)]

            inject_page_effects(slide, page_type,
                theme_name=theme_name,
                page_animations=a_spec,
                page_transition=t_name,
            )
            injected_pages += 1

        logger.info("动画/转场注入完成: %d 页（主题=%s, 转场=%s, 动画=%s）",
                    injected_pages,
                    theme_name or "none",
                    args.transitions if args.transitions else "off",
                    args.animations if args.animations else "off")

    def _apply_autofit(self, prs: Presentation, elements: list[dict]) -> None:
        """T008：auto_fit 自适应（复用现有 aippt.text_replacer.auto_fit）

        对所有有文本的 shape 执行字号自适应：
            1. 遍历每页 slide 的所有 shape
            2. 跳过无文本框的 shape（如纯背景矩形、连接线）
            3. 调用 auto_fit 几何估算，超容量时按比例缩小字号
            4. 字号下限 10pt（防止过小不可读）

        :param prs: Presentation 对象
        :param elements: 全部元素元数据列表（用于日志上下文）
        """
        try:
            from aippt.text_replacer import auto_fit
        except Exception as e:
            logger.warning("auto_fit 模块加载失败，跳过自适应: %s", e)
            return

        MIN_FONT_PT = 10  # 字号下限，防止过小不可读
        adjusted_count = 0

        for page_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                # 仅处理有文本框且有文本的 shape
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # 获取当前字号（取首个 run）
                tf = shape.text_frame
                if not tf.paragraphs or not tf.paragraphs[0].runs:
                    continue
                run = tf.paragraphs[0].runs[0]
                if run.font.size is None:
                    continue
                original_size_pt = run.font.size.pt

                # 跳过已小于下限的字号
                if original_size_pt <= MIN_FONT_PT:
                    continue

                # 调用现有 auto_fit 进行几何估算与缩小
                try:
                    auto_fit(
                        shape=shape,
                        new_text=text,
                        original_text=text,
                        slot_info=None,
                        page_str=str(page_num),
                        slot_name=shape.name or "auto",
                    )
                except Exception as e:
                    logger.debug("auto_fit 跳过 shape %s: %s", shape.name, e)
                    continue

                # 强制下限保护：若 auto_fit 缩到 10pt 以下，回拉到 10pt
                try:
                    new_size = run.font.size
                    if new_size and new_size.pt < MIN_FONT_PT:
                        run.font.size = Pt(MIN_FONT_PT)
                        adjusted_count += 1
                except Exception:
                    pass

        if adjusted_count:
            logger.info("auto_fit 下限保护: %d 个 shape 字号回拉到 %dpt",
                        adjusted_count, MIN_FONT_PT)

    def _fetch_page_assets(
        self,
        pages: list[dict],
        args: RenderArgs,
    ) -> dict[int, dict[str, str]]:
        from aippt.assets.asset_planner import build_asset_plan
        from aippt.assets.asset_fetcher import fetch_assets

        result: dict[int, dict[str, str]] = {}
        asset_cache = args.asset_cache_dir
        scene = args.extra.get("scene", "") if args.extra else ""
        style = args.theme or ""

        for page_data in pages:
            page_num = page_data.get("page_id", 0)
            if not page_num:
                continue

            if page_data.get("assets"):
                plan = page_data["assets"]
            else:
                plan = build_asset_plan(page_data, scene=scene, style=style)

            if not plan:
                continue

            assets = fetch_assets(plan, cache_dir=asset_cache,
                                 enable_photos=args.enable_assets)
            page_map: dict[str, str] = {}
            for a in assets:
                page_map[a.slot] = a.local_path

            if page_map:
                result[page_num] = page_map

        total = sum(len(v) for v in result.values())
        if total:
            logger.info("Asset fetch: %d assets across %d pages", total, len(result))
        return result

    def _apply_page_assets(
        self,
        slide: Slide,
        page_data: dict,
        asset_map: dict[str, str],
        theme: dict,
        ctx: LayoutContext,
    ) -> None:
        page_type = page_data.get("page_type", "numbered_list")
        if not asset_map:
            return

        area = safe_area(get_token(theme, "spacing.safe_margin_inch", 0.5))

        for slot, local_path in asset_map.items():
            if slot.startswith("img_"):
                if page_type == "cover":
                    rect = GridRect(left=8.5, top=0.8, width=4.0, height=5.9)
                    add_image_box(slide, rect, local_path, "hero_image", ctx, theme,
                                 mask="rounded", overlay=True)
                elif page_type in ("numbered_list", "catalog"):
                    rect = GridRect(left=area.left, top=area.top,
                                    width=area.width, height=2.0)
                    add_image_box(slide, rect, local_path, "hero_image", ctx, theme,
                                 mask="rounded", overlay=False)
                elif page_type == "ending":
                    rect = GridRect(left=0, top=0, width=13.333, height=7.5)
                    add_image_box(slide, rect, local_path, "hero_image", ctx, theme,
                                 overlay=True)

            elif slot.startswith("icon_"):
                idx = int(slot.split("_")[1]) if "_" in slot else 0
                items = page_data.get("items", []) or []
                if idx >= len(items):
                    continue
                icon_y = area.top + 0.6 + idx * 0.85
                rect = GridRect(left=area.left + 0.1, top=icon_y,
                                width=0.4, height=0.4)
                add_image_box(slide, rect, local_path, "icon", ctx, theme)
