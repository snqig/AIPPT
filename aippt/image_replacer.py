"""
图片替换子模块 - 渲染引擎拆分（T301）
功能：图片槽位替换，支持本地路径/URL，等比覆盖/等比包含两种填充模式
依赖：python-pptx + Pillow（可选，无 Pillow 时降级为 contain 模式）

槽位数据格式（slot_data[page_str]["image_data"]）：
    {
        "<slot_name>": {
            "path": "/abs/path/to/image.png",   # 本地路径（与 url 二选一）
            "url": "https://example.com/img.jpg", # URL（与 path 二选一）
            "fit": "cover"                        # cover/contain，默认 cover
        }
    }

匹配策略：
    按 picture shape 出现顺序依次替换（第 N 个 picture ↔ image_data 第 N 项）。
    若 image_data 项数 > picture 数量，多余项忽略并 warning。
"""
from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import Any, Optional
from urllib.parse import urlparse
from urllib.request import urlopen

from pptx.shapes.base import BaseShape
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

from aippt.logger import logger


# ==================== 图片下载 ====================
def resolve_image_source(spec: dict[str, Any], cache_dir: Optional[str] = None) -> str:
    """解析图片来源，URL 自动下载到本地缓存

    :param spec: 图片规格字典，含 path 或 url 字段
    :param cache_dir: URL 下载缓存目录，None 时用系统临时目录
    :return: 本地图片文件路径
    :raises FileNotFoundError: 本地路径不存在
    :raises ValueError: path 和 url 都未提供
    """
    path = spec.get("path")
    url = spec.get("url")

    if path:
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"图片文件不存在: {path}")
        return str(p.resolve())

    if url:
        return _download_image(url, cache_dir)

    raise ValueError("图片规格必须包含 path 或 url 字段")


def _download_image(url: str, cache_dir: Optional[str] = None) -> str:
    """下载 URL 图片到本地缓存

    :param url: 图片 URL
    :param cache_dir: 缓存目录，None 时用系统临时目录
    :return: 本地文件路径
    """
    cache = Path(cache_dir) if cache_dir else Path(tempfile.gettempdir()) / "aippt_img_cache"
    cache.mkdir(parents=True, exist_ok=True)

    # 从 URL 推导文件名
    parsed = urlparse(url)
    name = Path(parsed.path).name or "image"
    # 加 URL 哈希避免重名
    import hashlib
    url_hash = hashlib.md5(url.encode("utf-8")).hexdigest()[:8]
    local_path = cache / f"{url_hash}_{name}"

    if local_path.exists():
        return str(local_path)

    logger.info("下载图片: %s → %s", url, local_path)
    with urlopen(url, timeout=30) as resp, open(local_path, "wb") as f:
        f.write(resp.read())
    return str(local_path)


# ==================== 等比裁剪/适配 ====================
def compute_fit_rect(
    target_left: int, target_top: int, target_w: int, target_h: int,
    img_w: int, img_h: int, fit: str = "cover",
) -> tuple[int, int, int, int]:
    """计算图片填充后的位置与尺寸（EMU 单位）

    :param target_left: 目标区域左上角 x（EMU）
    :param target_top: 目标区域左上角 y（EMU）
    :param target_w: 目标区域宽（EMU）
    :param target_h: 目标区域高（EMU）
    :param img_w: 原图宽（像素）
    :param img_h: 原图高（像素）
    :param fit: cover（等比覆盖，裁剪超出） / contain（等比包含，留白）
    :return: (left, top, width, height) EMU 四元组
    """
    if img_w <= 0 or img_h <= 0:
        return target_left, target_top, target_w, target_h

    target_ratio = target_w / target_h
    img_ratio = img_w / img_h

    if fit == "contain":
        # 等比包含：图片完整显示，留白居中
        if img_ratio > target_ratio:
            # 图片更宽，以宽度为准
            new_w = target_w
            new_h = int(target_w / img_ratio)
        else:
            new_h = target_h
            new_w = int(target_h * img_ratio)
        # 居中
        new_left = target_left + (target_w - new_w) // 2
        new_top = target_top + (target_h - new_h) // 2
        return new_left, new_top, new_w, new_h
    else:
        # cover：等比覆盖，图片可能超出目标区域（由 crop 处理）
        # 这里返回目标区域尺寸，裁剪由 _crop_image 处理
        return target_left, target_top, target_w, target_h


def crop_image_cover(
    img_path: str, target_w: int, target_h: int, img_w: int, img_h: int,
) -> tuple[str, int, int, int, int]:
    """cover 模式：用 Pillow 裁剪图片到目标比例，返回裁剪后路径 + 显示尺寸

    :param img_path: 原图本地路径
    :param target_w: 目标宽（EMU）
    :param target_h: 目标高（EMU）
    :param img_w: 原图宽（像素）
    :param img_h: 原图高（像素）
    :return: (cropped_path, display_left_offset, display_top_offset, display_w, display_h)
             display_* 为裁剪后图片在目标区域中的显示尺寸（EMU，居中）
    """
    target_ratio = target_w / target_h
    img_ratio = img_w / img_h

    try:
        from PIL import Image
    except ImportError:
        # 无 Pillow 时降级为 contain
        logger.warning("未安装 Pillow，cover 模式降级为 contain")
        new_w, new_h = (target_w, int(target_w / img_ratio)) if img_ratio > target_ratio \
            else (int(target_h * img_ratio), target_h)
        offset_l = (target_w - new_w) // 2
        offset_t = (target_h - new_h) // 2
        return img_path, offset_l, offset_t, new_w, new_h

    # 计算裁剪框（保持目标比例）
    if img_ratio > target_ratio:
        # 图片更宽，裁剪左右
        new_w_px = int(img_h * target_ratio)
        new_h_px = img_h
        left = (img_w - new_w_px) // 2
        top = 0
    else:
        # 图片更高，裁剪上下
        new_w_px = img_w
        new_h_px = int(img_w / target_ratio)
        left = 0
        top = (img_h - new_h_px) // 2

    # 裁剪并保存到临时文件
    img = Image.open(img_path)
    cropped = img.crop((left, top, left + new_w_px, top + new_h_px))
    cropped_path = str(Path(img_path).with_suffix(".cropped" + Path(img_path).suffix))
    cropped.save(cropped_path)
    img.close()

    # 裁剪后图片正好填充目标区域
    return cropped_path, 0, 0, target_w, target_h


# ==================== 图片替换主入口 ====================
def replace_images(
    slide: Any,
    image_data: dict[str, dict[str, Any]],
    cache_dir: Optional[str] = None,
) -> int:
    """替换 slide 中的 picture shape 为新图片

    匹配策略：按 picture shape 出现顺序依次替换 image_data 中的项。
    image_data 顺序按 dict 插入顺序（Python 3.7+ 保证）。

    :param slide: python-pptx Slide 对象
    :param image_data: 图片规格字典 {slot_name: {path/url, fit}}
    :param cache_dir: URL 下载缓存目录
    :return: 成功替换的图片数量
    """
    if not image_data:
        return 0

    # 收集 slide 中的所有 picture shape
    picture_shapes = [
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]

    if not picture_shapes:
        logger.warning("当前 slide 无 picture shape，跳过图片替换（共 %d 项）", len(image_data))
        return 0

    replaced = 0
    items = list(image_data.items())
    for idx, (slot_name, spec) in enumerate(items):
        if idx >= len(picture_shapes):
            logger.warning("图片槽位 %s 无对应 picture shape（图片槽位多于 picture）", slot_name)
            break

        shape = picture_shapes[idx]
        try:
            _replace_single_image(shape, spec, slot_name, cache_dir)
            replaced += 1
        except Exception as e:
            logger.warning("图片槽位 %s 替换失败: %s", slot_name, e)

    return replaced


def _replace_single_image(
    shape: BaseShape,
    spec: dict[str, Any],
    slot_name: str,
    cache_dir: Optional[str],
) -> None:
    """替换单个 picture shape

    策略：记录原 shape 位置尺寸 → 删除原 shape → 用 add_picture 插入新图片
    （python-pptx 未提供直接 replace_image API，删除+新增最稳妥）

    :param shape: 原 picture shape
    :param spec: 图片规格 {path/url, fit}
    :param slot_name: 槽位名（用于日志）
    :param cache_dir: URL 下载缓存目录
    """
    slide = shape._element.getparent()  # slide 的 XML 元素
    # 通过 shape.part 获取 slide 对象（用于 add_picture）
    from pptx.slide import Slide
    slide_obj: Optional[Slide] = getattr(shape, "part", None) and shape.part.slide

    # 记录原位置和尺寸（EMU）
    orig_left = shape.left
    orig_top = shape.top
    orig_w = shape.width
    orig_h = shape.height

    # 解析图片来源
    fit = spec.get("fit", "cover")
    img_path = resolve_image_source(spec, cache_dir)

    # 获取原图尺寸（像素）
    img_w_px, img_h_px = _get_image_size(img_path)

    # 计算填充尺寸
    if fit == "cover":
        cropped_path, offset_l, offset_t, new_w, new_h = crop_image_cover(
            img_path, orig_w, orig_h, img_w_px, img_h_px,
        )
        final_path = cropped_path
        final_left = orig_left + offset_l
        final_top = orig_top + offset_t
    else:
        # contain
        final_left, final_top, new_w, new_h = compute_fit_rect(
            orig_left, orig_top, orig_w, orig_h, img_w_px, img_h_px, "contain",
        )
        final_path = img_path

    # 删除原 shape
    sp = shape._element
    sp.getparent().remove(sp)

    # 插入新图片
    if slide_obj is not None:
        pic = slide_obj.shapes.add_picture(
            final_path, final_left, final_top, new_w, new_h,
        )
        # 尝试保留原 shape name（便于动画匹配）
        try:
            pic.name = f"image_{slot_name}"
        except Exception:
            pass
    else:
        logger.warning("图片槽位 %s 无法获取 slide 对象，仅删除原图片未插入新图", slot_name)


def _get_image_size(img_path: str) -> tuple[int, int]:
    """获取图片尺寸（像素）

    优先用 Pillow，无 Pillow 时用 PIL 替代或裸读图片头。

    :param img_path: 图片本地路径
    :return: (width, height) 像素
    """
    try:
        from PIL import Image
        with Image.open(img_path) as img:
            return img.size
    except ImportError:
        pass
    # 兜底：用 python-pptx 的 Image 读尺寸
    try:
        from pptx.util import Emu
        from pptx.image import Image
        img = Image.from_file(img_path)
        return img.size  # (px, py)
    except Exception as e:
        logger.warning("无法读取图片尺寸 %s: %s，使用默认 800x600", img_path, e)
        return 800, 600
