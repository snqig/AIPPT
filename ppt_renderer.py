"""
PPT 渲染引擎 - 阶段二核心模块
功能：基于 meta 元数据，将结构化内容精准替换到模板，100% 保留原样式
依赖：python-pptx
"""
import json
import argparse
from pathlib import Path
from copy import deepcopy
from typing import Any, Optional
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.util import Pt, Inches

from aippt.config import DEFAULT_REMOVE_COPYRIGHT, DEFAULT_AUTO_FIT
from aippt.logger import logger

try:
    from ppt_transitions import inject_transition
    from ppt_animations import inject_animations, RECOMMENDED_ANIMATIONS
    _ANIM_MODULES_READY = True
except Exception:
    _ANIM_MODULES_READY = False
    _ANIM_IMPORT_ERROR = "animation/transition modules not available"


class PptRenderer:
    """PPT 渲染引擎：模板 + meta + 结构化内容 → 成品 PPT"""

    def __init__(self, template_path: str, meta_path: str) -> None:
        template = Path(template_path)
        if not template.exists():
            raise FileNotFoundError(f"模板文件不存在: {template_path}")
        self.template_path = str(template)
        meta_file = Path(meta_path)
        if not meta_file.exists():
            raise FileNotFoundError(f"meta 文件不存在: {meta_path}")
        with open(meta_file, 'r', encoding='utf-8') as f:
            self.meta: dict[str, Any] = json.load(f)

    def render(
        self,
        slot_data: dict[str, dict[str, str]],
        output_path: str,
        remove_copyright: bool = DEFAULT_REMOVE_COPYRIGHT,
        auto_fit: bool = DEFAULT_AUTO_FIT,
        transitions: Optional[Any] = None,
        animations: Optional[Any] = None,
        notes_map: dict[int, str] = None,
        animation_theme: str = None,
    ) -> str:
        prs = Presentation(self.template_path)
        page_slots = self.meta.get('page_slots', {})
        stats: dict[str, int] = {'replaced': 0, 'missed': 0, 'skipped': 0}

        for page_str, slots in page_slots.items():
            page_num = int(page_str)
            if page_num < 1 or page_num > len(prs.slides):
                continue
            slide = prs.slides[page_num - 1]
            page_input = slot_data.get(page_str, {})
            if not page_input:
                continue

            available_shapes = [
                s for s in self._iter_all_shapes(slide.shapes)
                if s.has_text_frame and s.text_frame.text.strip()
            ]
            used: set[int] = set()

            for slot_info in slots:
                slot_name = slot_info.get('slot', '')
                # chart_data/table_data 由专门逻辑处理，跳过文本替换
                if slot_name in ('chart_data', 'table_data'):
                    continue
                if slot_name not in page_input:
                    stats['skipped'] += 1
                    continue
                new_value = page_input[slot_name]
                if new_value is None:
                    stats['skipped'] += 1
                    continue
                shape = self._find_shape(available_shapes, slot_info, used)
                if shape is None:
                    stats['missed'] += 1
                    logger.warning("页%s 槽位 %s 未匹配到 shape", page_str, slot_name)
                    continue
                original_text = shape.text_frame.text
                self._replace_text(shape, str(new_value))
                if auto_fit:
                    self._auto_fit(shape, str(new_value), original_text, slot_info, page_str, slot_name)
                used.add(id(shape._element))
                stats['replaced'] += 1

            # 处理图表/表格形状（检测到 chart_data/table_data 时触发）
            if 'chart_data' in page_input or 'table_data' in page_input:
                for shape in self._iter_all_shapes(slide.shapes):
                    if 'chart_data' in page_input and hasattr(shape, 'has_chart') and shape.has_chart:
                        self._replace_chart_data(shape, page_input['chart_data'])
                        stats['replaced'] += 1
                    elif 'table_data' in page_input and hasattr(shape, 'has_table') and shape.has_table:
                        self._fill_dynamic_table(shape, page_input['table_data'])
                        stats['replaced'] += 1

        # 演讲者备注注入：在删除版权页之前进行，page_id 与原始页码对齐
        # notes_map 键为 page_id（从 1 开始），值为该页备注文本；默认 None 不注入
        if notes_map:
            for idx, slide in enumerate(prs.slides):
                page_id = idx + 1
                if page_id in notes_map:
                    self._inject_notes(slide, notes_map[page_id])

        removed_pages: list[int] = []
        if remove_copyright and self.meta.get('removable_pages'):
            removed_pages = self._remove_slides(prs, self.meta['removable_pages'])

        # 动画主题：在 _inject_effects 之前，根据主题为每页设置默认 transition/animations
        # 优先级：单页显式 transitions/animations（dict）> 主题 page_overrides > 主题 global_transition
        if animation_theme:
            transitions, animations = self._apply_animation_theme(
                prs, animation_theme, transitions, animations
            )

        if transitions or animations:
            if not _ANIM_MODULES_READY:
                logger.warning("动画/转场模块加载失败，跳过注入: %s", _ANIM_IMPORT_ERROR)
            else:
                self._inject_effects(prs, transitions, animations)

        prs.save(str(output_path))
        logger.info("渲染完成: %s", output_path)
        logger.info("替换 %d / 未匹配 %d / 跳过 %d", stats['replaced'], stats['missed'], stats['skipped'])
        if removed_pages:
            logger.info("已删除版权页: %s", removed_pages)
        return str(output_path)

    def _apply_animation_theme(
        self,
        prs: Any,
        theme_name: str,
        transitions: Any,
        animations: Any,
    ) -> tuple[Any, Any]:
        """根据动画主题为每页构建默认 transition/animations，并与显式配置合并

        优先级：显式 transitions/animations（dict）> 主题 page_overrides > 主题 global_transition

        :param prs: Presentation 对象
        :param theme_name: 主题名（business/tech/formal）
        :param transitions: 已有的 transitions 配置（"auto"/"none"/dict/None）
        :param animations: 已有的 animations 配置（"auto"/"none"/dict/None）
        :return: 合并后的 (transitions, animations)
        """
        try:
            from aippt.animation_themes import (
                get_theme, build_page_transition_spec, build_page_animations_spec,
                SLIDE_TYPE_TO_PAGE_TYPE,
            )
        except Exception as e:
            logger.warning("动画主题模块加载失败，跳过主题注入: %s", e)
            return transitions, animations

        try:
            theme = get_theme(theme_name)
        except KeyError as e:
            logger.warning("%s，跳过主题注入", e)
            return transitions, animations

        chapters = self.meta.get('chapters', [])
        total = len(prs.slides)
        overrides = theme.get("page_overrides", {})
        global_transition = theme.get("global_transition")

        # 构建主题默认的 transitions/animations dict
        theme_transitions: dict[str, dict] = {}
        theme_animations: dict[str, list] = {}
        for page_num in range(1, total + 1):
            slide_type = self._infer_slide_type(page_num, chapters, total)
            page_type = SLIDE_TYPE_TO_PAGE_TYPE.get(slide_type, "numbered_list")
            page_cfg = overrides.get(page_type, {})

            # transition：page_overrides > global_transition
            t_name = page_cfg.get("transition") or global_transition
            if t_name and t_name != "none":
                t_spec = build_page_transition_spec(t_name)
                if t_spec:
                    theme_transitions[str(page_num)] = t_spec

            # animations：仅 page_overrides 中有配置的页面才注入
            a_cfg = page_cfg.get("animations")
            if a_cfg:
                a_spec = build_page_animations_spec(page_type, a_cfg)
                if a_spec:
                    theme_animations[str(page_num)] = a_spec

        # 合并：显式 dict 覆盖主题默认（单页 outline 配置优先级最高）
        if isinstance(transitions, dict):
            theme_transitions.update(transitions)
        if isinstance(animations, dict):
            theme_animations.update(animations)

        logger.info("动画主题 %s 已应用: %d 页转场, %d 页动画",
                    theme_name, len(theme_transitions), len(theme_animations))
        return (theme_transitions if theme_transitions else transitions,
                theme_animations if theme_animations else animations)

    def _inject_notes(self, slide, notes_text: str) -> None:
        """
        注入演讲者备注到 slide

        使用 slide.notes_slide.notes_text_frame.text 写入备注文本；
        若 slide 无 notes_slide，python-pptx 会自动创建。

        :param slide: python-pptx 的 Slide 对象
        :param notes_text: 备注文本
        """
        try:
            slide.notes_slide.notes_text_frame.text = notes_text
        except Exception as e:
            logger.warning("演讲者备注注入失败: %s", e)

    def _inject_effects(
        self,
        prs: Any,
        transitions: Any,
        animations: Any,
    ) -> None:
        chapters = self.meta.get('chapters', [])
        total = len(prs.slides)

        for page_num, slide in enumerate(prs.slides, 1):
            if transitions:
                if transitions == "auto":
                    inject_transition(slide, {"type": "fade", "speed": "med"})
                elif isinstance(transitions, dict) and str(page_num) in transitions:
                    inject_transition(slide, transitions[str(page_num)])

            if animations:
                slide_type = self._infer_slide_type(page_num, chapters, total)

                if animations == "auto":
                    anim_spec = RECOMMENDED_ANIMATIONS.get(
                        slide_type, RECOMMENDED_ANIMATIONS.get("CONTENT", [])
                    )
                    if anim_spec:
                        inject_animations(slide, anim_spec, slide_type)
                elif isinstance(animations, dict) and str(page_num) in animations:
                    anim_spec = animations[str(page_num)]
                    if anim_spec:
                        inject_animations(slide, anim_spec, slide_type)

    def _infer_slide_type(self, page_num: int, chapters: list[dict[str, Any]], total_pages: int) -> str:
        # 优先识别复合页面：含 chart/table 形状的页面
        page_meta_entry = self.meta.get('page_meta', {}).get(str(page_num))
        if page_meta_entry:
            if page_meta_entry.get("has_chart"):
                return "CHART"
            if page_meta_entry.get("has_table"):
                return "TABLE"

        # 章节页匹配
        for ch in chapters:
            key = ch.get("key", "")
            page = ch.get("page")
            start_page = ch.get("start_page")
            end_page = ch.get("end_page")

            if key == "cover" and page == page_num:
                return "COVER"
            if key == "end" and page == page_num:
                return "END"
            # 章节分隔页（单页 chapter）
            if key and key.startswith("chapter_") and start_page == end_page and start_page == page_num:
                return "CHAPTER"
            # 章节首页也可能是分隔页
            if key and key.startswith("chapter_") and start_page == page_num and start_page != end_page:
                # 仅当 page_slots 中 title 含 PART 字样时算 CHAPTER
                page_slots = self.meta.get('page_slots', {}).get(str(page_num), [])
                for s in page_slots:
                    mt = s.get('match_text', '')
                    if 'PART' in mt.upper() or '第' in mt and '章' in mt:
                        return "CHAPTER"
                return "CONTENT"

        # 末页通常是结束页
        if page_num == total_pages:
            return "END"

        return "CONTENT"

    def _iter_all_shapes(self, shapes):
        for s in shapes:
            yield s
            if s.shape_type == 6:  # GROUP
                yield from self._iter_all_shapes(s.shapes)

    def _find_shape(self, shapes, slot_info: dict[str, Any], used: set[int]) -> Optional[BaseShape]:
        match_text = slot_info.get('match_text', '').strip()
        shape_name = slot_info.get('shape_name', '')
        all_shapes = list(self._iter_all_shapes(shapes))

        # 策略1：shape_name + match_text 双匹配（最精确）
        if shape_name and match_text:
            for s in all_shapes:
                if id(s._element) in used:
                    continue
                if not s.has_text_frame:
                    continue
                if s.name == shape_name and match_text in s.text_frame.text:
                    return s

        # 策略2：match_text 文本匹配
        if match_text:
            for s in all_shapes:
                if id(s._element) in used:
                    continue
                if not s.has_text_frame:
                    continue
                if match_text in s.text_frame.text:
                    return s

        # 策略3：shape_name 匹配（兜底）
        if shape_name:
            for s in all_shapes:
                if id(s._element) in used:
                    continue
                if not s.has_text_frame:
                    continue
                if s.name == shape_name:
                    return s

        return None

    def _replace_text(self, shape: BaseShape, new_text: str) -> None:
        tf = shape.text_frame
        paragraphs = list(tf.paragraphs)
        if not paragraphs:
            return

        first_para = paragraphs[0]

        # 保留第一个 run 的格式，仅改其 text，删除同段其余 run
        if first_para.runs:
            first_para.runs[0].text = new_text
            for run in first_para.runs[1:]:
                run._r.getparent().remove(run._r)
        else:
            run = first_para.add_run()
            run.text = new_text

        # 删除其余段落（保留首段段落属性）
        for para in paragraphs[1:]:
            para._p.getparent().remove(para._p)

    def _auto_fit(
        self,
        shape: BaseShape,
        new_text: str,
        original_text: str,
        slot_info: Optional[dict[str, Any]] = None,
        page_str: Optional[str] = None,
        slot_name: Optional[str] = None,
    ) -> None:
        tf = shape.text_frame
        if not tf.paragraphs or not tf.paragraphs[0].runs:
            return
        run = tf.paragraphs[0].runs[0]
        if run.font.size is None:
            return
        original_size_pt = run.font.size.pt

        # 策略1：优先使用 meta 中的 capacity（事前预警）
        capacity = None
        if slot_info and isinstance(slot_info.get('capacity'), dict):
            capacity = slot_info['capacity'].get('total_chars')

        # 策略2：回退到运行时几何估算
        if capacity is None:
            try:
                box_w = shape.width        # EMU
                box_h = shape.height       # EMU
            except Exception:
                return
            if not box_w or not box_h:
                return
            EMU_PER_PT = 12700
            char_w = original_size_pt * 0.55 * EMU_PER_PT
            line_h = original_size_pt * 1.2 * EMU_PER_PT
            chars_per_line = max(1, int(box_w / char_w))
            lines_available = max(1, int(box_h / line_h))
            capacity = chars_per_line * lines_available

        text_len = len(new_text)
        if text_len <= capacity:
            return

        if page_str and slot_name:
            logger.warning("页%s 槽位 %s: 文本长度 %d 超过容量 %d，自动缩字号", page_str, slot_name, text_len, capacity)

        # 按比例缩小字号，下限 8pt
        ratio = (capacity / text_len) ** 0.5
        new_size = max(8, int(original_size_pt * ratio))
        if new_size < original_size_pt:
            run.font.size = Pt(new_size)

    # ==================== 图表与表格动态扩展 ====================

    def _replace_chart_data(self, shape: BaseShape, chart_data_dict: dict[str, Any]) -> None:
        """替换图表数据，100% 保留模板样式（字体/颜色/坐标轴格式/图例位置）

        支持 4 类图表：bar（柱状图）、line（折线图）、pie（饼图）、radar（雷达图）
        多系列自动适配：M>N 仅替换前 N 系列；M<N 多余系列清空数据；M==N 直接替换
        使用 chart.replace_data(ChartData) API，保留图表样式（系列数与模板一致）

        :param shape: GraphicFrame 形状，含 chart
        :param chart_data_dict: {"categories": [...], "series": [{"name": "...", "data": [...]}]}
        """
        from pptx.chart.data import ChartData

        try:
            chart = shape.chart
        except Exception as e:
            logger.warning("无法访问图表数据: %s", e)
            return

        # 检查图表类型，映射到 bar/line/pie/radar 四大类
        ct_str = str(chart.chart_type).upper()
        if 'BAR' in ct_str or 'COLUMN' in ct_str:
            chart_category = 'bar'
        elif 'LINE' in ct_str:
            chart_category = 'line'
        elif 'PIE' in ct_str:
            chart_category = 'pie'
        elif 'RADAR' in ct_str:
            chart_category = 'radar'
        else:
            logger.warning("不支持的图表类型: %s，跳过数据替换", chart.chart_type)
            return

        categories = chart_data_dict.get('categories', [])
        new_series = chart_data_dict.get('series', [])
        if not categories or not new_series:
            logger.warning("图表数据为空（categories=%d, series=%d），跳过",
                           len(categories), len(new_series))
            return

        try:
            plot = chart.plots[0]
        except Exception as e:
            logger.warning("无法访问图表 plot: %s", e)
            return

        # 多系列自动适配：缓存 series 集合到本地变量，获取模板原有系列数
        chart_series = list(plot.series)
        n = len(chart_series)  # 模板原有系列数
        m = len(new_series)    # 新数据系列数

        if m > n:
            logger.warning("图表系列数不匹配：模板 %d 系列，新数据 %d 系列，仅替换前 %d 系列",
                           n, m, n)
        elif m < n:
            logger.warning("图表系列数不匹配：模板 %d 系列，新数据 %d 系列，多余 %d 系列清空数据",
                           n, m, n - m)

        # 构造 ChartData，保持系列数与模板一致（N 系列），保留图表样式
        chart_data = ChartData()
        chart_data.categories = categories
        for i in range(n):
            if i < m:
                # 替换为新数据
                ser = new_series[i]
                chart_data.add_series(ser.get('name', ''), ser.get('data', []))
            else:
                # 多余系列清空数据（置零，保留系列结构避免破坏图表样式）
                orig_name = ''
                try:
                    orig_name = chart_series[i].name or ''
                except Exception:
                    pass
                chart_data.add_series(orig_name, [0] * len(categories))

        # 替换图表数据（replace_data 保留图表类型、坐标轴格式、图例等样式）
        try:
            chart.replace_data(chart_data)
        except Exception as e:
            # replace_data 可能因外部 Excel 链接报错（.target_part undefined），
            # 但图表 XML 数据通常已更新（分类/系列值已写入），图表显示正常
            logger.warning("图表数据替换（Excel 同步异常，图表 XML 已更新）: %s", e)

        logger.info("图表数据替换完成（类型=%s, 系列=%d/%d, 分类=%d）",
                    chart_category, min(m, n), n, len(categories))

    def _fill_dynamic_table(self, shape: BaseShape, table_data: dict[str, Any]) -> None:
        """动态填充表格数据，保留第 0 行表头样式模板

        - 根据数据行数动态增减行（克隆最后一行保持样式一致性）
        - 填入表头数据（覆盖模板占位文本，保留样式）
        - 填入数据行，每行单元格继承样式（字体/对齐/背景色）
        - 自动行高适配

        :param shape: GraphicFrame 形状，含 table
        :param table_data: {"headers": [...], "rows": [[...], ...]}
        """
        try:
            table = shape.table
        except Exception as e:
            logger.warning("无法访问表格数据: %s", e)
            return

        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        if not headers:
            logger.warning("表格 headers 为空，跳过填充")
            return

        current_rows = len(table.rows)
        if current_rows == 0:
            logger.warning("表格无行数据，无法填充")
            return
        needed_rows = len(rows) + 1  # +1 为表头行

        # 动态增减行（python-pptx 的 _RowCollection 不支持负索引，需用正索引）
        if needed_rows > current_rows:
            # 新增行：以最后一行为模板克隆（保留行结构与单元格样式）
            last_tr = table.rows[current_rows - 1]._tr
            added = 0
            for _ in range(needed_rows - current_rows):
                new_tr = deepcopy(last_tr)
                last_tr.addnext(new_tr)
                last_tr = new_tr
                added += 1
            logger.info("表格新增 %d 行（%d → %d）", added, current_rows, needed_rows)
        elif needed_rows < current_rows:
            # 删除多余行（从末尾删除，保留表头）
            extra = current_rows - needed_rows
            for _ in range(extra):
                remaining = len(table.rows)
                last_tr = table.rows[remaining - 1]._tr
                last_tr.getparent().remove(last_tr)
            logger.info("表格删除 %d 行（%d → %d）", extra, current_rows, needed_rows)

        # 填入表头数据（第 0 行，保留单元格样式）
        col_count = len(table.columns)
        for col_idx, header_text in enumerate(headers):
            if col_idx < col_count:
                cell = table.cell(0, col_idx)
                self._set_cell_text(cell, str(header_text))

        # 填入数据行
        for row_idx, row_data in enumerate(rows, start=1):
            if row_idx >= len(table.rows):
                break
            # 自动行高适配：默认 0.4 英寸，内容较长时增至 0.6
            try:
                max_cell_len = max((len(str(c)) for c in row_data), default=0)
                row_height = Inches(0.6) if max_cell_len > 20 else Inches(0.4)
                table.rows[row_idx].height = row_height
            except Exception:
                pass
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < col_count:
                    cell = table.cell(row_idx, col_idx)
                    self._set_cell_text(cell, str(cell_text))

        logger.info("表格填充完成（表头=%d 列, 数据行=%d）", min(len(headers), col_count), len(rows))

    def _set_cell_text(self, cell: Any, text: str) -> None:
        """设置单元格文本，保留首个 run 的格式（字体/字号/颜色/粗体）

        :param cell: pptx table Cell 对象
        :param text: 要填入的文本
        """
        tf = cell.text_frame
        paragraphs = list(tf.paragraphs)
        if not paragraphs:
            cell.text = str(text)
            return
        first_para = paragraphs[0]
        if first_para.runs:
            # 保留首个 run 的格式，仅替换其 text，删除同段其余 run
            first_para.runs[0].text = str(text)
            for run in first_para.runs[1:]:
                run._r.getparent().remove(run._r)
        else:
            # 没有现有 run，直接设置（会创建新 run，继承段落/单元格级格式）
            cell.text = str(text)
        # 删除多余段落（保留首段段落属性）
        for para in paragraphs[1:]:
            para._p.getparent().remove(para._p)

    def _remove_slides(self, prs: Any, page_numbers: list[int]) -> list[int]:
        removed: list[int] = []
        for page_num in sorted(page_numbers, reverse=True):
            idx = page_num - 1
            if 0 <= idx < len(prs.slides):
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                xml_slides.remove(slides[idx])
                removed.append(page_num)
        return removed

    def duplicate_slide(self, prs: Any, slide_index: int):
        source = prs.slides[slide_index]
        new_slide = prs.slides.add_slide(source.slide_layout)
        # 清空新 slide 自带占位符
        for shape in list(new_slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
        # 深拷贝源 slide 所有 shape
        for shape in source.shapes:
            new_el = deepcopy(shape._element)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        return new_slide


# ==================== CLI 入口 ====================
def cmd_render(args):
    with open(args.data, 'r', encoding='utf-8') as f:
        slot_data = json.load(f)
    renderer = PptRenderer(args.template, args.meta)
    renderer.render(
        slot_data,
        args.output,
        remove_copyright=not args.keep_copyright,
        auto_fit=not args.no_fit
    )


def main():
    parser = argparse.ArgumentParser(description='PPT 渲染引擎')
    sub = parser.add_subparsers(dest='command', required=True)

    r = sub.add_parser('render', help='渲染生成 PPT')
    r.add_argument('--template', required=True, help='模板 pptx 路径')
    r.add_argument('--meta', required=True, help='meta.json 路径')
    r.add_argument('--data', required=True, help='内容 JSON 路径（{"页码":{"槽位名":"值"}}）')
    r.add_argument('--output', required=True, help='输出 pptx 路径')
    r.add_argument('--keep-copyright', action='store_true', help='保留版权页')
    r.add_argument('--no-fit', action='store_true', help='关闭字号自适应')
    r.set_defaults(func=cmd_render)

    args = parser.parse_args()
    args.func(args)


if __name__ == '__main__':
    main()
